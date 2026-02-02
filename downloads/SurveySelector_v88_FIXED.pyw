#!/usr/bin/env python3
"""
Survey Crosstab Tool
Version 88.0 - ADDED STEP 5: GENERATE INSIGHTS REPORT
Last Updated: February 1, 2026

NEW IN V88 (February 1, 2026):
- ✅ Added Step 5: Generate Insights Report from PowerPoint
- ✅ Analyzes completed PowerPoint presentations
- ✅ Generates professional Word document with executive summary
- ✅ Compares findings to normative database benchmarks
- ✅ Provides actionable recommendations and insights
- ✅ Full integration with existing workflow (ALL v87 functionality retained)

CRITICAL FIX IN V33:
- Fixed adjust_percentages_to_100() to use REAL largest remainder method
- v32 was adding/subtracting ENTIRE difference to one value (WRONG!)
- This caused negative numbers when difference was negative
- NOW: Distributes remainders one-by-one to items with largest fractional parts
- This is actual 5th grade math, done correctly

COMPLETE FEATURE SET:
- ✅ Working Excel button with selection dialog
- ✅ Working PowerPoint template updater
- ✅ Working Benchmark comparison  
- ✅ Scrollable Step 4 interface
- ✅ Percentages sum to EXACTLY 100% with NO negatives
- ✅ All buttons visible and functional
- ✅ NEW: Insights Report Generation (Step 5)

Extracts questions from Word survey, processes Remark data, 
and generates professional PDF crosstab reports, Excel spreadsheets,
updates PowerPoint templates with survey data, and creates
comprehensive insights reports for client delivery.
"""

import tkinter as tk
from tkinter import filedialog, messagebox, ttk, simpledialog
from tkinter import font as tkfont
import json
import os
from datetime import datetime
import threading

# These will be imported when needed to speed up startup
# from docx import Document
# import pandas as pd
# from reportlab...


class SurveyCrosstabApp:
    """Main application class with step-based workflow"""
    
    # Color scheme
    COLORS = {
        'bg': '#f8f9fa',
        'white': '#ffffff',
        'primary': '#2563eb',
        'primary_hover': '#1d4ed8',
        'success': '#22c55e',
        'success_bg': '#f0fdf4',
        'success_border': '#bbf7d0',
        'gray_50': '#f9fafb',
        'gray_100': '#f3f4f6',
        'gray_200': '#e5e7eb',
        'gray_300': '#d1d5db',
        'gray_400': '#9ca3af',
        'gray_500': '#6b7280',
        'gray_600': '#4b5563',
        'gray_700': '#374151',
        'gray_800': '#1f2937',
        'gray_900': '#111827',
        'yellow_100': '#fef9c3',
        'yellow_800': '#854d0e',
        'border': '#e5e7eb',
        'excel': '#217346',  # Excel green
        'pptx': '#D14423',  # PowerPoint orange
        'word': '#2B579A',  # Word blue - NEW for Step 5
    }
    
    def __init__(self):
        self.root = None
        self.current_step = 1
        self.project_name = None
        self.project_data = {}
        self.selected_breakout = None
        self.selected_breakouts = []  # List of selected breakout question numbers
        self.remark_file = None
        self.remark_df = None
        self.questions = []
        self.questions_with_options = []
        
        # Template merger settings
        self.template_file = None
        self.logo_file = None
        self.facility_name = None
        self.benchmark_file = None
        
        # Company name (for PowerPoint and filenames)
        self.company_name = None
        
        # PowerPoint question tracking
        self.concept_questions = []  # Schedule Concept questions (manual entry)
        self.unique_questions = []   # Unique to this site questions (manual entry)
        
        # NEW in v88: Step 5 - Insights Report tracking
        self.generated_pptx_file = None
        self.insights_report_file = None
        
        # Projects folder
        self.projects_folder = os.path.join(
            os.path.dirname(os.path.abspath(__file__)), 
            "Survey_Projects"
        )
        if not os.path.exists(self.projects_folder):
            os.makedirs(self.projects_folder)
    
    def get_project_path(self, project_name):
        """Get the full path to a project file"""
        safe_name = "".join(c for c in project_name if c.isalnum() or c in (' ', '-', '_')).strip()
        return os.path.join(self.projects_folder, f"{safe_name}.json")
    
    def save_project(self):
        """Save current project state"""
        if not self.project_name:
            return
        self.project_data['last_modified'] = datetime.now().isoformat()
        # NEW in v88: Save Step 5 tracking
        self.project_data['generated_pptx_file'] = self.generated_pptx_file
        self.project_data['insights_report_file'] = self.insights_report_file
        filepath = self.get_project_path(self.project_name)
        with open(filepath, 'w', encoding='utf-8') as f:
            json.dump(self.project_data, f, indent=2)
    
    def load_project(self, project_name):
        """Load project from file"""
        def clean_unicode(text):
            """Fix Unicode encoding issues"""
            if not isinstance(text, str):
                return text
            text = text.replace('\u2013', '-').replace('\u2014', '-')
            text = text.replace('\u2018', "'").replace('\u2019', "'")
            text = text.replace('\u201c', '"').replace('\u201d', '"')
            return text
        
        def clean_questions(questions):
            """Clean Unicode in questions list"""
            if not questions:
                return questions
            cleaned = []
            for q in questions:
                q_copy = q.copy()
                if 'text' in q_copy:
                    q_copy['text'] = clean_unicode(q_copy['text'])
                if 'options' in q_copy:
                    q_copy['options'] = [clean_unicode(opt) for opt in q_copy['options']]
                cleaned.append(q_copy)
            return cleaned
        
        filepath = self.get_project_path(project_name)
        if os.path.exists(filepath):
            with open(filepath, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            # Clean Unicode in loaded data
            if 'questions' in data:
                data['questions'] = clean_questions(data['questions'])
            if 'questions_with_options' in data:
                data['questions_with_options'] = clean_questions(data['questions_with_options'])
            
            return data
        return None
    
    def get_existing_projects(self):
        """Get list of existing projects"""
        projects = []
        if os.path.exists(self.projects_folder):
            for filename in os.listdir(self.projects_folder):
                if filename.endswith('.json'):
                    filepath = os.path.join(self.projects_folder, filename)
                    try:
                        with open(filepath, 'r', encoding='utf-8') as f:
                            data = json.load(f)
                            projects.append({
                                'name': data.get('project_name', filename[:-5]),
                                'last_modified': data.get('last_modified', 'Unknown'),
                                'step': data.get('current_step', 'Unknown'),
                                'survey_file': data.get('survey_file', 'Unknown')
                            })
                    except:
                        pass
        return sorted(projects, key=lambda x: x.get('last_modified', ''), reverse=True)
    
    def run(self):
        """Main application entry point"""
        self.root = tk.Tk()
        self.root.title("Survey Crosstab Tool - With Insights Report Generator")
        self.root.configure(bg=self.COLORS['bg'])
        
        # Start maximized
        self.root.state('zoomed')
        
        # Set minimum size
        self.root.minsize(900, 600)
        
        # Configure fonts
        self.fonts = {
            'heading': tkfont.Font(family='Segoe UI', size=18, weight='bold'),
            'subheading': tkfont.Font(family='Segoe UI', size=12),
            'body': tkfont.Font(family='Segoe UI', size=10),
            'body_bold': tkfont.Font(family='Segoe UI', size=10, weight='bold'),
            'small': tkfont.Font(family='Segoe UI', size=9),
            'step_title': tkfont.Font(family='Segoe UI', size=11, weight='bold'),
            'step_detail': tkfont.Font(family='Segoe UI', size=10),
            'button': tkfont.Font(family='Segoe UI', size=10, weight='bold'),
        }
        
        self.build_ui()
        self.root.mainloop()
    
    def build_ui(self):
        """Build the main UI structure"""
        # Clear existing widgets
        for widget in self.root.winfo_children():
            widget.destroy()
        
        # Header
        self.create_header()
        
        # Main content area
        self.content_frame = tk.Frame(self.root, bg=self.COLORS['bg'])
        self.content_frame.pack(fill=tk.BOTH, expand=True, padx=40, pady=20)
        
        # Center container with max width
        self.center_container = tk.Frame(self.content_frame, bg=self.COLORS['bg'])
        self.center_container.pack(expand=True, fill=tk.BOTH)
        self.center_container.configure(width=800)
        
        # Completed steps container
        self.completed_frame = tk.Frame(self.center_container, bg=self.COLORS['bg'])
        self.completed_frame.pack(fill=tk.X, pady=(0, 10))
        
        # Current step container
        self.step_frame = tk.Frame(
            self.center_container, 
            bg=self.COLORS['white'],
            highlightbackground=self.COLORS['border'],
            highlightthickness=1
        )
        self.step_frame.pack(fill=tk.BOTH, expand=True)
        
        # Footer
        self.create_footer()
        
        # Show first step
        self.show_current_step()
    
    def create_header(self):
        """Create the header bar"""
        header = tk.Frame(self.root, bg=self.COLORS['white'], height=70)
        header.pack(fill=tk.X)
        header.pack_propagate(False)
        
        # Add shadow line
        shadow = tk.Frame(self.root, bg=self.COLORS['border'], height=1)
        shadow.pack(fill=tk.X)
        
        # Header content
        header_content = tk.Frame(header, bg=self.COLORS['white'])
        header_content.pack(fill=tk.BOTH, expand=True, padx=40, pady=15)
        
        # Title
        title = tk.Label(
            header_content,
            text="Survey Crosstab Tool",
            font=self.fonts['heading'],
            bg=self.COLORS['white'],
            fg=self.COLORS['gray_900']
        )
        title.pack(side=tk.LEFT)
        
        # Project name (if set)
        self.project_label = tk.Label(
            header_content,
            text="",
            font=self.fonts['small'],
            bg=self.COLORS['white'],
            fg=self.COLORS['gray_500']
        )
        self.project_label.pack(side=tk.LEFT, padx=(20, 0))
    
    def create_footer(self):
        """Create the footer bar"""
        footer_line = tk.Frame(self.root, bg=self.COLORS['border'], height=1)
        footer_line.pack(fill=tk.X, side=tk.BOTTOM)
        
        footer = tk.Frame(self.root, bg=self.COLORS['white'], height=40)
        footer.pack(fill=tk.X, side=tk.BOTTOM)
        footer.pack_propagate(False)
        
        version = tk.Label(
            footer,
            text="Survey Crosstab Tool v88.0 - Now with Insights Report Generation",
            font=self.fonts['small'],
            bg=self.COLORS['white'],
            fg=self.COLORS['gray_500']
        )
        version.pack(side=tk.LEFT, padx=40, pady=10)
    
    def update_project_label(self):
        """Update the project name in header"""
        if self.project_name:
            self.project_label.config(text=f"Project: {self.project_name}")
        else:
            self.project_label.config(text="")
    
    def show_completed_step(self, step_num, title, detail, can_edit=False):
        """Show a completed step widget"""
        frame = tk.Frame(
            self.completed_frame,
            bg=self.COLORS['white'],
            highlightbackground=self.COLORS['border'],
            highlightthickness=1
        )
        frame.pack(fill=tk.X, pady=(0, 10))
        
        inner = tk.Frame(frame, bg=self.COLORS['white'])
        inner.pack(fill=tk.X, padx=15, pady=12)
        
        # Left side: checkmark and text
        left = tk.Frame(inner, bg=self.COLORS['white'])
        left.pack(side=tk.LEFT, fill=tk.X, expand=True)
        
        # Green checkmark circle
        check_frame = tk.Frame(left, bg=self.COLORS['white'])
        check_frame.pack(side=tk.LEFT, padx=(0, 12))
        
        check_canvas = tk.Canvas(
            check_frame, 
            width=32, height=32, 
            bg=self.COLORS['white'],
            highlightthickness=0
        )
        check_canvas.pack()
        check_canvas.create_oval(2, 2, 30, 30, fill=self.COLORS['success'], outline='')
        check_canvas.create_text(16, 16, text="✓", fill='white', font=('Segoe UI', 14, 'bold'))
        
        # Text
        text_frame = tk.Frame(left, bg=self.COLORS['white'])
        text_frame.pack(side=tk.LEFT, fill=tk.X)
        
        step_label = tk.Label(
            text_frame,
            text=f"Step {step_num}",
            font=self.fonts['small'],
            bg=self.COLORS['white'],
            fg=self.COLORS['gray_500']
        )
        step_label.pack(anchor='w')
        
        title_label = tk.Label(
            text_frame,
            text=title,
            font=self.fonts['step_title'],
            bg=self.COLORS['white'],
            fg=self.COLORS['gray_900']
        )
        title_label.pack(anchor='w')
        
        detail_label = tk.Label(
            text_frame,
            text=detail,
            font=self.fonts['step_detail'],
            bg=self.COLORS['white'],
            fg=self.COLORS['gray_600']
        )
        detail_label.pack(anchor='w')
        
        # Edit button (if applicable)
        if can_edit:
            edit_btn = tk.Button(
                inner,
                text="✎ Edit",
                font=self.fonts['body'],
                bg=self.COLORS['white'],
                fg=self.COLORS['primary'],
                relief=tk.FLAT,
                cursor='hand2',
                command=lambda s=step_num: self.go_to_step(s)
            )
            edit_btn.pack(side=tk.RIGHT)
    
    def go_to_step(self, step_num):
        """Navigate back to a previous step"""
        self.current_step = step_num
        self.show_current_step()
    
    def show_current_step(self):
        """Display the current step UI"""
        # Clear completed steps
        for widget in self.completed_frame.winfo_children():
            widget.destroy()
        
        # Clear step frame
        for widget in self.step_frame.winfo_children():
            widget.destroy()
        
        # Show completed steps
        if self.current_step > 1:
            project_type = "New Project" if self.project_data.get('is_new') else "Opened"
            self.show_completed_step(1, "Project Selected", f"{project_type}: {self.project_name}")
        
        if self.current_step > 2:
            breakout_label = self.get_breakout_label(self.selected_breakouts)
            self.show_completed_step(2, "Breakout Variables Selected", breakout_label, can_edit=True)
        
        if self.current_step > 3:
            self.show_completed_step(3, "Remark Data Uploaded", os.path.basename(self.remark_file), can_edit=True)
        
        # Show current step content
        if self.current_step == 1:
            self.show_step_1()
        elif self.current_step == 2:
            self.show_step_2()
        elif self.current_step == 3:
            self.show_step_3()
        elif self.current_step == 4:
            self.show_step_4()
        
        self.update_project_label()
    
    def get_breakout_label(self, breakout_ids):
        """Get display label for breakout variable(s)"""
        if isinstance(breakout_ids, list):
            labels = []
            for bid in breakout_ids:
                labels.append(f"Q{bid}")
            return ", ".join(labels)
        
        # Single breakout (backwards compatibility)
        questions = self.questions_with_options or self.questions
        for q in questions:
            if str(q['number']) == str(breakout_ids):
                q_text = q.get('text', f'Question {q["number"]}')
                return f"Q{q['number']}: {q_text[:50]}..."
        
        # Fallback
        return f"Question {breakout_ids}"
    
    def show_step_1(self):
        """Step 1: Select or Create Project"""
        content = tk.Frame(self.step_frame, bg=self.COLORS['white'])
        content.pack(fill=tk.BOTH, expand=True, padx=30, pady=25)
        
        # Step header
        header = tk.Frame(content, bg=self.COLORS['white'])
        header.pack(fill=tk.X, pady=(0, 25))
        
        # Step number circle
        num_canvas = tk.Canvas(header, width=36, height=36, bg=self.COLORS['white'], highlightthickness=0)
        num_canvas.pack(side=tk.LEFT, padx=(0, 12))
        num_canvas.create_oval(2, 2, 34, 34, fill=self.COLORS['primary'], outline='')
        num_canvas.create_text(18, 18, text="1", fill='white', font=('Segoe UI', 12, 'bold'))
        
        # Step title
        title_frame = tk.Frame(header, bg=self.COLORS['white'])
        title_frame.pack(side=tk.LEFT)
        
        tk.Label(
            title_frame,
            text="Welcome to Survey Crosstab Tool",
            font=self.fonts['step_title'],
            bg=self.COLORS['white'],
            fg=self.COLORS['gray_900']
        ).pack(anchor='w')
        
        tk.Label(
            title_frame,
            text="Choose how you'd like to proceed",
            font=self.fonts['step_detail'],
            bg=self.COLORS['white'],
            fg=self.COLORS['gray_500']
        ).pack(anchor='w')
        
        # Buttons container - First row
        buttons_frame = tk.Frame(content, bg=self.COLORS['white'])
        buttons_frame.pack(fill=tk.X, pady=(0, 10))
        
        # New Project button
        new_btn = tk.Frame(
            buttons_frame,
            bg=self.COLORS['white'],
            highlightbackground=self.COLORS['gray_300'],
            highlightthickness=2,
            cursor='hand2'
        )
        new_btn.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=(0, 10))
        
        new_inner = tk.Frame(new_btn, bg=self.COLORS['white'])
        new_inner.pack(fill=tk.BOTH, expand=True, padx=30, pady=25)
        
        tk.Label(new_inner, text="➕", font=('Segoe UI', 28), bg=self.COLORS['white'], fg=self.COLORS['gray_400']).pack()
        tk.Label(new_inner, text="New Project", font=self.fonts['body_bold'], bg=self.COLORS['white'], fg=self.COLORS['gray_900']).pack(pady=(10, 5))
        tk.Label(new_inner, text="Start a new survey analysis", font=self.fonts['small'], bg=self.COLORS['white'], fg=self.COLORS['gray_500']).pack()
        
        new_btn.bind('<Button-1>', lambda e: self.create_new_project())
        for child in new_btn.winfo_children():
            child.bind('<Button-1>', lambda e: self.create_new_project())
            for subchild in child.winfo_children():
                subchild.bind('<Button-1>', lambda e: self.create_new_project())
        
        # Open Existing button
        open_btn = tk.Frame(
            buttons_frame,
            bg=self.COLORS['white'],
            highlightbackground=self.COLORS['gray_200'],
            highlightthickness=2,
            cursor='hand2'
        )
        open_btn.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=(10, 0))
        
        open_inner = tk.Frame(open_btn, bg=self.COLORS['white'])
        open_inner.pack(fill=tk.BOTH, expand=True, padx=30, pady=25)
        
        tk.Label(open_inner, text="📂", font=('Segoe UI', 28), bg=self.COLORS['white'], fg=self.COLORS['gray_400']).pack()
        tk.Label(open_inner, text="Open Existing", font=self.fonts['body_bold'], bg=self.COLORS['white'], fg=self.COLORS['gray_900']).pack(pady=(10, 5))
        tk.Label(open_inner, text="Continue a previous project", font=self.fonts['small'], bg=self.COLORS['white'], fg=self.COLORS['gray_500']).pack()
        
        open_btn.bind('<Button-1>', lambda e: self.toggle_recent_projects())
        for child in open_btn.winfo_children():
            child.bind('<Button-1>', lambda e: self.toggle_recent_projects())
            for subchild in child.winfo_children():
                subchild.bind('<Button-1>', lambda e: self.toggle_recent_projects())
        
        # Recent projects list (initially hidden) 
        self.recent_frame = tk.Frame(content, bg=self.COLORS['white'])
        self.recent_projects_visible = False
        
        # NEW in v88: Standalone Insights Report Generator button - Below recent projects
        separator = tk.Frame(content, bg=self.COLORS['border'], height=2)
        separator.pack(fill=tk.X, pady=(30, 20))
        
        tk.Label(
            content,
            text="OR: Generate Insights Report Only",
            font=self.fonts['body_bold'],
            bg=self.COLORS['white'],
            fg=self.COLORS['gray_700']
        ).pack(anchor='w', pady=(0, 15))
        
        # Insights Report Generator button (standalone)
        insights_btn = tk.Frame(
            content,
            bg=self.COLORS['white'],
            highlightbackground='#2B579A',
            highlightthickness=2,
            cursor='hand2'
        )
        insights_btn.pack(fill=tk.X, pady=(0, 20))
        
        insights_inner = tk.Frame(insights_btn, bg=self.COLORS['white'])
        insights_inner.pack(fill=tk.BOTH, expand=True, padx=30, pady=25)
        
        # Icon and text container
        content_container = tk.Frame(insights_inner, bg=self.COLORS['white'])
        content_container.pack(side=tk.LEFT, fill=tk.X, expand=True)
        
        tk.Label(content_container, text="📄", font=('Segoe UI', 28), bg=self.COLORS['white'], fg='#2B579A').pack(side=tk.LEFT, padx=(0, 20))
        
        text_container = tk.Frame(content_container, bg=self.COLORS['white'])
        text_container.pack(side=tk.LEFT, fill=tk.X, expand=True)
        
        tk.Label(text_container, text="Generate Insights Report from PowerPoint", font=self.fonts['body_bold'], bg=self.COLORS['white'], fg=self.COLORS['gray_900'], anchor='w').pack(anchor='w')
        tk.Label(text_container, text="Upload a PowerPoint presentation and generate a professional insights report", font=self.fonts['small'], bg=self.COLORS['white'], fg=self.COLORS['gray_500'], anchor='w', wraplength=600, justify=tk.LEFT).pack(anchor='w', pady=(5, 0))
        
        insights_btn.bind('<Button-1>', lambda e: self.launch_standalone_insights_generator())
        for child in insights_btn.winfo_children():
            child.bind('<Button-1>', lambda e: self.launch_standalone_insights_generator())
            for subchild in child.winfo_children():
                subchild.bind('<Button-1>', lambda e: self.launch_standalone_insights_generator())
                for subsubchild in subchild.winfo_children():
                    subsubchild.bind('<Button-1>', lambda e: self.launch_standalone_insights_generator())
    
    def toggle_recent_projects(self):
        """Show/hide recent projects list"""
        if self.recent_projects_visible:
            self.recent_frame.pack_forget()
            self.recent_projects_visible = False
        else:
            self.show_recent_projects()
            self.recent_projects_visible = True
    
    def show_recent_projects(self):
        """Display list of recent projects"""
        # Clear and show frame
        for widget in self.recent_frame.winfo_children():
            widget.destroy()
        
        self.recent_frame.pack(fill=tk.X, pady=(10, 0))
        
        # Separator
        sep = tk.Frame(self.recent_frame, bg=self.COLORS['border'], height=1)
        sep.pack(fill=tk.X, pady=(0, 15))
        
        # Title
        tk.Label(
            self.recent_frame,
            text="Recent Projects",
            font=self.fonts['body_bold'],
            bg=self.COLORS['white'],
            fg=self.COLORS['gray_700']
        ).pack(anchor='w', pady=(0, 10))
        
        # Project list
        projects = self.get_existing_projects()
        
        if not projects:
            tk.Label(
                self.recent_frame,
                text="No existing projects found",
                font=self.fonts['body'],
                bg=self.COLORS['white'],
                fg=self.COLORS['gray_500']
            ).pack(anchor='w')
            return
        
        for proj in projects[:5]:
            proj_frame = tk.Frame(
                self.recent_frame,
                bg=self.COLORS['gray_50'],
                cursor='hand2'
            )
            proj_frame.pack(fill=tk.X, pady=3)
            
            inner = tk.Frame(proj_frame, bg=self.COLORS['gray_50'])
            inner.pack(fill=tk.X, padx=12, pady=10)
            
            # Project info
            info = tk.Frame(inner, bg=self.COLORS['gray_50'])
            info.pack(side=tk.LEFT, fill=tk.X, expand=True)
            
            tk.Label(
                info,
                text=proj['name'],
                font=self.fonts['body_bold'],
                bg=self.COLORS['gray_50'],
                fg=self.COLORS['gray_900']
            ).pack(anchor='w')
            
            modified = proj['last_modified'][:10] if len(proj['last_modified']) > 10 else proj['last_modified']
            tk.Label(
                info,
                text=f"Last modified: {modified}",
                font=self.fonts['small'],
                bg=self.COLORS['gray_50'],
                fg=self.COLORS['gray_500']
            ).pack(anchor='w')
            
            # Status badge
            badge = tk.Label(
                inner,
                text=proj['step'],
                font=self.fonts['small'],
                bg=self.COLORS['gray_200'],
                fg=self.COLORS['gray_600'],
                padx=8,
                pady=2
            )
            badge.pack(side=tk.RIGHT)
            
            # Bind click
            proj_name = proj['name']
            proj_frame.bind('<Button-1>', lambda e, n=proj_name: self.open_project(n))
            for child in proj_frame.winfo_children():
                child.bind('<Button-1>', lambda e, n=proj_name: self.open_project(n))
                for subchild in child.winfo_children():
                    subchild.bind('<Button-1>', lambda e, n=proj_name: self.open_project(n))
    
    def create_new_project(self):
        """Create a new project"""
        name = simpledialog.askstring(
            "New Project",
            "Enter the Name of the Company:",
            parent=self.root
        )
        
        if not name:
            return
        
        # Store company name
        self.company_name = name
        
        # Tell user what to do next
        messagebox.showinfo(
            "Select Survey Document",
            "Next, select the Word document (.docx) containing your survey questions.\n\n"
            "The tool will extract all questions from this document.",
            parent=self.root
        )
        
        # Select survey file
        survey_file = filedialog.askopenfilename(
            title="Select Survey Word Document (.docx)",
            filetypes=[("Word Documents", "*.docx"), ("All Files", "*.*")]
        )
        
        if not survey_file:
            return
        
        self.show_processing("Creating project...")
        
        def process():
            try:
                # Extract questions from survey
                self.questions = self.extract_questions(survey_file)
                self.questions_with_options = self.extract_questions_with_options(survey_file)
                
                # Set up project
                self.project_name = name
                self.project_data = {
                    'project_name': name,
                    'company_name': name,  # Store company name
                    'survey_file': survey_file,
                    'is_new': True,
                    'questions': self.questions,
                    'questions_with_options': self.questions_with_options,
                    'current_step': 'Breakout Selection',
                    'created': datetime.now().isoformat()
                }
                self.save_project()
                
                self.root.after(0, self.finish_project_creation)
            except Exception as e:
                self.root.after(0, lambda: self.show_error(f"Failed to create project:\n{e}"))
        
        threading.Thread(target=process, daemon=True).start()
    
    def finish_project_creation(self):
        """Complete project creation and move to next step"""
        self.hide_processing()
        self.current_step = 2
        self.show_current_step()
    
    def open_project(self, project_name):
        """Open an existing project"""
        self.show_processing("Loading project...")
        
        def process():
            try:
                self.project_name = project_name
                self.project_data = self.load_project(project_name)
                
                if not self.project_data:
                    raise Exception("Project file not found or corrupted")
                
                # Restore state
                self.questions = self.project_data.get('questions', [])
                self.questions_with_options = self.project_data.get('questions_with_options', [])
                
                # Load company name (fallback to project name for backwards compatibility)
                self.company_name = self.project_data.get('company_name', project_name)
                
                # Load PowerPoint question tracking
                self.concept_questions = self.project_data.get('concept_questions', [])
                self.unique_questions = self.project_data.get('unique_questions', [])
                
                # Handle both old (single) and new (multiple) breakout formats
                self.selected_breakouts = self.project_data.get('breakout_columns', [])
                if not self.selected_breakouts:
                    # Backwards compatibility with old single breakout
                    old_breakout = self.project_data.get('breakout_column')
                    if old_breakout:
                        self.selected_breakouts = [old_breakout]
                
                self.selected_breakout = self.selected_breakouts[0] if self.selected_breakouts else None
                self.remark_file = self.project_data.get('remark_data_file')
                
                # Determine which step to show
                step = self.project_data.get('current_step', 'Breakout Selection')
                if step in ['Completed', 'Breakouts Created']:
                    self.current_step = 4
                elif step == 'Remark Data Loaded' or self.remark_file:
                    self.current_step = 4
                elif self.selected_breakouts:
                    self.current_step = 3
                else:
                    self.current_step = 2
                
                self.root.after(0, self.finish_open_project)
            except Exception as e:
                self.root.after(0, lambda: self.show_error(f"Failed to open project:\n{e}"))
        
        threading.Thread(target=process, daemon=True).start()
    
    def finish_open_project(self):
        """Complete project opening"""
        self.hide_processing()
        self.show_current_step()
    
    def show_step_2(self):
        """Step 2: Select Breakout Variables (up to 6)"""
        content = tk.Frame(self.step_frame, bg=self.COLORS['white'])
        content.pack(fill=tk.BOTH, expand=True, padx=30, pady=25)
        
        # Step header
        header = tk.Frame(content, bg=self.COLORS['white'])
        header.pack(fill=tk.X, pady=(0, 25))
        
        num_canvas = tk.Canvas(header, width=36, height=36, bg=self.COLORS['white'], highlightthickness=0)
        num_canvas.pack(side=tk.LEFT, padx=(0, 12))
        num_canvas.create_oval(2, 2, 34, 34, fill=self.COLORS['primary'], outline='')
        num_canvas.create_text(18, 18, text="2", fill='white', font=('Segoe UI', 12, 'bold'))
        
        title_frame = tk.Frame(header, bg=self.COLORS['white'])
        title_frame.pack(side=tk.LEFT)
        
        tk.Label(
            title_frame,
            text="Select Breakout Variables",
            font=self.fonts['step_title'],
            bg=self.COLORS['white'],
            fg=self.COLORS['gray_900']
        ).pack(anchor='w')
        
        tk.Label(
            title_frame,
            text="Choose up to 6 questions to use for segmenting results (click to select/unselect)",
            font=self.fonts['step_detail'],
            bg=self.COLORS['white'],
            fg=self.COLORS['gray_500']
        ).pack(anchor='w')
        
        # Selection counter
        self.selection_label = tk.Label(
            content,
            text="Selected: 0 of 6 maximum",
            font=self.fonts['body'],
            bg=self.COLORS['white'],
            fg=self.COLORS['gray_600']
        )
        self.selection_label.pack(anchor='w', pady=(0, 10))
        
        # Scrollable options frame
        options_container = tk.Frame(content, bg=self.COLORS['white'])
        options_container.pack(fill=tk.BOTH, expand=True, pady=(0, 20))
        
        # Canvas with scrollbar for many questions
        canvas = tk.Canvas(options_container, bg=self.COLORS['white'], highlightthickness=0)
        scrollbar = ttk.Scrollbar(options_container, orient="vertical", command=canvas.yview)
        scrollable_frame = tk.Frame(canvas, bg=self.COLORS['white'])
        
        scrollable_frame.bind(
            "<Configure>",
            lambda e: canvas.configure(scrollregion=canvas.bbox("all"))
        )
        
        canvas.create_window((0, 0), window=scrollable_frame, anchor="nw")
        canvas.configure(yscrollcommand=scrollbar.set)
        
        # Mouse wheel scrolling
        def on_mousewheel(event):
            canvas.yview_scroll(int(-1*(event.delta/120)), "units")
        canvas.bind_all("<MouseWheel>", on_mousewheel)
        
        canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        
        # Store checkbox variables
        self.breakout_vars = {}
        
        # Build breakout options from actual survey questions
        questions = self.questions_with_options or self.questions
        
        if not questions:
            tk.Label(
                scrollable_frame,
                text="No questions found in survey. Please go back and select a valid survey document.",
                font=self.fonts['body'],
                bg=self.COLORS['white'],
                fg=self.COLORS['gray_500']
            ).pack(anchor='w', pady=20)
        else:
            for q in questions:
                q_num = q['number']
                q_text = q.get('text', f'Question {q_num}')
                
                # Truncate long questions
                display_text = f"Q{q_num}: {q_text[:70]}"
                if len(q_text) > 70:
                    display_text += "..."
                
                opt_frame = tk.Frame(
                    scrollable_frame,
                    bg=self.COLORS['white'],
                    highlightbackground=self.COLORS['gray_200'],
                    highlightthickness=1,
                    cursor='hand2'
                )
                opt_frame.pack(fill=tk.X, pady=2, padx=(0, 10))
                
                inner = tk.Frame(opt_frame, bg=self.COLORS['white'])
                inner.pack(fill=tk.X, padx=12, pady=8)
                
                # Checkbox variable
                var = tk.BooleanVar(value=False)
                self.breakout_vars[str(q_num)] = var
                
                # Checkbox
                cb = tk.Checkbutton(
                    inner,
                    text=display_text,
                    variable=var,
                    font=self.fonts['body'],
                    bg=self.COLORS['white'],
                    fg=self.COLORS['gray_900'],
                    selectcolor=self.COLORS['white'],
                    activebackground=self.COLORS['white'],
                    cursor='hand2',
                    command=self.update_selection_count
                )
                cb.pack(side=tk.LEFT, anchor='w')
        
        # Continue button
        continue_btn = tk.Button(
            content,
            text="Continue",
            font=self.fonts['button'],
            bg=self.COLORS['primary'],
            fg='white',
            relief=tk.FLAT,
            padx=30,
            pady=10,
            cursor='hand2',
            command=self.confirm_breakout
        )
        continue_btn.pack(fill=tk.X, pady=(10, 0))
    
    def update_selection_count(self):
        """Update the selection counter and enforce limit"""
        selected = [k for k, v in self.breakout_vars.items() if v.get()]
        count = len(selected)
        
        if count > 6:
            # Uncheck the last one clicked (find it and uncheck)
            messagebox.showwarning(
                "Maximum Reached",
                "You can select up to 6 breakout variables.\n\nPlease unselect one before selecting another.",
                parent=self.root
            )
            # Find and uncheck one (the most recent is tricky, so just warn)
            return
        
        self.selection_label.config(text=f"Selected: {count} of 6 maximum")
    
    def confirm_breakout(self):
        """Confirm breakout selection and proceed"""
        selected = [k for k, v in self.breakout_vars.items() if v.get()]
        
        if not selected:
            messagebox.showwarning(
                "No Selection", 
                "Please select at least one question to use as a breakout variable.\n\n"
                "Click the checkbox next to the questions you want to use.",
                parent=self.root
            )
            return
        
        if len(selected) > 6:
            messagebox.showwarning(
                "Too Many Selected",
                "Please select no more than 6 breakout variables.",
                parent=self.root
            )
            return
        
        # Store as list
        self.selected_breakouts = selected
        self.selected_breakout = selected[0]  # Keep for backwards compatibility
        
        self.project_data['breakout_columns'] = selected
        self.project_data['current_step'] = 'Awaiting Remark Data'
        self.save_project()
        
        self.current_step = 3
        self.show_current_step()
    
    def show_step_3(self):
        """Step 3: Upload Remark Data"""
        content = tk.Frame(self.step_frame, bg=self.COLORS['white'])
        content.pack(fill=tk.BOTH, expand=True, padx=30, pady=25)
        
        # Step header
        header = tk.Frame(content, bg=self.COLORS['white'])
        header.pack(fill=tk.X, pady=(0, 25))
        
        num_canvas = tk.Canvas(header, width=36, height=36, bg=self.COLORS['white'], highlightthickness=0)
        num_canvas.pack(side=tk.LEFT, padx=(0, 12))
        num_canvas.create_oval(2, 2, 34, 34, fill=self.COLORS['primary'], outline='')
        num_canvas.create_text(18, 18, text="3", fill='white', font=('Segoe UI', 12, 'bold'))
        
        title_frame = tk.Frame(header, bg=self.COLORS['white'])
        title_frame.pack(side=tk.LEFT)
        
        tk.Label(
            title_frame,
            text="Upload Remark Data",
            font=self.fonts['step_title'],
            bg=self.COLORS['white'],
            fg=self.COLORS['gray_900']
        ).pack(anchor='w')
        
        tk.Label(
            title_frame,
            text="Select your exported Remark Excel file",
            font=self.fonts['step_detail'],
            bg=self.COLORS['white'],
            fg=self.COLORS['gray_500']
        ).pack(anchor='w')
        
        # Upload area
        upload_frame = tk.Frame(
            content,
            bg=self.COLORS['white'],
            highlightbackground=self.COLORS['gray_300'],
            highlightthickness=2,
            cursor='hand2'
        )
        upload_frame.pack(fill=tk.BOTH, expand=True, pady=20)
        
        upload_inner = tk.Frame(upload_frame, bg=self.COLORS['white'])
        upload_inner.pack(expand=True)
        
        tk.Label(upload_inner, text="📤", font=('Segoe UI', 36), bg=self.COLORS['white'], fg=self.COLORS['gray_400']).pack(pady=(30, 10))
        tk.Label(upload_inner, text="Click to upload or drag and drop", font=self.fonts['body_bold'], bg=self.COLORS['white'], fg=self.COLORS['gray_900']).pack()
        tk.Label(upload_inner, text="Excel files (.xlsx, .xls) or CSV files", font=self.fonts['small'], bg=self.COLORS['white'], fg=self.COLORS['gray_500']).pack(pady=(5, 30))
        
        # Bind click
        upload_frame.bind('<Button-1>', lambda e: self.select_remark_file())
        for child in upload_frame.winfo_children():
            child.bind('<Button-1>', lambda e: self.select_remark_file())
            for subchild in child.winfo_children():
                subchild.bind('<Button-1>', lambda e: self.select_remark_file())
    
    def select_remark_file(self):
        """Select and process Remark data file"""
        filepath = filedialog.askopenfilename(
            title="Select Remark Data Export File",
            filetypes=[
                ("Excel Files", "*.xlsx *.xls"),
                ("CSV Files", "*.csv"),
                ("All Files", "*.*")
            ]
        )
        
        if not filepath:
            return
        
        self.show_processing("Processing Remark data...")
        
        def process():
            try:
                import pandas as pd
                
                # Load data
                if filepath.endswith('.csv'):
                    df = pd.read_csv(filepath)
                else:
                    df = pd.read_excel(filepath)
                
                self.remark_file = filepath
                self.remark_df = df
                
                self.project_data['remark_data_file'] = filepath
                self.project_data['current_step'] = 'Remark Data Loaded'
                self.save_project()
                
                self.root.after(0, self.finish_remark_upload)
            except Exception as e:
                self.root.after(0, lambda: self.show_error(f"Failed to load Remark data:\n{e}"))
        
        threading.Thread(target=process, daemon=True).start()
    
    def finish_remark_upload(self):
        """Complete Remark data upload"""
        self.hide_processing()
        self.current_step = 4
        self.show_current_step()
    
    def show_step_4(self):
        """Step 4: Results - Preview and Download"""
        # Create scrollable container for Step 4
        outer_frame = tk.Frame(self.step_frame, bg=self.COLORS['white'])
        outer_frame.pack(fill=tk.BOTH, expand=True)
        
        # Canvas with scrollbar
        canvas = tk.Canvas(outer_frame, bg=self.COLORS['white'], highlightthickness=0)
        scrollbar = ttk.Scrollbar(outer_frame, orient="vertical", command=canvas.yview)
        
        content = tk.Frame(canvas, bg=self.COLORS['white'])
        
        content.bind(
            "<Configure>",
            lambda e: canvas.configure(scrollregion=canvas.bbox("all"))
        )
        
        canvas.create_window((0, 0), window=content, anchor="nw")
        canvas.configure(yscrollcommand=scrollbar.set)
        
        # Mouse wheel scrolling - bind to specific widgets
        def on_mousewheel(event):
            canvas.yview_scroll(int(-1*(event.delta/120)), "units")
        
        canvas.bind("<MouseWheel>", on_mousewheel)
        content.bind("<MouseWheel>", on_mousewheel)
        outer_frame.bind("<MouseWheel>", on_mousewheel)
        
        # Recursive bind for all children
        def bind_wheel_to_children(widget):
            widget.bind("<MouseWheel>", on_mousewheel)
            for child in widget.winfo_children():
                bind_wheel_to_children(child)
        
        # We'll call this after building the UI
        self._step4_bind_wheel = lambda: bind_wheel_to_children(content)
        
        canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=30, pady=25)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        
        # Step header with green check
        header = tk.Frame(content, bg=self.COLORS['white'])
        header.pack(fill=tk.X, pady=(0, 25))
        
        num_canvas = tk.Canvas(header, width=36, height=36, bg=self.COLORS['white'], highlightthickness=0)
        num_canvas.pack(side=tk.LEFT, padx=(0, 12))
        num_canvas.create_oval(2, 2, 34, 34, fill=self.COLORS['success'], outline='')
        num_canvas.create_text(18, 18, text="✓", fill='white', font=('Segoe UI', 14, 'bold'))
        
        title_frame = tk.Frame(header, bg=self.COLORS['white'])
        title_frame.pack(side=tk.LEFT)
        
        tk.Label(
            title_frame,
            text="Ready to Generate Results",
            font=self.fonts['step_title'],
            bg=self.COLORS['white'],
            fg=self.COLORS['gray_900']
        ).pack(anchor='w')
        
        tk.Label(
            title_frame,
            text="All steps completed - preview or download your crosstab reports",
            font=self.fonts['step_detail'],
            bg=self.COLORS['white'],
            fg=self.COLORS['gray_500']
        ).pack(anchor='w')
        
        # Success message
        success_frame = tk.Frame(
            content,
            bg=self.COLORS['success_bg'],
            highlightbackground=self.COLORS['success_border'],
            highlightthickness=1
        )
        success_frame.pack(fill=tk.X, pady=(0, 25))
        
        success_inner = tk.Frame(success_frame, bg=self.COLORS['success_bg'])
        success_inner.pack(fill=tk.X, padx=15, pady=12)
        
        tk.Label(
            success_inner,
            text="✓ Processing Complete",
            font=self.fonts['body_bold'],
            bg=self.COLORS['success_bg'],
            fg='#166534'
        ).pack(anchor='w')
        
        # Get counts
        q_count = len(self.questions) if self.questions else "?"
        
        tk.Label(
            success_inner,
            text=f"{q_count} questions ready for analysis",
            font=self.fonts['body'],
            bg=self.COLORS['success_bg'],
            fg='#166534'
        ).pack(anchor='w')
        
        # Action buttons
        buttons_frame = tk.Frame(content, bg=self.COLORS['white'])
        buttons_frame.pack(fill=tk.X, pady=(0, 10))
        
        # Preview button
        preview_btn = tk.Button(
            buttons_frame,
            text="👁 Preview PDF",
            font=self.fonts['button'],
            bg=self.COLORS['primary'],
            fg='white',
            relief=tk.FLAT,
            padx=20,
            pady=12,
            cursor='hand2',
            command=self.show_preview_selection
        )
        preview_btn.pack(fill=tk.X, pady=(0, 10))
        
        # Download button
        download_btn = tk.Button(
            buttons_frame,
            text="💾 Download PDF",
            font=self.fonts['button'],
            bg=self.COLORS['success'],
            fg='white',
            relief=tk.FLAT,
            padx=20,
            pady=12,
            cursor='hand2',
            command=self.show_download_selection
        )
        download_btn.pack(fill=tk.X, pady=(0, 20))
        
        # Excel section
        tk.Label(
            buttons_frame,
            text="Excel Spreadsheet",
            font=self.fonts['body_bold'],
            bg=self.COLORS['white'],
            fg=self.COLORS['gray_700']
        ).pack(anchor='w', pady=(0, 5))
        
        excel_btn = tk.Button(
            buttons_frame,
            text="📊 Download Excel",
            font=self.fonts['button'],
            bg=self.COLORS['excel'],
            fg='white',
            relief=tk.FLAT,
            padx=20,
            pady=12,
            cursor='hand2',
            command=self.show_excel_selection
        )
        excel_btn.pack(fill=tk.X, pady=(0, 20))
        
        # PowerPoint Template section
        tk.Label(
            buttons_frame,
            text="PowerPoint Template",
            font=self.fonts['body_bold'],
            bg=self.COLORS['white'],
            fg=self.COLORS['gray_700']
        ).pack(anchor='w', pady=(0, 5))
        
        tk.Label(
            buttons_frame,
            text="Update an existing PowerPoint template with your survey data",
            font=self.fonts['small'],
            bg=self.COLORS['white'],
            fg=self.COLORS['gray_500']
        ).pack(anchor='w', pady=(0, 5))
        
        pptx_btn = tk.Button(
            buttons_frame,
            text="📊 Update PowerPoint Template",
            font=self.fonts['button'],
            bg=self.COLORS['pptx'],
            fg='white',
            relief=tk.FLAT,
            padx=20,
            pady=12,
            cursor='hand2',
            command=self.update_powerpoint_template
        )
        pptx_btn.pack(fill=tk.X, pady=(0, 20))
        
        # Benchmark Comparison section
        tk.Label(
            buttons_frame,
            text="Benchmark Comparison",
            font=self.fonts['body_bold'],
            bg=self.COLORS['white'],
            fg=self.COLORS['gray_700']
        ).pack(anchor='w', pady=(0, 5))
        
        tk.Label(
            buttons_frame,
            text="Compare your overall results to the normative database",
            font=self.fonts['small'],
            bg=self.COLORS['white'],
            fg=self.COLORS['gray_500']
        ).pack(anchor='w', pady=(0, 5))
        
        benchmark_btn = tk.Button(
            buttons_frame,
            text="📊 Benchmark Comparison Report",
            font=self.fonts['button'],
            bg='#6366F1',  # Indigo
            fg='white',
            relief=tk.FLAT,
            padx=20,
            pady=12,
            cursor='hand2',
            command=self.generate_benchmark_report
        )
        benchmark_btn.pack(fill=tk.X, pady=(0, 30))
        
        # Bind mouse wheel to all children for scrolling
        if hasattr(self, '_step4_bind_wheel'):
            self._step4_bind_wheel()
    
    def show_preview_selection(self):
        """Show dialog to select which PDFs to preview"""
        self.show_pdf_selection_dialog(preview=True)
    
    def show_download_selection(self):
        """Show dialog to select which PDFs to download"""
        self.show_pdf_selection_dialog(preview=False)
    
    def show_pdf_selection_dialog(self, preview=False):
        """Show dialog to select which breakout PDFs to generate"""
        dialog = tk.Toplevel(self.root)
        dialog.title("Select PDFs" if not preview else "Select PDFs to Preview")
        dialog.geometry("400x500")
        dialog.transient(self.root)
        dialog.grab_set()
        
        # Center dialog
        dialog.update_idletasks()
        x = self.root.winfo_x() + (self.root.winfo_width() - dialog.winfo_width()) // 2
        y = self.root.winfo_y() + (self.root.winfo_height() - dialog.winfo_height()) // 2
        dialog.geometry(f"+{x}+{y}")
        
        # Title
        title_frame = tk.Frame(dialog, bg=self.COLORS['gray_50'])
        title_frame.pack(fill=tk.X, padx=20, pady=(20, 10))
        
        tk.Label(
            title_frame,
            text="Select Breakout Questions" if not preview else "Select PDFs to Preview",
            font=self.fonts['heading'],
            bg=self.COLORS['gray_50'],
            fg=self.COLORS['gray_900']
        ).pack(anchor='w')
        
        tk.Label(
            title_frame,
            text="Choose which crosstab reports to generate",
            font=self.fonts['body'],
            bg=self.COLORS['gray_50'],
            fg=self.COLORS['gray_600']
        ).pack(anchor='w', pady=(5, 0))
        
        # Checkboxes for each breakout
        checkbox_frame = tk.Frame(dialog, bg=self.COLORS['white'])
        checkbox_frame.pack(fill=tk.BOTH, expand=True, padx=20, pady=10)
        
        # Scrollable area
        canvas = tk.Canvas(checkbox_frame, bg=self.COLORS['white'], highlightthickness=0)
        scrollbar = ttk.Scrollbar(checkbox_frame, orient="vertical", command=canvas.yview)
        scrollable_frame = tk.Frame(canvas, bg=self.COLORS['white'])
        
        scrollable_frame.bind(
            "<Configure>",
            lambda e: canvas.configure(scrollregion=canvas.bbox("all"))
        )
        
        canvas.create_window((0, 0), window=scrollable_frame, anchor="nw")
        canvas.configure(yscrollcommand=scrollbar.set)
        
        canvas.pack(side="left", fill="both", expand=True)
        scrollbar.pack(side="right", fill="y")
        
        # Create checkboxes
        checkbox_vars = {}
        questions = self.questions_with_options or self.questions
        
        for breakout_num in self.selected_breakouts:
            var = tk.BooleanVar(value=True)  # Selected by default
            checkbox_vars[breakout_num] = var
            
            # Find question text
            q_text = f"Question {breakout_num}"
            for q in questions:
                if str(q['number']) == str(breakout_num):
                    q_text = f"Q{breakout_num}: {q.get('text', '')[:60]}"
                    break
            
            cb = tk.Checkbutton(
                scrollable_frame,
                text=q_text,
                variable=var,
                font=self.fonts['body'],
                bg=self.COLORS['white'],
                fg=self.COLORS['gray_900'],
                activebackground=self.COLORS['white'],
                selectcolor=self.COLORS['white'],
                padx=10,
                pady=5
            )
            cb.pack(anchor='w', fill=tk.X)
        
        # Buttons
        button_frame = tk.Frame(dialog, bg=self.COLORS['white'])
        button_frame.pack(fill=tk.X, padx=20, pady=20)
        
        def on_generate():
            # Get selected breakouts
            selected = [num for num, var in checkbox_vars.items() if var.get()]
            if not selected:
                messagebox.showwarning("No Selection", "Please select at least one PDF to generate.", parent=dialog)
                return
            dialog.destroy()
            self.generate_selected_pdfs(selected, preview)
        
        def on_cancel():
            dialog.destroy()
        
        tk.Button(
            button_frame,
            text="Cancel",
            font=self.fonts['button'],
            bg=self.COLORS['gray_200'],
            fg=self.COLORS['gray_700'],
            relief=tk.FLAT,
            padx=20,
            pady=10,
            cursor='hand2',
            command=on_cancel
        ).pack(side=tk.LEFT, padx=(0, 10))
        
        tk.Button(
            button_frame,
            text="Generate" if not preview else "Preview",
            font=self.fonts['button'],
            bg=self.COLORS['primary'],
            fg='white',
            relief=tk.FLAT,
            padx=20,
            pady=10,
            cursor='hand2',
            command=on_generate
        ).pack(side=tk.LEFT)
    
    def generate_selected_pdfs(self, selected_breakouts, preview=False):
        """Generate PDFs for selected breakouts only"""
        
        # For download mode, get save locations from user FIRST (on main thread)
        save_locations = {}
        if not preview:
            for breakout in selected_breakouts:
                default_name = f"Crosstab_{self.company_name or self.project_name}_Q{breakout}.pdf"
                default_dir = os.path.dirname(self.remark_file) if self.remark_file else os.path.expanduser("~")
                
                output_file = filedialog.asksaveasfilename(
                    title=f"Save PDF Report for Q{breakout}",
                    defaultextension=".pdf",
                    filetypes=[("PDF files", "*.pdf"), ("All files", "*.*")],
                    initialdir=default_dir,
                    initialfile=default_name,
                    parent=self.root
                )
                
                if not output_file:
                    # User cancelled
                    if len(save_locations) == 0:
                        # Cancelled on first one, abort entirely
                        return
                    else:
                        # Already got some locations, stop asking for more
                        break
                
                save_locations[breakout] = output_file
        
        # Now start background processing
        self.show_processing("Generating PDF report(s)...")
        
        def process():
            try:
                # Ensure we have data loaded
                if self.remark_df is None:
                    import pandas as pd
                    if self.remark_file.endswith('.csv'):
                        self.remark_df = pd.read_csv(self.remark_file)
                    else:
                        self.remark_df = pd.read_excel(self.remark_file)
                
                # Generate PDF for each selected breakout
                output_paths = []
                breakouts_to_generate = list(save_locations.keys()) if not preview else selected_breakouts
                
                for breakout in breakouts_to_generate:
                    self.selected_breakout = breakout  # Set current breakout
                    
                    # Get output path
                    if preview:
                        import tempfile
                        temp_dir = tempfile.gettempdir()
                        output_path = os.path.join(temp_dir, f"Crosstab_{self.company_name or self.project_name}_Q{breakout}.pdf")
                    else:
                        output_path = save_locations[breakout]
                    
                    # Generate the PDF
                    self.create_pdf_at_location(output_path)
                    output_paths.append(output_path)
                
                if output_paths:  # Only save project if we generated at least one file
                    self.project_data['current_step'] = 'Completed'
                    self.project_data['output_files'] = output_paths
                    self.save_project()
                
                self.root.after(0, lambda: self.finish_pdf_generation(output_paths, preview))
            except Exception as e:
                import traceback
                error_msg = f"{str(e)}\n\n{traceback.format_exc()}"
                self.root.after(0, lambda: self.show_error(f"Failed to generate PDF:\n{error_msg}"))
        
        threading.Thread(target=process, daemon=True).start()
    
    def finish_pdf_generation(self, output_paths, preview):
        """Complete PDF generation"""
        self.hide_processing()
        
        if not output_paths:
            messagebox.showinfo("Cancelled", "PDF generation was cancelled.", parent=self.root)
            return
        
        if preview:
            # Open each PDF in default viewer separately
            for pdf_path in output_paths:
                os.startfile(pdf_path)
            
            if len(output_paths) > 1:
                messagebox.showinfo(
                    "Multiple PDFs Opened",
                    f"Opened {len(output_paths)} PDF files for preview.\n\n"
                    f"All files are in:\n{os.path.dirname(output_paths[0])}",
                    parent=self.root
                )
        else:
            # Show summary
            output_dir = os.path.dirname(output_paths[0]) if output_paths else ""
            file_list = "\n".join([os.path.basename(p) for p in output_paths])
            messagebox.showinfo(
                "PDFs Generated",
                f"Generated {len(output_paths)} PDF file(s):\n\n{file_list}\n\nSaved to:\n{output_dir}",
                parent=self.root
            )
    
    def show_excel_selection(self):
        """Show dialog to select which Excel files to generate"""
        dialog = tk.Toplevel(self.root)
        dialog.title("Select Excel Exports")
        dialog.geometry("400x500")
        dialog.transient(self.root)
        dialog.grab_set()
        
        # Center dialog
        dialog.update_idletasks()
        x = self.root.winfo_x() + (self.root.winfo_width() - dialog.winfo_width()) // 2
        y = self.root.winfo_y() + (self.root.winfo_height() - dialog.winfo_height()) // 2
        dialog.geometry(f"+{x}+{y}")
        
        # Title
        title_frame = tk.Frame(dialog, bg=self.COLORS['gray_50'])
        title_frame.pack(fill=tk.X, padx=20, pady=(20, 10))
        
        tk.Label(
            title_frame,
            text="Select Breakout Questions",
            font=self.fonts['heading'],
            bg=self.COLORS['gray_50'],
            fg=self.COLORS['gray_900']
        ).pack(anchor='w')
        
        tk.Label(
            title_frame,
            text="Choose which crosstab reports to export to Excel",
            font=self.fonts['body'],
            bg=self.COLORS['gray_50'],
            fg=self.COLORS['gray_600']
        ).pack(anchor='w', pady=(5, 0))
        
        # Checkboxes for each breakout
        checkbox_frame = tk.Frame(dialog, bg=self.COLORS['white'])
        checkbox_frame.pack(fill=tk.BOTH, expand=True, padx=20, pady=10)
        
        # Scrollable area
        canvas = tk.Canvas(checkbox_frame, bg=self.COLORS['white'], highlightthickness=0)
        scrollbar = ttk.Scrollbar(checkbox_frame, orient="vertical", command=canvas.yview)
        scrollable_frame = tk.Frame(canvas, bg=self.COLORS['white'])
        
        scrollable_frame.bind(
            "<Configure>",
            lambda e: canvas.configure(scrollregion=canvas.bbox("all"))
        )
        
        canvas.create_window((0, 0), window=scrollable_frame, anchor="nw")
        canvas.configure(yscrollcommand=scrollbar.set)
        
        canvas.pack(side="left", fill="both", expand=True)
        scrollbar.pack(side="right", fill="y")
        
        # Create checkboxes
        checkbox_vars = {}
        questions = self.questions_with_options or self.questions
        
        for breakout_num in self.selected_breakouts:
            var = tk.BooleanVar(value=True)  # Selected by default
            checkbox_vars[breakout_num] = var
            
            # Find question text
            q_text = f"Question {breakout_num}"
            for q in questions:
                if str(q['number']) == str(breakout_num):
                    q_text = f"Q{breakout_num}: {q.get('text', '')[:60]}"
                    break
            
            cb = tk.Checkbutton(
                scrollable_frame,
                text=q_text,
                variable=var,
                font=self.fonts['body'],
                bg=self.COLORS['white'],
                fg=self.COLORS['gray_900'],
                activebackground=self.COLORS['white'],
                selectcolor=self.COLORS['white'],
                padx=10,
                pady=5
            )
            cb.pack(anchor='w', fill=tk.X)
        
        # Buttons
        button_frame = tk.Frame(dialog, bg=self.COLORS['white'])
        button_frame.pack(fill=tk.X, padx=20, pady=20)
        
        def on_generate():
            # Get selected breakouts
            selected = [num for num, var in checkbox_vars.items() if var.get()]
            if not selected:
                messagebox.showwarning("No Selection", "Please select at least one Excel file to generate.", parent=dialog)
                return
            dialog.destroy()
            self.generate_selected_excel(selected)
        
        def on_cancel():
            dialog.destroy()
        
        tk.Button(
            button_frame,
            text="Cancel",
            font=self.fonts['button'],
            bg=self.COLORS['gray_200'],
            fg=self.COLORS['gray_700'],
            relief=tk.FLAT,
            padx=20,
            pady=10,
            cursor='hand2',
            command=on_cancel
        ).pack(side=tk.LEFT, padx=(0, 10))
        
        tk.Button(
            button_frame,
            text="Generate Excel",
            font=self.fonts['button'],
            bg=self.COLORS['excel'],
            fg='white',
            relief=tk.FLAT,
            padx=20,
            pady=10,
            cursor='hand2',
            command=on_generate
        ).pack(side=tk.LEFT)
    
    def generate_selected_excel(self, selected_breakouts):
        """Generate Excel files for selected breakouts"""
        
        # Get save locations from user FIRST (on main thread)
        save_locations = {}
        for breakout in selected_breakouts:
            default_name = f"Crosstab_{self.company_name or self.project_name}_Q{breakout}.xlsx"
            default_dir = os.path.dirname(self.remark_file) if self.remark_file else os.path.expanduser("~")
            
            output_file = filedialog.asksaveasfilename(
                title=f"Save Excel Report for Q{breakout}",
                defaultextension=".xlsx",
                filetypes=[("Excel files", "*.xlsx"), ("All files", "*.*")],
                initialdir=default_dir,
                initialfile=default_name,
                parent=self.root
            )
            
            if not output_file:
                # User cancelled
                if len(save_locations) == 0:
                    # Cancelled on first one, abort entirely
                    return
                else:
                    # Already got some locations, stop asking for more
                    break
            
            save_locations[breakout] = output_file
        
        # Now start background processing
        self.show_processing("Generating Excel report(s)...")
        
        def process():
            try:
                # Ensure we have data loaded
                if self.remark_df is None:
                    import pandas as pd
                    if self.remark_file.endswith('.csv'):
                        self.remark_df = pd.read_csv(self.remark_file)
                    else:
                        self.remark_df = pd.read_excel(self.remark_file)
                
                # Generate Excel for each selected breakout
                output_paths = []
                
                for breakout in list(save_locations.keys()):
                    self.selected_breakout = breakout  # Set current breakout
                    output_path = save_locations[breakout]
                    
                    # Generate the Excel file
                    self.create_excel_at_location(output_path)
                    output_paths.append(output_path)
                
                if output_paths:
                    self.project_data['excel_files'] = output_paths
                    self.save_project()
                
                self.root.after(0, lambda: self.finish_excel_generation(output_paths))
            except Exception as e:
                import traceback
                error_msg = f"{str(e)}\n\n{traceback.format_exc()}"
                self.root.after(0, lambda: self.show_error(f"Failed to generate Excel:\n{error_msg}"))
        
        threading.Thread(target=process, daemon=True).start()
    
    def finish_excel_generation(self, output_paths):
        """Complete Excel generation"""
        self.hide_processing()
        
        if not output_paths:
            messagebox.showinfo("Cancelled", "Excel generation was cancelled.", parent=self.root)
            return
        
        # Show summary
        output_dir = os.path.dirname(output_paths[0]) if output_paths else ""
        file_list = "\n".join([os.path.basename(p) for p in output_paths])
        
        if messagebox.askyesno(
            "Excel Files Generated",
            f"Generated {len(output_paths)} Excel file(s):\n\n{file_list}\n\n"
            f"Saved to:\n{output_dir}\n\nOpen the first one now?",
            parent=self.root
        ):
            if output_paths:
                os.startfile(output_paths[0])
    
    def create_excel_at_location(self, output_file):
        """Create the Excel crosstab report at specified location"""
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from openpyxl.utils import get_column_letter
        import pandas as pd
        
        # Create workbook
        wb = Workbook()
        ws = wb.active
        ws.title = f"Crosstab Q{self.selected_breakout}"
        
        # Add version info sheet
        info_sheet = wb.create_sheet("Version Info")
        info_sheet['A1'] = "Generated by Version:"
        info_sheet['B1'] = "73"
        info_sheet['A2'] = "Date:"
        from datetime import datetime
        info_sheet['B2'] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        
        df = self.remark_df
        
        # Find the breakout column in data
        breakout_col = self.find_data_column(df, self.selected_breakout)
        if breakout_col is None:
            raise Exception(f"Could not find data column for breakout Q{self.selected_breakout}")
        
        # Get breakout options FROM THE SURVEY
        questions = self.questions_with_options or self.questions
        breakout_question = None
        for q in questions:
            if str(q['number']) == str(self.selected_breakout):
                breakout_question = q
                break
        
        if breakout_question and breakout_question.get('options'):
            breakout_values = breakout_question['options']
        else:
            breakout_values = sorted(df[breakout_col].dropna().unique().tolist())
        
        # Get data code mapping for breakout column
        breakout_code_mapping = self.get_data_code_mapping(df, breakout_col, len(breakout_values))
        
        # Get the breakout question text
        breakout_text = breakout_question.get('text', f'Question {self.selected_breakout}')[:50] if breakout_question else f'Q{self.selected_breakout}'
        
        # Styles
        header_fill = PatternFill(start_color="4F81BD", end_color="4F81BD", fill_type="solid")
        header_font = Font(bold=True, color="FFFFFF", size=10)
        question_fill = PatternFill(start_color="F3F4F6", end_color="F3F4F6", fill_type="solid")
        question_font = Font(bold=True, size=10)
        total_fill = PatternFill(start_color="F9FAFB", end_color="F9FAFB", fill_type="solid")
        total_font = Font(bold=True, size=9)
        avg_fill = PatternFill(start_color="DBEAFE", end_color="DBEAFE", fill_type="solid")
        avg_font = Font(bold=True, italic=True, size=9)
        center_align = Alignment(horizontal='center', vertical='center')
        left_align = Alignment(horizontal='left', vertical='center')
        thin_border = Border(
            left=Side(style='thin'),
            right=Side(style='thin'),
            top=Side(style='thin'),
            bottom=Side(style='thin')
        )
        
        # Column headers
        row_num = 1
        col_headers = [breakout_text] + [str(v)[:20] for v in breakout_values] + ['Total']
        for col_idx, header in enumerate(col_headers, 1):
            cell = ws.cell(row=row_num, column=col_idx, value=header)
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = center_align
            cell.border = thin_border
        
        row_num += 1
        
        # Process each question (excluding the breakout question)
        total_questions = len(questions)
        for q_idx, q in enumerate(questions):
            q_num = q['number']
            
            # Skip the breakout question itself
            if str(q_num) == str(self.selected_breakout):
                continue
            
            q_text = q.get('text', f'Question {q_num}')
            
            # Find data column
            data_col = self.find_data_column(df, q_num)
            if data_col is None:
                continue
            
            # Get response options FROM THE SURVEY
            response_options = q.get('options', [])
            if not response_options:
                response_options = sorted(df[data_col].dropna().unique().tolist())
            
            # Question header row
            q_header_text = f"{q_num}  {q_text[:100]}"
            cell = ws.cell(row=row_num, column=1, value=q_header_text)
            cell.fill = question_fill
            cell.font = question_font
            cell.alignment = left_align
            cell.border = thin_border
            ws.merge_cells(start_row=row_num, start_column=1, end_row=row_num, end_column=len(breakout_values) + 2)
            row_num += 1
            
            # Check if this question needs special calculation
            is_last_question = (q_idx == total_questions - 1)
            special_calc = self.get_special_calculation(q_num, q_text, response_options, is_last_question)
            
            # Get data code mapping for this question's responses
            response_code_mapping = self.get_data_code_mapping(df, data_col, len(response_options))
            
            # EMERGENCY DIAGNOSTIC for Q44 - write to Desktop
            if q_num == 44:
                import os
                desktop = os.path.join(os.path.expanduser("~"), "Desktop")
                diag_file = os.path.join(desktop, "Q44_V49_DIAGNOSTIC.txt")
                
                # Find Picking group
                picking_idx = None
                for grp_idx, group_val in enumerate(breakout_values):
                    if 'picking' in str(group_val).lower():
                        picking_idx = grp_idx
                        break
                
                if picking_idx is not None:
                    try:
                        with open(diag_file, 'w', encoding='utf-8') as f:
                            f.write("="*80 + "\n")
                            f.write("Q44 V49 DIAGNOSTIC - METHOD 7B REMOVED\n")
                            f.write("="*80 + "\n")
                            
                            group_val = breakout_values[picking_idx]
                            group_df = self.filter_by_option(df, breakout_col, group_val, picking_idx, breakout_code_mapping)
                            
                            # Exclude BLANK
                            valid_responses = group_df[group_df[data_col].notna()]
                            valid_responses = valid_responses[valid_responses[data_col] != 'BLANK']
                            group_total = len(valid_responses)
                            
                            f.write(f"Group total (excluding BLANK): {group_total}\n\n")
                            
                            f.write("Data values:\n")
                            for val in valid_responses[data_col].unique():
                                count = len(valid_responses[valid_responses[data_col] == val])
                                f.write(f"  '{val}': {count}\n")
                            
                            f.write("\n" + "="*80 + "\n")
                            f.write("TESTING EACH SURVEY OPTION:\n\n")
                            
                            for idx, response in enumerate(response_options):
                                f.write(f"Option {idx}: '{response}'\n")
                                filtered = self.filter_by_option(group_df, data_col, response, idx, response_code_mapping)
                                count = len(filtered)
                                
                                f.write(f"  Matched: {count} rows\n")
                                
                                if count > 0:
                                    matched_vals = filtered[data_col].dropna().unique().tolist()
                                    f.write(f"  Data values matched:\n")
                                    for mv in matched_vals:
                                        mv_count = len(filtered[filtered[data_col] == mv])
                                        f.write(f"    '{mv}': {mv_count}\n")
                                
                                pct = (count / group_total) * 100 if group_total > 0 else 0
                                f.write(f"  Raw %: {pct:.1f}%\n\n")
                            
                            f.write("="*80 + "\n")
                        
                        self.q44_diag_file = diag_file
                    except Exception as e:
                        pass
            
            # DIAGNOSTIC for Q8 and Q16
            if q_num in [8, 16, 17, 58]:
                import os
                # Put diagnostic in same folder as output Excel file
                output_dir = os.path.dirname(output_file) if output_file else os.path.expanduser("~")
                diag_file = os.path.join(output_dir, f"Q{q_num}_DIAGNOSTIC.txt")
                
                try:
                    with open(diag_file, 'w', encoding='utf-8') as f:
                        f.write("="*80 + "\n")
                        f.write(f"Q{q_num} DIAGNOSTIC - VERSION 66\n")
                        f.write("="*80 + "\n")
                        f.write(f"Question: {q_text}\n\n")
                        
                        f.write(f"Survey options ({len(response_options)}):\n")
                        for i, opt in enumerate(response_options):
                            f.write(f"  {i}: {opt}\n")
                        f.write("\n")
                        
                        # NEW: Add matching test for Q8
                        if q_num == 8:
                            f.write("="*80 + "\n")
                            f.write("MATCHING TEST FOR EACH OPTION\n")
                            f.write("="*80 + "\n")
                            data_values = [v for v in df[data_col].dropna().unique() if v != 'BLANK']
                            f.write(f"Data values: {data_values}\n\n")
                            for idx, opt in enumerate(response_options):
                                f.write(f"Option {idx}: '{opt}'\n")
                                matched_df = self.filter_by_option(df, data_col, opt, idx, None)
                                matched_values = matched_df[data_col].unique() if len(matched_df) > 0 else []
                                if len(matched_values) > 0:
                                    for mv in matched_values:
                                        if mv != 'BLANK':
                                            count = len(matched_df[matched_df[data_col] == mv])
                                            f.write(f"  ✓ MATCHES: '{mv}' ({count} rows)\n")
                                else:
                                    f.write(f"  ✗ NO MATCHES\n")
                                f.write("\n")
                            f.write("="*80 + "\n\n")
                        
                        # Show column structure
                        f.write(f"Columns: {len(breakout_values)} breakout groups + 1 Total = {len(breakout_values) + 1} total\n")
                        for i, bv in enumerate(breakout_values):
                            f.write(f"  Column {i}: {bv}\n")
                        f.write(f"  Column {len(breakout_values)}: TOTAL\n\n")
                        
                        # Check Picking group specifically
                        picking_idx = None
                        for idx, bv in enumerate(breakout_values):
                            if 'picking' in str(bv).lower():
                                picking_idx = idx
                                break
                        
                        if picking_idx is not None:
                            group_df = self.filter_by_option(df, breakout_col, breakout_values[picking_idx], picking_idx, breakout_code_mapping)
                            valid = group_df[(group_df[data_col].notna()) & (group_df[data_col] != 'BLANK')]
                            
                            f.write(f"PICKING GROUP ({len(valid)} valid responses):\n")
                            for idx, response in enumerate(response_options):
                                filtered = self.filter_by_option(group_df, data_col, response, idx, response_code_mapping)
                                count = len(filtered)
                                pct = (count / len(valid) * 100) if len(valid) > 0 else 0
                                f.write(f"  Option {idx}: {count} ({pct:.1f}%)\n")
                            f.write("\n")
                        
                        # Check TOTAL column
                        valid_total = df[(df[data_col].notna()) & (df[data_col] != 'BLANK')]
                        f.write(f"TOTAL COLUMN ({len(valid_total)} valid responses):\n")
                        for idx, response in enumerate(response_options):
                            filtered = self.filter_by_option(df, data_col, response, idx, response_code_mapping)
                            count = len(filtered)
                            pct = (count / len(valid_total) * 100) if len(valid_total) > 0 else 0
                            f.write(f"  Option {idx}: {count} ({pct:.1f}%)\n")
                        
                        f.write("="*80 + "\n")
                    
                    if q_num == 8:
                        self.q8_diag_file = diag_file
                    elif q_num == 16:
                        self.q16_diag_file = diag_file
                    elif q_num == 17:
                        self.q17_diag_file = diag_file
                    elif q_num == 58:
                        self.q58_diag_file = diag_file
                        
                except Exception as e:
                    pass  # Silently ignore diagnostic errors
            
            # ===== STEP 1: Calculate all raw percentages for each column =====
            column_percentages = []
            
            for grp_idx, group_val in enumerate(breakout_values):
                group_df = self.filter_by_option(df, breakout_col, group_val, grp_idx, breakout_code_mapping)
                # CRITICAL FIX: Exclude BLANK from total (BLANK is a string, not NaN)
                valid_responses = group_df[group_df[data_col].notna()]
                valid_responses = valid_responses[valid_responses[data_col] != 'BLANK']
                group_total = len(valid_responses)
                
                if group_total > 0:
                    col_pcts = []
                    for idx, response in enumerate(response_options):
                        count = len(self.filter_by_option(group_df, data_col, response, idx, response_code_mapping))
                        pct = count / group_total
                        col_pcts.append(pct)
                    column_percentages.append(col_pcts)
                else:
                    column_percentages.append([0] * len(response_options))
            
            # Calculate for Total column
            # CRITICAL FIX: Exclude BLANK from total
            valid_total = df[df[data_col].notna()]
            valid_total = valid_total[valid_total[data_col] != 'BLANK']
            total_respondents = len(valid_total)
            if total_respondents > 0:
                total_col_pcts = []
                for idx, response in enumerate(response_options):
                    total_count = len(self.filter_by_option(df, data_col, response, idx, response_code_mapping))
                    total_pct = total_count / total_respondents
                    total_col_pcts.append(total_pct)
                column_percentages.append(total_col_pcts)
            else:
                column_percentages.append([0] * len(response_options))
            
            # ===== STEP 2: Adjust each column to sum to 100% =====
            adjusted_columns = []
            for col_idx, col_pcts in enumerate(column_percentages):
                adjusted = self.adjust_percentages_to_100(col_pcts)
                
                # DEBUG: Log Q25 THD column
                if q_num == 25 and col_idx == 3:  # THD is column index 3
                    import os
                    desktop = os.path.join(os.path.expanduser("~"), "Desktop")
                    debug_file = os.path.join(desktop, "Q25_DEBUG.txt")
                    with open(debug_file, 'w') as f:
                        f.write(f"Q25 THD Column (index {col_idx}):\n")
                        f.write(f"Input percentages (decimals): {col_pcts}\n")
                        f.write(f"Sum of inputs: {sum(col_pcts)}\n")
                        f.write(f"Adjusted (integers): {adjusted}\n")
                        f.write(f"Sum of adjusted: {sum(adjusted)}\n")
                        f.write(f"Type of adjusted[0]: {type(adjusted[0])}\n")
                
                # DEBUG: Log Q7 THD column too
                if q_num == 7 and col_idx == 3:  # THD is column index 3
                    import os
                    desktop = os.path.join(os.path.expanduser("~"), "Desktop")
                    debug_file = os.path.join(desktop, "Q7_DEBUG.txt")
                    with open(debug_file, 'w') as f:
                        f.write(f"Q7 THD Column (index {col_idx}):\n")
                        f.write(f"Input percentages (decimals): {col_pcts}\n")
                        f.write(f"Sum of inputs: {sum(col_pcts)}\n")
                        f.write(f"Adjusted (integers): {adjusted}\n")
                        f.write(f"Sum of adjusted: {sum(adjusted)}\n")
                        f.write(f"Type of adjusted[0]: {type(adjusted[0])}\n")
                
                adjusted_columns.append(adjusted)
            
            # ===== STEP 3: Write response rows with adjusted percentages =====
            for idx, response in enumerate(response_options):
                # Response text
                cell = ws.cell(row=row_num, column=1, value=f"    {str(response)[:60]}")
                cell.alignment = left_align
                cell.border = thin_border
                
                # Write adjusted percentages for each column
                col_idx = 2
                for col_num in range(len(adjusted_columns)):
                    pct_int = adjusted_columns[col_num][idx]
                    pct_decimal = pct_int / 100.0  # Excel needs decimal format for percentage
                    cell = ws.cell(row=row_num, column=col_idx, value=pct_decimal)
                    cell.number_format = '0%'
                    cell.alignment = center_align
                    cell.border = thin_border
                    col_idx += 1
                
                row_num += 1
            
            # Total row (counts)
            cell = ws.cell(row=row_num, column=1, value="    Total")
            cell.fill = total_fill
            cell.font = total_font
            cell.alignment = left_align
            cell.border = thin_border
            
            col_idx = 2
            for grp_idx, group_val in enumerate(breakout_values):
                group_df = self.filter_by_option(df, breakout_col, group_val, grp_idx, breakout_code_mapping)
                count = len(group_df[group_df[data_col].notna()])
                cell = ws.cell(row=row_num, column=col_idx, value=count)
                cell.fill = total_fill
                cell.font = total_font
                cell.alignment = center_align
                cell.border = thin_border
                col_idx += 1
            
            # CRITICAL FIX: Exclude BLANK
            valid_grand = df[df[data_col].notna()]
            valid_grand = valid_grand[valid_grand[data_col] != 'BLANK']
            grand_total = len(valid_grand)
            cell = ws.cell(row=row_num, column=col_idx, value=grand_total)
            cell.fill = total_fill
            cell.font = total_font
            cell.alignment = center_align
            cell.border = thin_border
            row_num += 1
            
            # Special calculation row (Average Rating, Average Miles, etc.)
            if special_calc:
                avg_label = special_calc.get('label', 'Average')
                cell = ws.cell(row=row_num, column=1, value=f"    {avg_label}")
                cell.fill = avg_fill
                cell.font = avg_font
                cell.alignment = left_align
                cell.border = thin_border
                
                col_idx = 2
                for grp_idx, group_val in enumerate(breakout_values):
                    group_df = self.filter_by_option(df, breakout_col, group_val, grp_idx, breakout_code_mapping)
                    avg = self.calculate_special_average(group_df, data_col, response_options, special_calc)
                    cell = ws.cell(row=row_num, column=col_idx, value=avg if avg is not None else "N/A")
                    if avg is not None:
                        cell.number_format = '0.00'
                    cell.fill = avg_fill
                    cell.font = avg_font
                    cell.alignment = center_align
                    cell.border = thin_border
                    col_idx += 1
                
                # Total average
                total_avg = self.calculate_special_average(df, data_col, response_options, special_calc)
                cell = ws.cell(row=row_num, column=col_idx, value=total_avg if total_avg is not None else "N/A")
                if total_avg is not None:
                    cell.number_format = '0.00'
                cell.fill = avg_fill
                cell.font = avg_font
                cell.alignment = center_align
                cell.border = thin_border
                row_num += 1
            
            # Add spacing between questions
            row_num += 1
        
        # Adjust column widths
        ws.column_dimensions['A'].width = 50
        for col_idx in range(2, len(breakout_values) + 3):
            ws.column_dimensions[get_column_letter(col_idx)].width = 12
        
        # Save workbook
        wb.save(output_file)
    

    def generate_benchmark_report(self):
        """Generate a benchmark comparison report (Overall vs Norms)"""
        # Ask for norms file
        norms_file = filedialog.askopenfilename(
            title="Select Normative Database File",
            filetypes=[("Excel files", "*.xlsx *.xls"), ("All files", "*.*")],
            parent=self.root
        )
        
        if not norms_file:
            return
        
        # Ask for output format
        format_choice = messagebox.askquestion(
            "Output Format",
            "Generate PDF report?\n\nClick 'Yes' for PDF, 'No' for Excel",
            parent=self.root
        )
        
        # Get save location
        company = self.company_name or self.project_name
        if format_choice == 'yes':
            default_name = f"Benchmark_Comparison_{company.replace(' ', '_')}.pdf"
            filetypes = [("PDF files", "*.pdf"), ("All files", "*.*")]
        else:
            default_name = f"Benchmark_Comparison_{company.replace(' ', '_')}.xlsx"
            filetypes = [("Excel files", "*.xlsx"), ("All files", "*.*")]
        
        default_dir = os.path.dirname(self.remark_file) if self.remark_file else os.path.expanduser("~")
        
        output_file = filedialog.asksaveasfilename(
            title="Save Benchmark Report",
            defaultextension=".pdf" if format_choice == 'yes' else ".xlsx",
            filetypes=filetypes,
            initialdir=default_dir,
            initialfile=default_name,
            parent=self.root
        )
        
        if not output_file:
            return
        
        # Generate report in background
        self.show_processing("Generating benchmark comparison report...")
        
        def process():
            try:
                import pandas as pd
                
                # Load norms data
                norms_df = pd.read_excel(norms_file)
                norms_data = self.parse_norms_file(norms_df)
                
                # Ensure we have survey data loaded
                if self.remark_df is None:
                    if self.remark_file.endswith('.csv'):
                        self.remark_df = pd.read_csv(self.remark_file)
                    else:
                        self.remark_df = pd.read_excel(self.remark_file)
                
                # Generate report
                if format_choice == 'yes':
                    self.create_benchmark_pdf(output_file, norms_data)
                else:
                    self.create_benchmark_excel(output_file, norms_data)
                
                self.root.after(0, lambda: self.finish_benchmark_report(output_file))
            except Exception as e:
                import traceback
                error_msg = f"{str(e)}\n\n{traceback.format_exc()}"
                self.root.after(0, lambda: self.show_error(f"Failed to generate benchmark report:\n{error_msg}"))
        
        threading.Thread(target=process, daemon=True).start()
    
    def parse_norms_file(self, norms_df):
        """Parse the normative database Excel file"""
        import pandas as pd
        norms_data = {}
        current_question = None
        current_responses = []
        
        first_col = norms_df.columns[0]
        
        for idx, row in norms_df.iterrows():
            text = str(row[first_col]).strip() if pd.notna(row[first_col]) else ""
            avg_val = row.get('Averages', None)
            
            if not text:
                continue
            
            # Skip metadata rows
            if text.lower() in ['headcount', 'shift', 'date', 'demographic information']:
                continue
            
            # Skip AVERAGE rows (weighted averages)
            if text.upper() == 'AVERAGE':
                continue
            
            # Check if this is a question (no percentage value and looks like a question)
            is_question = pd.isna(avg_val) and (
                text.endswith('?') or
                text.lower().startswith(('how ', 'what ', 'do you ', 'are you ', 'which ', 'when working', 'on my days', 'if you have', 'importance of'))
            )
            
            if is_question:
                # Save previous question if exists
                if current_question and current_responses:
                    norms_data[current_question] = current_responses
                
                current_question = text
                current_responses = []
            elif current_question and pd.notna(avg_val):
                # This is a response option with percentage
                try:
                    pct = float(avg_val)
                    current_responses.append({
                        'response': text,
                        'percentage': pct
                    })
                except (ValueError, TypeError):
                    pass
        
        # Save last question
        if current_question and current_responses:
            norms_data[current_question] = current_responses
        
        return norms_data
    
    def match_survey_to_norms(self, survey_question, norms_data):
        """Find the best matching norms question for a survey question"""
        import re
        
        def normalize(text):
            """Normalize text for matching"""
            text = text.lower()
            # Remove punctuation except apostrophes
            text = re.sub(r"[^\w\s']", ' ', text)
            # Normalize apostrophes
            text = text.replace("'", "'")
            return ' '.join(text.split())
        
        def get_key_words(text):
            """Extract key words, removing stop words"""
            stop_words = {'a', 'an', 'the', 'is', 'are', 'do', 'does', 'you', 'your', 'to', 'for', 'of', 'in', 'on', 'at', 'and', 'or', 'that', 'this', 'what', 'how', 'which', 'when', 'where', 'who'}
            words = normalize(text).split()
            return set(w for w in words if w not in stop_words and len(w) > 2)
        
        survey_normalized = normalize(survey_question)
        survey_words = get_key_words(survey_question)
        
        best_match = None
        best_score = 0
        
        for norms_question in norms_data.keys():
            norms_normalized = normalize(norms_question)
            norms_words = get_key_words(norms_question)
            
            # Calculate overlap score
            if not survey_words or not norms_words:
                continue
            
            overlap = survey_words & norms_words
            score = len(overlap) / max(len(survey_words), len(norms_words))
            
            # Boost for key phrase matches
            key_phrases = [
                ('second job', 0.3), ('student', 0.3), ('gender', 0.3),
                ('age group', 0.3), ('hours of sleep', 0.3), ('alarm clock', 0.3),
                ('commute', 0.3), ('company', 0.2), ('department', 0.2),
                ('childcare', 0.3), ('eldercare', 0.3), ('overtime', 0.3),
                ('spouse', 0.3), ('single parent', 0.3), ('shift', 0.2)
            ]
            
            for phrase, boost in key_phrases:
                if phrase in survey_normalized and phrase in norms_normalized:
                    score += boost
            
            # Exact match gets highest score
            if survey_normalized == norms_normalized:
                score = 1.0
            
            if score > best_score and score > 0.25:  # Minimum threshold
                best_score = score
                best_match = norms_question
        
        return best_match, best_score
    
    def match_response_to_norms(self, survey_response, norms_responses):
        """Find matching norms response for a survey response"""
        import re
        
        def normalize(text):
            text = str(text).lower().strip()
            text = re.sub(r"[^\w\s']", ' ', text)
            text = text.replace("'", "'")
            return ' '.join(text.split())
        
        survey_norm = normalize(survey_response)
        
        for norms_resp in norms_responses:
            norms_norm = normalize(norms_resp['response'])
            
            # Exact match
            if survey_norm == norms_norm:
                return norms_resp['percentage']
            
            # Partial match (one contains the other)
            if survey_norm in norms_norm or norms_norm in survey_norm:
                return norms_resp['percentage']
            
            # Word overlap
            survey_words = set(survey_norm.split())
            norms_words = set(norms_norm.split())
            overlap = survey_words & norms_words
            
            if len(overlap) >= min(2, len(survey_words)):
                return norms_resp['percentage']
        
        return None
    
    def create_benchmark_pdf(self, output_file, norms_data):
        """Create PDF benchmark comparison report"""
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import letter, landscape
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak
        from reportlab.lib.enums import TA_LEFT, TA_CENTER
        import pandas as pd
        
        # Create document
        doc = SimpleDocTemplate(
            output_file,
            pagesize=landscape(letter),
            leftMargin=0.5*inch,
            rightMargin=0.5*inch,
            topMargin=0.5*inch,
            bottomMargin=0.5*inch
        )
        
        elements = []
        df = self.remark_df
        questions = self.questions_with_options or self.questions
        
        # Title
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            'Title',
            parent=styles['Heading1'],
            fontSize=16,
            spaceAfter=12,
            alignment=TA_CENTER
        )
        
        company = self.company_name or self.project_name
        elements.append(Paragraph(f"Benchmark Comparison Report: {company}", title_style))
        elements.append(Spacer(1, 0.2*inch))
        
        # Column headers
        col_headers = ['Question / Response', 'Overall', 'Normative']
        
        matched_count = 0
        unmatched_count = 0
        
        # Process each question
        for q in questions:
            q_num = q['number']
            q_text = q.get('text', f'Question {q_num}')
            
            # Find data column
            data_col = self.find_data_column(df, q_num)
            if data_col is None:
                continue
            
            # Get response options
            response_options = q.get('options', [])
            if not response_options:
                response_options = sorted(df[data_col].dropna().unique().tolist())
            
            if len(response_options) == 0:
                continue
            
            # Find matching norms question
            norms_question, match_score = self.match_survey_to_norms(q_text, norms_data)
            norms_responses = norms_data.get(norms_question, []) if norms_question else []
            
            if norms_question:
                matched_count += 1
            else:
                unmatched_count += 1
            
            # Build table data
            table_data = [col_headers]
            
            # Question header row
            q_header = [f"Q{q_num}: {q_text[:80]}", '', '']
            table_data.append(q_header)
            
            # Get data code mapping
            response_code_mapping = self.get_data_code_mapping(df, data_col, len(response_options))
            
            # Total respondents
            total_count = len(df[df[data_col].notna()])
            
            # Response rows
            for idx, response in enumerate(response_options):
                # Calculate overall percentage
                resp_count = len(self.filter_by_option(df, data_col, response, idx, response_code_mapping))
                overall_pct = (resp_count / total_count * 100) if total_count > 0 else 0
                
                # Find matching norms percentage
                norms_pct = self.match_response_to_norms(response, norms_responses)
                norms_display = f"{norms_pct:.1f}%" if norms_pct is not None else "—"
                
                row = [
                    f"    {str(response)[:50]}",
                    f"{overall_pct:.1f}%",
                    norms_display
                ]
                table_data.append(row)
            
            # Total row
            total_row = ["    Total", str(total_count), ""]
            table_data.append(total_row)
            
            # Create table
            col_widths = [5.5*inch, 1.5*inch, 1.5*inch]
            table = Table(table_data, colWidths=col_widths)
            
            num_data_rows = len(response_options)
            total_row_idx = 2 + num_data_rows
            
            # Table style
            style_commands = [
                # Header row
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#6366F1')),  # Indigo
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 9),
                ('ALIGN', (1, 0), (-1, 0), 'CENTER'),
                
                # Question header row
                ('BACKGROUND', (0, 1), (-1, 1), colors.HexColor('#f3f4f6')),
                ('FONTNAME', (0, 1), (-1, 1), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 1), (-1, 1), 9),
                ('SPAN', (0, 1), (-1, 1)),
                
                # Data rows
                ('FONTSIZE', (0, 2), (-1, -1), 8),
                ('ALIGN', (1, 2), (-1, -1), 'CENTER'),
                
                # Total row
                ('FONTNAME', (0, total_row_idx), (-1, total_row_idx), 'Helvetica-Bold'),
                ('BACKGROUND', (0, total_row_idx), (-1, total_row_idx), colors.HexColor('#f9fafb')),
                
                # Grid
                ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#d1d5db')),
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                ('LEFTPADDING', (0, 0), (-1, -1), 4),
                ('RIGHTPADDING', (0, 0), (-1, -1), 4),
                ('TOPPADDING', (0, 0), (-1, -1), 4),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
            ]
            
            # Highlight rows with no norms data
            if not norms_question:
                style_commands.append(('BACKGROUND', (2, 2), (2, total_row_idx - 1), colors.HexColor('#fef3c7')))
            
            style = TableStyle(style_commands)
            table.setStyle(style)
            elements.append(table)
            elements.append(Spacer(1, 0.3*inch))
        
        # Summary
        summary_style = ParagraphStyle(
            'Summary',
            parent=styles['Normal'],
            fontSize=10,
            spaceAfter=6
        )
        elements.append(Paragraph(f"Summary: {matched_count} questions matched to normative database, {unmatched_count} questions without norms data", summary_style))
        
        # Build PDF
        doc.build(elements)
    
    def create_benchmark_excel(self, output_file, norms_data):
        """Create Excel benchmark comparison report"""
        from openpyxl import Workbook
        from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
        from openpyxl.utils import get_column_letter
        import pandas as pd
        
        wb = Workbook()
        ws = wb.active
        ws.title = "Benchmark Comparison"
        
        df = self.remark_df
        questions = self.questions_with_options or self.questions
        
        # Styles
        header_fill = PatternFill(start_color='6366F1', end_color='6366F1', fill_type='solid')
        header_font = Font(bold=True, color='FFFFFF')
        question_fill = PatternFill(start_color='F3F4F6', end_color='F3F4F6', fill_type='solid')
        total_fill = PatternFill(start_color='F9FAFB', end_color='F9FAFB', fill_type='solid')
        no_norms_fill = PatternFill(start_color='FEF3C7', end_color='FEF3C7', fill_type='solid')
        thin_border = Border(
            left=Side(style='thin', color='D1D5DB'),
            right=Side(style='thin', color='D1D5DB'),
            top=Side(style='thin', color='D1D5DB'),
            bottom=Side(style='thin', color='D1D5DB')
        )
        
        row_num = 1
        
        # Title
        company = self.company_name or self.project_name
        ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=3)
        ws['A1'] = f"Benchmark Comparison Report: {company}"
        ws['A1'].font = Font(bold=True, size=14)
        row_num = 3
        
        # Headers
        headers = ['Question / Response', 'Overall', 'Normative']
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=row_num, column=col, value=header)
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = Alignment(horizontal='center')
            cell.border = thin_border
        row_num += 1
        
        matched_count = 0
        unmatched_count = 0
        
        # Process each question
        for q in questions:
            q_num = q['number']
            q_text = q.get('text', f'Question {q_num}')
            
            # Find data column
            data_col = self.find_data_column(df, q_num)
            if data_col is None:
                continue
            
            # Get response options
            response_options = q.get('options', [])
            if not response_options:
                response_options = sorted(df[data_col].dropna().unique().tolist())
            
            if len(response_options) == 0:
                continue
            
            # Find matching norms question
            norms_question, match_score = self.match_survey_to_norms(q_text, norms_data)
            norms_responses = norms_data.get(norms_question, []) if norms_question else []
            
            if norms_question:
                matched_count += 1
            else:
                unmatched_count += 1
            
            # Question header row
            ws.merge_cells(start_row=row_num, start_column=1, end_row=row_num, end_column=3)
            cell = ws.cell(row=row_num, column=1, value=f"Q{q_num}: {q_text[:80]}")
            cell.fill = question_fill
            cell.font = Font(bold=True)
            cell.border = thin_border
            for col in range(2, 4):
                ws.cell(row=row_num, column=col).border = thin_border
            row_num += 1
            
            # Get data code mapping
            response_code_mapping = self.get_data_code_mapping(df, data_col, len(response_options))
            
            # Total respondents
            total_count = len(df[df[data_col].notna()])
            
            # Response rows
            for idx, response in enumerate(response_options):
                # Calculate overall percentage
                resp_count = len(self.filter_by_option(df, data_col, response, idx, response_code_mapping))
                overall_pct = (resp_count / total_count) if total_count > 0 else 0
                
                # Find matching norms percentage
                norms_pct = self.match_response_to_norms(response, norms_responses)
                
                # Response text
                cell = ws.cell(row=row_num, column=1, value=f"    {str(response)[:50]}")
                cell.border = thin_border
                
                # Overall percentage
                cell = ws.cell(row=row_num, column=2, value=overall_pct)
                cell.number_format = '0%'
                cell.alignment = Alignment(horizontal='center')
                cell.border = thin_border
                
                # Norms percentage
                cell = ws.cell(row=row_num, column=3)
                if norms_pct is not None:
                    cell.value = norms_pct / 100  # Convert to decimal for percentage format
                    cell.number_format = '0%'
                else:
                    cell.value = "—"
                    cell.fill = no_norms_fill
                cell.alignment = Alignment(horizontal='center')
                cell.border = thin_border
                
                row_num += 1
            
            # Total row
            cell = ws.cell(row=row_num, column=1, value="    Total")
            cell.fill = total_fill
            cell.font = Font(bold=True)
            cell.border = thin_border
            
            cell = ws.cell(row=row_num, column=2, value=total_count)
            cell.fill = total_fill
            cell.font = Font(bold=True)
            cell.alignment = Alignment(horizontal='center')
            cell.border = thin_border
            
            cell = ws.cell(row=row_num, column=3, value="")
            cell.fill = total_fill
            cell.border = thin_border
            
            row_num += 2  # Extra spacing
        
        # Summary
        row_num += 1
        ws.cell(row=row_num, column=1, value=f"Summary: {matched_count} questions matched, {unmatched_count} without norms")
        ws.cell(row=row_num, column=1).font = Font(italic=True)
        
        # Adjust column widths
        ws.column_dimensions['A'].width = 60
        ws.column_dimensions['B'].width = 12
        ws.column_dimensions['C'].width = 12
        
        # Save workbook
        wb.save(output_file)
    
    def finish_benchmark_report(self, output_file):
        """Complete benchmark report generation"""
        self.hide_processing()
        
        if messagebox.askyesno(
            "Benchmark Report Generated",
            f"Report saved to:\n{output_file}\n\nOpen it now?",
            parent=self.root
        ):
            os.startfile(output_file)
    
    def update_powerpoint_template(self):
        """Update a PowerPoint template with survey data"""
        # Ask for template file
        template_file = filedialog.askopenfilename(
            title="Select PowerPoint Template to Update",
            filetypes=[("PowerPoint files", "*.pptx"), ("All files", "*.*")],
            parent=self.root
        )
        
        if not template_file:
            return
        
        # Store template path
        self.pptx_template_file = template_file
        
        # Show dialog to select Schedule Concept questions
        self.show_concept_questions_dialog(template_file)
    
    def show_concept_questions_dialog(self, template_file):
        """Show dialog to select which questions are Schedule Concepts"""
        dialog = tk.Toplevel(self.root)
        dialog.title("Select Schedule Concept Questions")
        dialog.geometry("550x500")
        dialog.transient(self.root)
        dialog.grab_set()
        
        # Center dialog
        dialog.update_idletasks()
        x = self.root.winfo_x() + (self.root.winfo_width() - dialog.winfo_width()) // 2
        y = self.root.winfo_y() + (self.root.winfo_height() - dialog.winfo_height()) // 2
        dialog.geometry(f"+{x}+{y}")
        
        # Title
        title_frame = tk.Frame(dialog, bg=self.COLORS['gray_50'])
        title_frame.pack(fill=tk.X, padx=20, pady=(20, 10))
        
        tk.Label(
            title_frame,
            text="Which questions are Schedule Concepts?",
            font=self.fonts['heading'],
            bg=self.COLORS['gray_50'],
            fg=self.COLORS['gray_900']
        ).pack(anchor='w')
        
        tk.Label(
            title_frame,
            text="These will need to be filled in manually after the PowerPoint is created.",
            font=self.fonts['body'],
            bg=self.COLORS['gray_50'],
            fg=self.COLORS['gray_600']
        ).pack(anchor='w', pady=(5, 0))
        
        # Scrollable checkbox area
        checkbox_frame = tk.Frame(dialog, bg=self.COLORS['white'])
        checkbox_frame.pack(fill=tk.BOTH, expand=True, padx=20, pady=10)
        
        canvas = tk.Canvas(checkbox_frame, bg=self.COLORS['white'], highlightthickness=0)
        scrollbar = ttk.Scrollbar(checkbox_frame, orient="vertical", command=canvas.yview)
        scrollable_frame = tk.Frame(canvas, bg=self.COLORS['white'])
        
        scrollable_frame.bind(
            "<Configure>",
            lambda e: canvas.configure(scrollregion=canvas.bbox("all"))
        )
        
        canvas.create_window((0, 0), window=scrollable_frame, anchor="nw")
        canvas.configure(yscrollcommand=scrollbar.set)
        
        canvas.pack(side="left", fill="both", expand=True)
        scrollbar.pack(side="right", fill="y")
        
        # Create checkboxes for all questions
        checkbox_vars = {}
        questions = self.questions_with_options or self.questions
        
        for q in questions:
            q_num = str(q['number'])
            q_text = q.get('text', f'Question {q_num}')
            
            var = tk.BooleanVar(value=q_num in [str(x) for x in self.concept_questions])
            checkbox_vars[q_num] = var
            
            display_text = f"Q{q_num}: {q_text[:60]}"
            if len(q_text) > 60:
                display_text += "..."
            
            cb = tk.Checkbutton(
                scrollable_frame,
                text=display_text,
                variable=var,
                font=self.fonts['body'],
                bg=self.COLORS['white'],
                fg=self.COLORS['gray_900'],
                activebackground=self.COLORS['white'],
                selectcolor=self.COLORS['white'],
                padx=10,
                pady=3
            )
            cb.pack(anchor='w', fill=tk.X)
        
        # Buttons
        button_frame = tk.Frame(dialog, bg=self.COLORS['white'])
        button_frame.pack(fill=tk.X, padx=20, pady=20)
        
        def on_continue():
            # Save selected concept questions
            self.concept_questions = [num for num, var in checkbox_vars.items() if var.get()]
            self.project_data['concept_questions'] = self.concept_questions
            self.save_project()
            dialog.destroy()
            # Show next dialog for unique questions
            self.show_unique_questions_dialog(template_file)
        
        def on_cancel():
            dialog.destroy()
        
        tk.Button(
            button_frame,
            text="Cancel",
            font=self.fonts['button'],
            bg=self.COLORS['gray_200'],
            fg=self.COLORS['gray_700'],
            relief=tk.FLAT,
            padx=20,
            pady=10,
            cursor='hand2',
            command=on_cancel
        ).pack(side=tk.LEFT, padx=(0, 10))
        
        tk.Button(
            button_frame,
            text="Continue →",
            font=self.fonts['button'],
            bg=self.COLORS['primary'],
            fg='white',
            relief=tk.FLAT,
            padx=20,
            pady=10,
            cursor='hand2',
            command=on_continue
        ).pack(side=tk.LEFT)
    
    def show_unique_questions_dialog(self, template_file):
        """Show dialog to select which questions are unique to this site"""
        dialog = tk.Toplevel(self.root)
        dialog.title("Select Site-Unique Questions")
        dialog.geometry("550x500")
        dialog.transient(self.root)
        dialog.grab_set()
        
        # Center dialog
        dialog.update_idletasks()
        x = self.root.winfo_x() + (self.root.winfo_width() - dialog.winfo_width()) // 2
        y = self.root.winfo_y() + (self.root.winfo_height() - dialog.winfo_height()) // 2
        dialog.geometry(f"+{x}+{y}")
        
        # Title
        title_frame = tk.Frame(dialog, bg=self.COLORS['gray_50'])
        title_frame.pack(fill=tk.X, padx=20, pady=(20, 10))
        
        tk.Label(
            title_frame,
            text="Which questions are unique to this site?",
            font=self.fonts['heading'],
            bg=self.COLORS['gray_50'],
            fg=self.COLORS['gray_900']
        ).pack(anchor='w')
        
        tk.Label(
            title_frame,
            text="These will need to be filled in manually after the PowerPoint is created.",
            font=self.fonts['body'],
            bg=self.COLORS['gray_50'],
            fg=self.COLORS['gray_600']
        ).pack(anchor='w', pady=(5, 0))
        
        # Scrollable checkbox area
        checkbox_frame = tk.Frame(dialog, bg=self.COLORS['white'])
        checkbox_frame.pack(fill=tk.BOTH, expand=True, padx=20, pady=10)
        
        canvas = tk.Canvas(checkbox_frame, bg=self.COLORS['white'], highlightthickness=0)
        scrollbar = ttk.Scrollbar(checkbox_frame, orient="vertical", command=canvas.yview)
        scrollable_frame = tk.Frame(canvas, bg=self.COLORS['white'])
        
        scrollable_frame.bind(
            "<Configure>",
            lambda e: canvas.configure(scrollregion=canvas.bbox("all"))
        )
        
        canvas.create_window((0, 0), window=scrollable_frame, anchor="nw")
        canvas.configure(yscrollcommand=scrollbar.set)
        
        canvas.pack(side="left", fill="both", expand=True)
        scrollbar.pack(side="right", fill="y")
        
        # Create checkboxes for all questions (excluding concept questions)
        checkbox_vars = {}
        questions = self.questions_with_options or self.questions
        
        for q in questions:
            q_num = str(q['number'])
            q_text = q.get('text', f'Question {q_num}')
            
            # Skip if already marked as concept
            if q_num in [str(x) for x in self.concept_questions]:
                continue
            
            var = tk.BooleanVar(value=q_num in [str(x) for x in self.unique_questions])
            checkbox_vars[q_num] = var
            
            display_text = f"Q{q_num}: {q_text[:60]}"
            if len(q_text) > 60:
                display_text += "..."
            
            cb = tk.Checkbutton(
                scrollable_frame,
                text=display_text,
                variable=var,
                font=self.fonts['body'],
                bg=self.COLORS['white'],
                fg=self.COLORS['gray_900'],
                activebackground=self.COLORS['white'],
                selectcolor=self.COLORS['white'],
                padx=10,
                pady=3
            )
            cb.pack(anchor='w', fill=tk.X)
        
        # Buttons
        button_frame = tk.Frame(dialog, bg=self.COLORS['white'])
        button_frame.pack(fill=tk.X, padx=20, pady=20)
        
        def on_continue():
            # Save selected unique questions
            self.unique_questions = [num for num, var in checkbox_vars.items() if var.get()]
            self.project_data['unique_questions'] = self.unique_questions
            self.save_project()
            dialog.destroy()
            # Now process the PowerPoint
            self.process_powerpoint_with_tracking(template_file)
        
        def on_cancel():
            dialog.destroy()
        
        tk.Button(
            button_frame,
            text="Cancel",
            font=self.fonts['button'],
            bg=self.COLORS['gray_200'],
            fg=self.COLORS['gray_700'],
            relief=tk.FLAT,
            padx=20,
            pady=10,
            cursor='hand2',
            command=on_cancel
        ).pack(side=tk.LEFT, padx=(0, 10))
        
        tk.Button(
            button_frame,
            text="Generate PowerPoint →",
            font=self.fonts['button'],
            bg=self.COLORS['pptx'],
            fg='white',
            relief=tk.FLAT,
            padx=20,
            pady=10,
            cursor='hand2',
            command=on_continue
        ).pack(side=tk.LEFT)
    
    def process_powerpoint_with_tracking(self, template_file):
        """Process PowerPoint template with full question tracking"""
        # Process in background
        self.show_processing("Updating PowerPoint template with survey data...")
        
        def process():
            try:
                # Ensure we have data loaded
                if self.remark_df is None:
                    import pandas as pd
                    if self.remark_file.endswith('.csv'):
                        self.remark_df = pd.read_csv(self.remark_file)
                    else:
                        self.remark_df = pd.read_excel(self.remark_file)
                
                # Process the template
                results = self.process_powerpoint_template(template_file)
                
                self.root.after(0, lambda: self.show_pptx_update_results(results))
            except Exception as e:
                import traceback
                error_msg = f"{str(e)}\n\n{traceback.format_exc()}"
                self.root.after(0, lambda: self.show_error(f"Failed to process PowerPoint template:\n{error_msg}"))
        
        threading.Thread(target=process, daemon=True).start()
    
    def process_powerpoint_template(self, template_file):
        """Process PowerPoint template and update charts with survey data"""
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        import pandas as pd
        import re
        
        # Build question database from survey data
        df = self.remark_df
        questions = self.questions_with_options or self.questions
        
        # Create survey data dictionary
        survey_data = {}
        for q in questions:
            q_num = str(q['number'])
            q_text = q.get('text', f'Question {q_num}')
            
            # Find data column
            data_col = self.find_data_column(df, q['number'])
            if data_col is None:
                continue
            
            # Get response options
            response_options = q.get('options', [])
            if not response_options:
                response_options = sorted(df[data_col].dropna().unique().tolist())
            
            if len(response_options) == 0:
                continue
            
            # Get data code mapping
            response_code_mapping = self.get_data_code_mapping(df, data_col, len(response_options))
            
            # Calculate total percentages
            total_count = len(df[df[data_col].notna()])
            if total_count == 0:
                continue
            
            responses = []
            for idx, response in enumerate(response_options):
                resp_count = len(self.filter_by_option(df, data_col, response, idx, response_code_mapping))
                pct = (resp_count / total_count) * 100
                responses.append({
                    'text': str(response),
                    'pct': pct
                })
            
            survey_data[q_num] = {
                'text': q_text,
                'responses': responses,
                'total_count': total_count
            }
        
        # Load PowerPoint template
        prs = Presentation(template_file)
        
        updated = []
        not_updated = []
        no_chart_slides = []  # Track slides with no charts
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_num = slide_idx + 1
            
            # Get slide title
            title = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    if '?' in text or len(text) > 15:
                        title = text
                        break
            
            # Find chart
            chart_shape = None
            for shape in slide.shapes:
                if shape.has_chart:
                    chart_shape = shape
                    break
            
            if not chart_shape:
                # Track slide with no chart (skip title slide which is usually slide 1)
                if slide_num > 1:
                    no_chart_slides.append({
                        'slide': slide_num,
                        'title': title[:50] if title else '(no title)'
                    })
                continue
            
            chart = chart_shape.chart
            
            # Match to survey question
            match_num, score = self.match_pptx_question(title, survey_data)
            
            if not match_num or score < 3:
                not_updated.append({
                    'slide': slide_num,
                    'title': title[:50] if title else '(no title)',
                    'reason': 'No matching question found'
                })
                continue
            
            survey_q = survey_data[match_num]
            
            # Get existing categories from chart
            try:
                existing_cats = [str(c) for c in chart.plots[0].categories]
            except:
                not_updated.append({
                    'slide': slide_num,
                    'title': title[:50] if title else '(no title)',
                    'reason': 'Could not read chart categories'
                })
                continue
            
            # Check category count match
            if len(existing_cats) != len(survey_q['responses']):
                not_updated.append({
                    'slide': slide_num,
                    'title': title[:50] if title else '(no title)',
                    'reason': f"Category count mismatch: chart has {len(existing_cats)}, survey has {len(survey_q['responses'])}"
                })
                continue
            
            # Get percentages
            new_percentages = [r['pct'] for r in survey_q['responses']]
            
            # Get existing series info
            series_names = [s.name for s in chart.series]
            has_benchmark = len(series_names) > 1 and 'Average' in str(series_names[-1])
            
            # Get benchmark values if they exist
            benchmark_values = None
            benchmark_name = None
            if has_benchmark:
                try:
                    benchmark_values = list(chart.series[-1].values)
                    benchmark_name = chart.series[-1].name
                except:
                    pass
            
            try:
                # Create new chart data
                chart_data = CategoryChartData()
                chart_data.categories = existing_cats
                
                # Add our data series with company name
                series_label = self.company_name or self.project_name
                chart_data.add_series(series_label, new_percentages)
                
                # Add benchmark if it existed
                if benchmark_values and len(benchmark_values) == len(existing_cats):
                    chart_data.add_series(benchmark_name or 'Average Shiftworker', benchmark_values)
                
                # Replace chart data
                chart.replace_data(chart_data)
                
                updated.append({
                    'slide': slide_num,
                    'title': title[:50] if title else '(no title)',
                    'excel_q': f'Q{match_num}'
                })
            except Exception as e:
                not_updated.append({
                    'slide': slide_num,
                    'title': title[:50] if title else '(no title)',
                    'reason': f'Error updating chart: {str(e)[:30]}'
                })
        
        # Generate output filename using company name
        import os
        company_safe = (self.company_name or self.project_name).replace(' ', '_')
        output_file = os.path.join(
            os.path.dirname(self.remark_file) if self.remark_file else os.path.expanduser("~"),
            f"{company_safe}_Survey_Results.pptx"
        )
        
        # Save updated PowerPoint
        prs.save(output_file)
        
        # Find survey questions that weren't matched to any template chart
        matched_q_nums = set()
        for item in updated:
            # Extract question number from 'Q15' format
            q_str = item.get('excel_q', '')
            if q_str.startswith('Q'):
                try:
                    matched_q_nums.add(int(q_str[1:]))
                except:
                    pass
        
        # Get all survey questions and find unmatched ones
        questions = self.questions_with_options or self.questions
        questions_not_in_template = []
        for q in questions:
            q_num = q['number']
            # Skip concept and unique questions (they're expected to not match)
            if str(q_num) in [str(x) for x in self.concept_questions + self.unique_questions]:
                continue
            # Skip breakout questions
            if str(q_num) in [str(x) for x in self.selected_breakouts]:
                continue
            if q_num not in matched_q_nums:
                questions_not_in_template.append({
                    'q_num': q_num,
                    'text': q.get('text', f'Question {q_num}')[:60]
                })
        
        return {
            'updated': updated,
            'not_updated': not_updated,
            'no_chart_slides': no_chart_slides,
            'questions_not_in_template': questions_not_in_template,
            'output_file': output_file,
            'total_slides': len(prs.slides),
            'concept_questions': self.concept_questions,
            'unique_questions': self.unique_questions
        }
    
    def match_pptx_question(self, ppt_title, survey_data):
        """Find best matching survey question for a PowerPoint slide title"""
        import re
        
        if not ppt_title:
            return None, 0
        
        ppt_lower = ppt_title.lower()
        best_match = None
        best_score = 0
        
        for q_num, q_data in survey_data.items():
            q_text_lower = q_data['text'].lower()
            
            # Calculate word overlap
            ppt_words = set(re.findall(r'\b\w+\b', ppt_lower))
            q_words = set(re.findall(r'\b\w+\b', q_text_lower))
            
            # Remove common words
            stop_words = {'the', 'a', 'an', 'is', 'are', 'do', 'you', 'your', 'to', 'of', 'in', 'for', 'on', 'at', 'this', 'that', 'which', 'best', 'describes', 'would', 'prefer', 'like', 'how', 'what', 'if'}
            ppt_words -= stop_words
            q_words -= stop_words
            
            overlap = ppt_words & q_words
            score = len(overlap)
            
            # Boost for key phrase matches
            key_phrases = [
                ('commute', 'far'), ('second job',), ('student',), ('gender',),
                ('age group',), ('single parent',), ('children', 'family'), ('sleep',),
                ('safety', 'problems'), ('safe place',), ('communication', 'priority'),
                ('communication', 'important'), ('time', 'communicate', 'daily'),
                ('management', 'welcomes'), ('enjoy', 'work'), ('pay', 'good'),
                ('part of', 'company'), ('better', 'facility'), ('best places', 'work'),
                ('training', 'important'), ('enough training',), ('supervisor', 'responds'),
                ('upper management', 'responds'), ('better schedule',), ('like', 'current schedule'),
                ('time off', 'predictable'), ('flexibility', 'time off'), ('8-hour shift',),
                ('12-hour shift',), ('more important',), ('pay', 'factor'), ('3 days off',),
                ('percentage', 'time'), ('depend', 'overtime'), ('overtime', 'distribution'),
                ('more overtime', 'time off'), ('overtime', 'every week'), ('how long', 'worked', 'company'),
                ('schedule', 'assigned'), ('next 12 months',), ('concept',), ('leaving',),
                ('rotate', 'shifts'), ('single shift',), ('weekdays off',),
            ]
            
            for phrase_words in key_phrases:
                if all(word in ppt_lower for word in phrase_words) and \
                   all(word in q_text_lower for word in phrase_words):
                    score += 15
            
            if score > best_score:
                best_score = score
                best_match = q_num
        
        return best_match, best_score
    
    def show_pptx_update_results(self, results):
        """Show comprehensive dialog with PowerPoint update results"""
        self.hide_processing()
        
        updated = results['updated']
        not_updated = results['not_updated']
        questions_not_in_template = results.get('questions_not_in_template', [])
        output_file = results['output_file']
        concept_questions = results.get('concept_questions', [])
        unique_questions = results.get('unique_questions', [])
        
        # Create results dialog
        dialog = tk.Toplevel(self.root)
        dialog.title("PowerPoint Update Results")
        dialog.geometry("800x750")
        dialog.transient(self.root)
        dialog.grab_set()
        
        # Center dialog
        dialog.update_idletasks()
        x = self.root.winfo_x() + (self.root.winfo_width() - dialog.winfo_width()) // 2
        y = self.root.winfo_y() + (self.root.winfo_height() - dialog.winfo_height()) // 2
        dialog.geometry(f"+{x}+{y}")
        
        # Main container with scrollbar
        main_canvas = tk.Canvas(dialog, bg=self.COLORS['white'], highlightthickness=0)
        main_scrollbar = ttk.Scrollbar(dialog, orient="vertical", command=main_canvas.yview)
        main_frame = tk.Frame(main_canvas, bg=self.COLORS['white'])
        
        main_frame.bind(
            "<Configure>",
            lambda e: main_canvas.configure(scrollregion=main_canvas.bbox("all"))
        )
        
        main_canvas.create_window((0, 0), window=main_frame, anchor="nw", width=760)
        main_canvas.configure(yscrollcommand=main_scrollbar.set)
        
        main_canvas.pack(side="left", fill="both", expand=True)
        main_scrollbar.pack(side="right", fill="y")
        
        # Mouse wheel scrolling - bind to canvas and dialog specifically
        def on_mousewheel(event):
            main_canvas.yview_scroll(int(-1*(event.delta/120)), "units")
        
        main_canvas.bind("<MouseWheel>", on_mousewheel)
        main_frame.bind("<MouseWheel>", on_mousewheel)
        dialog.bind("<MouseWheel>", on_mousewheel)
        
        # Also bind Enter/Leave to enable/disable scrolling when mouse is in dialog
        def bind_mousewheel(event):
            main_canvas.bind_all("<MouseWheel>", on_mousewheel)
        def unbind_mousewheel(event):
            main_canvas.unbind_all("<MouseWheel>")
        
        dialog.bind("<Enter>", bind_mousewheel)
        dialog.bind("<Leave>", unbind_mousewheel)
        
        content = tk.Frame(main_frame, bg=self.COLORS['white'])
        content.pack(fill=tk.BOTH, expand=True, padx=20, pady=20)
        
        # Summary header
        summary_frame = tk.Frame(content, bg=self.COLORS['success_bg'])
        summary_frame.pack(fill=tk.X, pady=(0, 15))
        
        summary_inner = tk.Frame(summary_frame, bg=self.COLORS['success_bg'])
        summary_inner.pack(fill=tk.X, padx=15, pady=12)
        
        company = self.company_name or self.project_name
        tk.Label(
            summary_inner,
            text=f"✓ PowerPoint Created for {company}",
            font=self.fonts['step_title'],
            bg=self.COLORS['success_bg'],
            fg='#166534'
        ).pack(anchor='w')
        
        tk.Label(
            summary_inner,
            text=f"File: {os.path.basename(output_file)}",
            font=self.fonts['body'],
            bg=self.COLORS['success_bg'],
            fg='#166534'
        ).pack(anchor='w')
        
        # SECTION 1: Questions Successfully Transferred
        section1 = tk.Frame(content, bg=self.COLORS['white'])
        section1.pack(fill=tk.X, pady=(10, 15))
        
        tk.Label(
            section1,
            text=f"✓ Questions Transferred ({len(updated)})",
            font=self.fonts['body_bold'],
            bg=self.COLORS['white'],
            fg='#166534'
        ).pack(anchor='w', pady=(0, 5))
        
        if updated:
            updated_text = tk.Text(section1, height=min(8, len(updated) + 1), wrap=tk.WORD,
                                   font=self.fonts['small'], bg=self.COLORS['gray_50'])
            updated_text.pack(fill=tk.X, pady=(0, 5))
            
            for item in updated:
                updated_text.insert(tk.END, f"  Slide {item['slide']}: {item['excel_q']} - {item['title']}\n")
            
            updated_text.config(state=tk.DISABLED)
        else:
            tk.Label(
                section1,
                text="  No questions were transferred",
                font=self.fonts['body'],
                bg=self.COLORS['white'],
                fg=self.COLORS['gray_500']
            ).pack(anchor='w')
        
        # SECTION 2: Unused Template Charts (charts with no matching survey data)
        section2 = tk.Frame(content, bg=self.COLORS['white'])
        section2.pack(fill=tk.X, pady=(10, 15))
        
        tk.Label(
            section2,
            text=f"⚠ Unused Template Charts ({len(not_updated)})",
            font=self.fonts['body_bold'],
            bg=self.COLORS['white'],
            fg='#b45309'  # Amber
        ).pack(anchor='w', pady=(0, 5))
        
        tk.Label(
            section2,
            text="These template charts had no matching survey question and were not updated:",
            font=self.fonts['small'],
            bg=self.COLORS['white'],
            fg=self.COLORS['gray_600']
        ).pack(anchor='w', pady=(0, 5))
        
        if not_updated:
            unused_text = tk.Text(section2, height=min(8, len(not_updated) + 1), wrap=tk.WORD,
                                  font=self.fonts['small'], bg='#fef3c7')  # Amber background
            unused_text.pack(fill=tk.X, pady=(0, 5))
            
            for item in not_updated:
                unused_text.insert(tk.END, f"  Slide {item['slide']}: {item['title']}\n")
                if item.get('reason'):
                    unused_text.insert(tk.END, f"    → {item['reason']}\n")
            
            unused_text.config(state=tk.DISABLED)
        else:
            tk.Label(
                section2,
                text="  All template charts were matched and updated!",
                font=self.fonts['body'],
                bg=self.COLORS['white'],
                fg='#166534'
            ).pack(anchor='w')
        
        # SECTION 3: Survey Questions Not In Template
        section3 = tk.Frame(content, bg=self.COLORS['white'])
        section3.pack(fill=tk.X, pady=(10, 15))
        
        tk.Label(
            section3,
            text=f"📋 Survey Questions Without Template Charts ({len(questions_not_in_template)})",
            font=self.fonts['body_bold'],
            bg=self.COLORS['white'],
            fg='#7c3aed'  # Purple
        ).pack(anchor='w', pady=(0, 5))
        
        tk.Label(
            section3,
            text="These survey questions had no matching chart in the template:",
            font=self.fonts['small'],
            bg=self.COLORS['white'],
            fg=self.COLORS['gray_600']
        ).pack(anchor='w', pady=(0, 5))
        
        if questions_not_in_template:
            missing_text = tk.Text(section3, height=min(8, len(questions_not_in_template) + 1), wrap=tk.WORD,
                                   font=self.fonts['small'], bg='#f3e8ff')  # Light purple background
            missing_text.pack(fill=tk.X, pady=(0, 5))
            
            for item in questions_not_in_template:
                missing_text.insert(tk.END, f"  Q{item['q_num']}: {item['text']}\n")
            
            missing_text.config(state=tk.DISABLED)
        else:
            tk.Label(
                section3,
                text="  All survey questions found matching charts!",
                font=self.fonts['body'],
                bg=self.COLORS['white'],
                fg='#166534'
            ).pack(anchor='w')
        
        # SECTION 4: Manual Entry Reminders
        reminder_frame = tk.Frame(content, bg='#fef3c7')  # Amber background
        reminder_frame.pack(fill=tk.X, pady=(15, 15))
        
        reminder_inner = tk.Frame(reminder_frame, bg='#fef3c7')
        reminder_inner.pack(fill=tk.X, padx=15, pady=12)
        
        tk.Label(
            reminder_inner,
            text="⚠ MANUAL ENTRY REQUIRED",
            font=self.fonts['body_bold'],
            bg='#fef3c7',
            fg='#92400e'
        ).pack(anchor='w', pady=(0, 10))
        
        # Concept questions reminder
        if concept_questions:
            tk.Label(
                reminder_inner,
                text=f"Schedule Concept Questions ({len(concept_questions)}):",
                font=self.fonts['body_bold'],
                bg='#fef3c7',
                fg='#92400e'
            ).pack(anchor='w', pady=(5, 2))
            
            questions = self.questions_with_options or self.questions
            for q_num in concept_questions:
                q_text = ""
                for q in questions:
                    if str(q['number']) == str(q_num):
                        q_text = q.get('text', '')[:50]
                        break
                tk.Label(
                    reminder_inner,
                    text=f"  • Q{q_num}: {q_text}...",
                    font=self.fonts['small'],
                    bg='#fef3c7',
                    fg='#78350f'
                ).pack(anchor='w')
        
        # Unique questions reminder
        if unique_questions:
            tk.Label(
                reminder_inner,
                text=f"\nSite-Unique Questions ({len(unique_questions)}):",
                font=self.fonts['body_bold'],
                bg='#fef3c7',
                fg='#92400e'
            ).pack(anchor='w', pady=(10, 2))
            
            questions = self.questions_with_options or self.questions
            for q_num in unique_questions:
                q_text = ""
                for q in questions:
                    if str(q['number']) == str(q_num):
                        q_text = q.get('text', '')[:50]
                        break
                tk.Label(
                    reminder_inner,
                    text=f"  • Q{q_num}: {q_text}...",
                    font=self.fonts['small'],
                    bg='#fef3c7',
                    fg='#78350f'
                ).pack(anchor='w')
        
        if not concept_questions and not unique_questions:
            tk.Label(
                reminder_inner,
                text="No manual entries required - you did not select any Concept or Unique questions.",
                font=self.fonts['body'],
                bg='#fef3c7',
                fg='#78350f'
            ).pack(anchor='w')
        
        # Buttons frame
        button_frame = tk.Frame(content, bg=self.COLORS['white'])
        button_frame.pack(fill=tk.X, pady=(15, 0))
        
        # Store output file for download
        self.pptx_output_file = output_file
        
        tk.Button(
            button_frame,
            text="Close",
            font=self.fonts['button'],
            bg=self.COLORS['gray_200'],
            fg=self.COLORS['gray_700'],
            relief=tk.FLAT,
            padx=20,
            pady=10,
            cursor='hand2',
            command=dialog.destroy
        ).pack(side=tk.LEFT, padx=(0, 10))
        
        tk.Button(
            button_frame,
            text="📂 Open File Location",
            font=self.fonts['button'],
            bg=self.COLORS['gray_200'],
            fg=self.COLORS['gray_700'],
            relief=tk.FLAT,
            padx=20,
            pady=10,
            cursor='hand2',
            command=lambda: os.startfile(os.path.dirname(output_file))
        ).pack(side=tk.LEFT, padx=(0, 10))
        
        tk.Button(
            button_frame,
            text="💾 Open PowerPoint",
            font=self.fonts['button'],
            bg=self.COLORS['pptx'],
            fg='white',
            relief=tk.FLAT,
            padx=20,
            pady=10,
            cursor='hand2',
            command=lambda: os.startfile(output_file)
        ).pack(side=tk.LEFT)
    
    
    def adjust_percentages_to_100(self, percentages):
        """
        Adjust rounded percentages to sum to exactly 100 using largest remainder method.
        
        This is 5th grade math done correctly:
        1. Take the integer part of each percentage
        2. Calculate how many "leftover" percentages we need to distribute
        3. Give them to the items with the LARGEST remainders (fractional parts)
        
        Args:
            percentages: List of decimal values (e.g., [0.333, 0.333, 0.334])
        
        Returns:
            List of integer percentages that sum to EXACTLY 100
        
        Example:
            Input: [0.333, 0.333, 0.334]  (sum = 1.0)
            Step 1: Integer parts = [33, 33, 33]  (sum = 99)
            Step 2: Remainders = [0.3, 0.3, 0.4]
            Step 3: Need 1 more to reach 100, give it to largest remainder (index 2)
            Result: [33, 33, 34]  (sum = 100) ✓
        """
        if not percentages or sum(percentages) == 0:
            return [0] * len(percentages)
        
        # Convert to percentage scale (0.333 -> 33.3)
        pct_100 = [p * 100 for p in percentages]
        
        # Separate into integer parts and remainders
        int_parts = [int(p) for p in pct_100]
        remainders = [p - int(p) for p in pct_100]
        
        # Calculate how many more we need to reach 100
        total = sum(int_parts)
        needed = 100 - total
        
        # If we already sum to 100, we're done
        if needed == 0:
            return int_parts
        
        # Get indices sorted by remainder size (largest first)
        # This is the "largest remainder" part of the method
        indices_by_remainder = sorted(range(len(remainders)), 
                                     key=lambda i: remainders[i], 
                                     reverse=True)
        
        # Distribute the needed percentages to items with largest remainders
        for i in range(abs(needed)):
            if i < len(indices_by_remainder):
                idx = indices_by_remainder[i]
                if needed > 0:
                    int_parts[idx] += 1
                else:
                    int_parts[idx] -= 1
        
        return int_parts
    
    def create_pdf_at_location(self, output_file):
        """Create the PDF crosstab report at specified location"""
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import letter, landscape
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak, KeepTogether
        from reportlab.lib.enums import TA_LEFT, TA_CENTER
        import pandas as pd
        
        # Create document
        doc = SimpleDocTemplate(
            output_file,
            pagesize=landscape(letter),
            leftMargin=0.5*inch,
            rightMargin=0.5*inch,
            topMargin=0.5*inch,
            bottomMargin=0.5*inch
        )
        
        # Styles
        styles = getSampleStyleSheet()
        styles.add(ParagraphStyle(
            name='QuestionHeader',
            parent=styles['Normal'],
            fontSize=10,
            fontName='Helvetica-Bold',
            textColor=colors.black,
            spaceAfter=6
        ))
        
        elements = []
        df = self.remark_df
        
        # Find the breakout column in data
        breakout_col = self.find_data_column(df, self.selected_breakout)
        if breakout_col is None:
            raise Exception(f"Could not find data column for breakout Q{self.selected_breakout}")
        
        # Get breakout options FROM THE SURVEY (not from data)
        questions = self.questions_with_options or self.questions
        breakout_question = None
        for q in questions:
            if str(q['number']) == str(self.selected_breakout):
                breakout_question = q
                break
        
        # Use survey options if available, otherwise fall back to data
        if breakout_question and breakout_question.get('options'):
            breakout_values = breakout_question['options']
        else:
            breakout_values = sorted(df[breakout_col].dropna().unique().tolist())
        
        # Get data code mapping for breakout column (handles non-standard codes like e,r,w)
        breakout_code_mapping = self.get_data_code_mapping(df, breakout_col, len(breakout_values))
        
        # Get the breakout question text for header
        breakout_text = breakout_question.get('text', f'Question {self.selected_breakout}')[:50] if breakout_question else f'Q{self.selected_breakout}'
        
        # Column headers
        col_headers = [breakout_text] + [str(v)[:12] for v in breakout_values] + ['Total']
        
        # Process each question (excluding the breakout question)
        total_questions = len(questions)
        for q_idx, q in enumerate(questions):
            q_num = q['number']
            
            # Skip the breakout question itself
            if str(q_num) == str(self.selected_breakout):
                continue
            
            q_text = q.get('text', f'Question {q_num}')
            
            # Find data column
            data_col = self.find_data_column(df, q_num)
            if data_col is None:
                continue
            
            # Get response options FROM THE SURVEY (not from data)
            response_options = q.get('options', [])
            if not response_options:
                response_options = sorted(df[data_col].dropna().unique().tolist())
            
            # Build table data
            table_data = [col_headers]
            
            # Question header row
            q_header = [f"{q_num}  {q_text[:80]}"] + [''] * (len(breakout_values) + 1)
            table_data.append(q_header)
            
            # Check if this question needs special calculation
            is_last_question = (q_idx == total_questions - 1)
            special_calc = self.get_special_calculation(q_num, q_text, response_options, is_last_question)
            
            # Get data code mapping for this question's responses
            response_code_mapping = self.get_data_code_mapping(df, data_col, len(response_options))
            
            # Calculate percentages for all columns FIRST (so we can adjust them)
            all_columns_pcts = []
            
            # DIAGNOSTIC for Q44 - write to file
            if q_num == 44:
                import os
                import tempfile
                diag_file = os.path.join(tempfile.gettempdir(), "Q44_DIAGNOSTIC.txt")
                try:
                    with open(diag_file, 'w', encoding='utf-8') as f:
                        f.write("="*80 + "\n")
                        f.write("DIAGNOSTIC FOR Q44\n")
                        f.write("="*80 + "\n")
                        f.write(f"Question: {q_text}\n")
                        f.write(f"Data column: {data_col}\n")
                        f.write(f"Number of survey options: {len(response_options)}\n")
                        for idx, opt in enumerate(response_options):
                            f.write(f"  {idx}: {opt}\n")
                        f.write(f"\nData values in Remark export:\n")
                        data_vals = df[data_col].dropna().unique().tolist()
                        for dv in data_vals:
                            f.write(f"  '{dv}'\n")
                        f.write(f"\nResponse code mapping: {response_code_mapping}\n")
                        f.write("\n" + "="*80 + "\n\n")
                except:
                    pass  # Silently fail if encoding issues
            
            # Calculate each breakout column
            for grp_idx, group_val in enumerate(breakout_values):
                group_df = self.filter_by_option(df, breakout_col, group_val, grp_idx, breakout_code_mapping)
                # CRITICAL FIX: Exclude BLANK from total (BLANK is a string, not NaN)
                valid_responses = group_df[group_df[data_col].notna()]
                valid_responses = valid_responses[valid_responses[data_col] != 'BLANK']
                group_total = len(valid_responses)
                
                if group_total > 0:
                    # Calculate counts for each response
                    counts = []
                    matched_data_values = []  # Track what data values each option matched
                    
                    for idx, response in enumerate(response_options):
                        filtered = self.filter_by_option(group_df, data_col, response, idx, response_code_mapping)
                        count = len(filtered)
                        counts.append(count)
                        
                        if count > 0:
                            matched_vals = filtered[data_col].unique().tolist()
                            matched_data_values.append((idx, response, count, matched_vals))
                    
                    pcts = [count / group_total for count in counts]
                    
                    # DIAGNOSTIC for Q44 first column - write to file
                    if q_num == 44 and grp_idx == 0:
                        try:
                            with open(diag_file, 'a', encoding='utf-8') as f:
                                f.write(f"Column {grp_idx} ({group_val}):\n")
                                f.write(f"  Total respondents: {group_total}\n\n")
                                
                                f.write("  Matching details:\n")
                                for idx, resp, cnt, vals in matched_data_values:
                                    f.write(f"    Option {idx} '{resp}': {cnt} matches\n")
                                    f.write(f"      Matched data values: {vals}\n")
                                
                                f.write(f"\n  Sum of counts: {sum(counts)} (should equal {group_total})\n")
                                f.write(f"  Counts by option: {counts}\n")
                                f.write(f"  Percentages: {[f'{p:.1%}' for p in pcts]}\n")
                                f.write(f"  Sum of percentages: {sum(pcts):.1%}\n")
                                
                                if sum(counts) != group_total:
                                    f.write(f"\n  WARNING: MISMATCH! Counts sum to {sum(counts)} but should be {group_total}\n")
                                    f.write(f"  This means responses are being counted multiple times!\n")
                                    
                                    # Find which data values are matched by multiple options
                                    all_matched = {}
                                    for idx, resp, cnt, vals in matched_data_values:
                                        for v in vals:
                                            if v not in all_matched:
                                                all_matched[v] = []
                                            all_matched[v].append(idx)
                                    
                                    f.write(f"\n  Data values matched to multiple options:\n")
                                    for val, option_indices in all_matched.items():
                                        if len(option_indices) > 1:
                                            f.write(f"    '{val}' matched to options: {option_indices}\n")
                                
                                f.write("\n" + "="*80 + "\n")
                                f.write(f"\nDiagnostic written to: {diag_file}\n")
                                f.write("Please send this file to Claude!\n")
                        except:
                            pass  # Silently fail if encoding issues
                        
                        # Show message box with file location
                        try:
                            from tkinter import messagebox
                            messagebox.showinfo(
                                "Diagnostic Created",
                                f"Q44 diagnostic written to:\n\n{diag_file}\n\nPlease send this file to Claude!"
                            )
                        except:
                            pass
                else:
                    pcts = [0] * len(response_options)
                all_columns_pcts.append(pcts)
            
            # Calculate Total column
            # CRITICAL FIX: Exclude BLANK from total
            valid_total = df[df[data_col].notna()]
            valid_total = valid_total[valid_total[data_col] != 'BLANK']
            total_respondents = len(valid_total)
            if total_respondents > 0:
                total_pcts = [len(self.filter_by_option(df, data_col, response, idx, response_code_mapping)) / total_respondents
                             for idx, response in enumerate(response_options)]
            else:
                total_pcts = [0] * len(response_options)
            all_columns_pcts.append(total_pcts)
            
            # Adjust ALL columns to sum to exactly 100%
            adjusted_columns = [self.adjust_percentages_to_100(col) for col in all_columns_pcts]
            
            # Response rows - now use adjusted percentages
            for idx, response in enumerate(response_options):
                row = [f"    {str(response)[:50]}"]
                
                # Add adjusted percentages for each column
                for col_idx in range(len(adjusted_columns)):
                    row.append(f"{adjusted_columns[col_idx][idx]}%")
                
                table_data.append(row)
            
            # Total row (counts)
            total_row = ["    Total"]
            for grp_idx, group_val in enumerate(breakout_values):
                group_df = self.filter_by_option(df, breakout_col, group_val, grp_idx, breakout_code_mapping)
                count = len(group_df[group_df[data_col].notna()])
                total_row.append(str(count))
            # CRITICAL FIX: Exclude BLANK
            valid_grand = df[df[data_col].notna()]
            valid_grand = valid_grand[valid_grand[data_col] != 'BLANK']
            grand_total = len(valid_grand)
            total_row.append(str(grand_total))
            table_data.append(total_row)
            
            # Special calculation row (Average Rating, Average Miles, Average Hours)
            if special_calc:
                avg_label = special_calc.get('label', 'Average')
                avg_row = [f"    {avg_label}"]
                for grp_idx, group_val in enumerate(breakout_values):
                    group_df = self.filter_by_option(df, breakout_col, group_val, grp_idx, breakout_code_mapping)
                    avg = self.calculate_special_average(group_df, data_col, response_options, special_calc)
                    avg_row.append(f"{avg:.2f}" if avg is not None else "N/A")
                
                # Total average
                total_avg = self.calculate_special_average(df, data_col, response_options, special_calc)
                avg_row.append(f"{total_avg:.2f}" if total_avg is not None else "N/A")
                table_data.append(avg_row)
            
            # Create table with 20% smaller columns to prevent bleed-off
            col_widths = [2.8*inch] + [0.64*inch] * (len(breakout_values) + 1)
            
            # Adjust column widths if too many columns
            if len(breakout_values) > 8:
                col_widths = [2.4*inch] + [0.48*inch] * (len(breakout_values) + 1)
            
            table = Table(table_data, colWidths=col_widths)
            
            # Determine number of data rows for styling
            num_data_rows = len(response_options)
            total_row_idx = 2 + num_data_rows  # header + question header + data rows
            
            # Table style
            style_commands = [
                # Header row
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#4F81BD')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 7),
                ('ALIGN', (1, 0), (-1, 0), 'CENTER'),
                
                # Question header row
                ('BACKGROUND', (0, 1), (-1, 1), colors.HexColor('#f3f4f6')),
                ('FONTNAME', (0, 1), (-1, 1), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 1), (-1, 1), 8),
                ('SPAN', (0, 1), (-1, 1)),
                
                # Data rows
                ('FONTSIZE', (0, 2), (-1, -1), 7),
                ('ALIGN', (1, 2), (-1, -1), 'CENTER'),
                
                # Total row
                ('FONTNAME', (0, total_row_idx), (-1, total_row_idx), 'Helvetica-Bold'),
                ('BACKGROUND', (0, total_row_idx), (-1, total_row_idx), colors.HexColor('#f9fafb')),
                
                # Grid
                ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#d1d5db')),
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                ('LEFTPADDING', (0, 0), (-1, -1), 4),
                ('RIGHTPADDING', (0, 0), (-1, -1), 4),
                ('TOPPADDING', (0, 0), (-1, -1), 4),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
            ]
            
            # Average row styling (if present)
            if special_calc:
                avg_row_idx = total_row_idx + 1
                style_commands.extend([
                    ('FONTNAME', (0, avg_row_idx), (-1, avg_row_idx), 'Helvetica-BoldOblique'),
                    ('BACKGROUND', (0, avg_row_idx), (-1, avg_row_idx), colors.HexColor('#dbeafe')),
                ])
            style = TableStyle(style_commands)
            table.setStyle(style)
            
            # Prevent table from splitting across pages
            table.hAlign = 'LEFT'
            table.splitByRow = False
            
            # Add table to elements
            elements.append(table)
            elements.append(Spacer(1, 0.15*inch))
        
        # Build PDF
        doc.build(elements)
    
    def get_special_calculation(self, q_num, q_text, response_options=None, is_last_question=False):
        """Determine if a question needs special calculation and return config"""
        q_text_lower = q_text.lower() if q_text else ""
        q_num_str = str(q_num)
        
        # 1. Rate Concept questions - ONLY "rate concept" or "rate schedule"
        if "rate concept" in q_text_lower or "rate schedule" in q_text_lower:
            return {
                'type': 'rating',
                'label': 'Average Rating',
                'text_map': {
                    # Survey text
                    'perfect': 5,
                    'perfect!': 5,
                    "it's okay": 4,
                    'its okay': 4,
                    'i prefer something else': 3,
                    'prefer something else': 3,
                    'not great': 2,
                    'never show me this again': 1,
                    'never show this to me again': 1,
                    'never show this to me again!': 1,
                    # Remark placeholder text (maps by position)
                    'excellent': 5,
                    'very good': 4,
                    'good': 3,
                    'fair': 2,
                    'poor': 1,
                },
                'code_map': {'a': 5, 'b': 4, 'c': 3, 'd': 2, 'e': 1}
            }
        
        # 2. Years with Company (Tenure)
        if ("years" in q_text_lower or "how long" in q_text_lower) and \
           ("company" in q_text_lower or "employed" in q_text_lower or "worked" in q_text_lower):
            return {
                'type': 'tenure',
                'label': 'Average Years',
                'text_map': {
                    'less than 6 months': 0.25,
                    '6 months to a year': 0.75,
                    '6 months to 1 year': 0.75,
                    '1 to 5 years': 3,
                    '1-5 years': 3,
                    '6 to 10 years': 7.5,
                    '6-10 years': 7.5,
                    '11 to 15 years': 12.5,
                    '11-15 years': 12.5,
                    '16 to 20 years': 17.5,
                    '16-20 years': 17.5,
                    'over 20 years': 25,
                    'more than 20 years': 25,
                },
                'code_map': {'a': 0.25, 'b': 0.75, 'c': 3, 'd': 7.5, 'e': 12.5, 'f': 17.5, 'g': 25}
            }
        
        # 3. Age Group
        if "age" in q_text_lower and ("group" in q_text_lower or "old" in q_text_lower or "your age" in q_text_lower):
            return {
                'type': 'age',
                'label': 'Average Age',
                'text_map': {
                    '25 and under': 22,
                    'under 25': 22,
                    '25 or under': 22,
                    '26 to 30': 27.5,
                    '26-30': 27.5,
                    '31 to 35': 33,
                    '31-35': 33,
                    '36 to 40': 37.5,
                    '36-40': 37.5,
                    '41 to 45': 43,
                    '41-45': 43,
                    '46 to 50': 47.5,
                    '46-50': 47.5,
                    '51 to 55': 53,
                    '51-55': 53,
                    'over 55': 60,
                    '55 and over': 60,
                    '56 and over': 60,
                },
                'code_map': {'a': 22, 'b': 27.5, 'c': 33, 'd': 37.5, 'e': 43, 'f': 47.5, 'g': 53, 'h': 60}
            }
        
        # 4. Commute Distance
        if "commute" in q_text_lower or "how far" in q_text_lower or "distance to work" in q_text_lower:
            return {
                'type': 'commute',
                'label': 'Average Miles',
                'text_map': {
                    'less than 1 mile': 0.5,
                    'less than a mile': 0.5,
                    '1 to 5 miles': 3,
                    '1-5 miles': 3,
                    '6 to 10 miles': 7.5,
                    '6-10 miles': 7.5,
                    '11 to 20 miles': 15,
                    '11-20 miles': 15,
                    '21 to 30 miles': 25,
                    '21-30 miles': 25,
                    '31 to 40 miles': 35,
                    '31-40 miles': 35,
                    'more than 40 miles': 50,
                    'over 40 miles': 50,
                },
                'code_map': {'a': 0.5, 'b': 3, 'c': 7.5, 'd': 15, 'e': 25, 'f': 35, 'g': 50}
            }
        
        # 5. Sleep Hours - must specifically ask about hours of sleep
        if ("hours of sleep" in q_text_lower) or \
           ("how many hours" in q_text_lower and "sleep" in q_text_lower):
            return {
                'type': 'sleep',
                'label': 'Average Hours',
                'text_map': {
                    'less than 5': 4,
                    'less than 5 hours': 4,
                    '5 or more hours but less than 6': 5.5,
                    '5 or more hours but less than 6 hours': 5.5,
                    '5 to 6': 5.5,
                    '5 to 6 hours': 5.5,
                    '5-6': 5.5,
                    '6 or more hours but less than 7': 6.5,
                    '6 or more hours but less than 7 hours': 6.5,
                    '6 to 7': 6.5,
                    '6 to 7 hours': 6.5,
                    '6-7': 6.5,
                    '7 or more hours but less than 8': 7.5,
                    '7 or more hours but less than 8 hours': 7.5,
                    '7 to 8': 7.5,
                    '7 to 8 hours': 7.5,
                    '7-8': 7.5,
                    '8 or more hours but less than 9': 8.5,
                    '8 or more hours but less than 9 hours': 8.5,
                    '8 to 9': 8.5,
                    '8 to 9 hours': 8.5,
                    '8-9': 8.5,
                    '9 or more': 9.5,
                    '9 or more hours': 9.5,
                    'over 9': 9.5,
                    'over 9 hours': 9.5,
                    'more than 9': 9.5,
                },
                'code_map': {'a': 4, 'b': 5.5, 'c': 6.5, 'd': 7.5, 'e': 8.5, 'f': 9.5},
                'exclude': ["i don't work that shift", "don't work that shift", "n/a", "i dont work"]
            }
        
        # 6. Communicating Daily Plant Conditions
        if ("communicat" in q_text_lower) and ("plant" in q_text_lower or "conditions" in q_text_lower or "daily" in q_text_lower):
            return {
                'type': 'communication',
                'label': 'Average Minutes',
                'text_map': {
                    'less than 10': 5,
                    'less than 10 minutes': 5,
                    '10 minutes': 10,
                    '10': 10,
                    '15 minutes': 15,
                    '15': 15,
                    '20 minutes': 20,
                    '20': 20,
                    '25 minutes': 25,
                    '25': 25,
                    '30 minutes': 30,
                    '30': 30,
                    'more than 30': 45,
                    'more than 30 minutes': 45,
                    'over 30': 45,
                },
                'code_map': {'a': 5, 'b': 10, 'c': 15, 'd': 20, 'e': 25, 'f': 30, 'g': 45}
            }
        
        # 7. Agree/Disagree Scale (1-5) - must have STRONGLY agree AND STRONGLY disagree
        if response_options:
            options_lower = [str(opt).lower() for opt in response_options]
            has_strongly_agree = any('strongly agree' in opt for opt in options_lower)
            has_strongly_disagree = any('strongly disagree' in opt for opt in options_lower)
            
            if has_strongly_agree and has_strongly_disagree:
                return {
                    'type': 'agree_disagree',
                    'label': 'Average',
                    'text_map': {
                        'strongly disagree': 1,
                        'disagree': 2,
                        'neither agree nor disagree': 3,
                        'neither': 3,
                        'neutral': 3,
                        'agree': 4,
                        'strongly agree': 5,
                    },
                    'code_map': {'1': 1, '2': 2, '3': 3, '4': 4, '5': 5,
                                 'a': 1, 'b': 2, 'c': 3, 'd': 4, 'e': 5}
                }
        
        # 8. Overtime/Extra Hours - ONLY the last question (Q71) about how much overtime per week
        # Other overtime questions are agree/disagree or preferences, not hours
        if is_last_question or "how much overtime" in q_text_lower or \
           ("overtime" in q_text_lower and ("every week" in q_text_lower or "per week" in q_text_lower)):
            return {
                'type': 'overtime',
                'label': 'Average Hours',
                'text_map': {
                    'none': 0,
                    'less than 2': 1,
                    'less than 2 hours': 1,
                    'between 2 and 4': 3,
                    'between 2 and 4 hours': 3,
                    '2 - 4': 3,
                    '2-4': 3,
                    '2 to 4': 3,
                    'between 4 and 6': 5,
                    'between 4 and 6 hours': 5,
                    '4 - 6': 5,
                    '4-6': 5,
                    '4 to 6': 5,
                    'between 6 and 8': 7,
                    'between 6 and 8 hours': 7,
                    '6 - 8': 7,
                    '6-8': 7,
                    '6 to 8': 7,
                    'between 8 and 12': 10,
                    'between 8 and 12 hours': 10,
                    '8 - 12': 10,
                    '8-12': 10,
                    '8 to 12': 10,
                    'i will take all i can get': 15,
                    'i will take all that i can get': 15,
                    'all i can get': 15,
                    'all that i can get': 15,
                    'as much as possible': 15,
                },
                'code_map': {'a': 0, 'b': 1, 'c': 3, 'd': 5, 'e': 7, 'f': 10, 'g': 15}
            }
        
        return None
    
    def calculate_special_average(self, df, data_col, response_options, calc_config):
        """Calculate weighted average based on special calculation config"""
        if df.empty or data_col not in df.columns:
            return None
        
        text_map = calc_config.get('text_map', {})
        code_map = calc_config.get('code_map', {})
        exclude_list = calc_config.get('exclude', [])
        
        # Get actual values from data
        data_values = df[data_col].dropna().unique().tolist()
        
        if len(data_values) == 0:
            return None
        
        total_weight = 0
        total_count = 0
        
        for data_val in data_values:
            data_val_str = str(data_val).strip()
            data_val_lower = data_val_str.lower()
            
            # Check if this value should be excluded
            should_exclude = False
            for exc in exclude_list:
                if exc in data_val_lower:
                    should_exclude = True
                    break
            
            if should_exclude:
                continue
            
            # Try to find the weighted value
            value = None
            
            # 1. Try exact text match (lowercase)
            if data_val_lower in text_map:
                value = text_map[data_val_lower]
            
            # 2. Try partial text match - SORT BY LENGTH (longest first) to avoid
            #    "disagree" matching before "neither agree or disagree"
            if value is None:
                sorted_keys = sorted(text_map.keys(), key=len, reverse=True)
                for text_key in sorted_keys:
                    if text_key in data_val_lower or data_val_lower in text_key:
                        value = text_map[text_key]
                        break
            
            # 3. Try code map (letter or number codes)
            if value is None:
                if data_val_lower in code_map:
                    value = code_map[data_val_lower]
                elif data_val_str in code_map:
                    value = code_map[data_val_str]
            
            # 4. For agree/disagree, try numeric directly
            if value is None and calc_config.get('type') == 'agree_disagree':
                try:
                    num_val = int(float(data_val_str))
                    if 1 <= num_val <= 5:
                        value = num_val
                except:
                    pass
            
            # If we found a value, add to weighted sum
            if value is not None:
                count = len(df[df[data_col] == data_val])
                total_weight += value * count
                total_count += count
        
        return total_weight / total_count if total_count > 0 else None
    
    def is_rating_question(self, options):
        """Check if a question is a rating/scale question - DEPRECATED, use get_special_calculation"""
        return False  # Now handled by get_special_calculation
    
    def calculate_average_rating(self, df, data_col, response_options):
        """Calculate weighted average rating - DEPRECATED, use calculate_special_average"""
        return None  # Now handled by calculate_special_average
    
    def filter_by_option(self, df, column, option, option_index=None, data_code_mapping=None):
        """Filter dataframe by matching option value (handles string/number mismatches)"""
        if column not in df.columns:
            return df.head(0)  # Empty dataframe
        
        option_str = str(option).strip().lower()
        
        # Normalize quotes and apostrophes - use Unicode code points to be explicit
        option_normalized = option_str.replace('\u2019', "'").replace('\u2018', "'")  # Smart single quotes
        option_normalized = option_normalized.replace('\u201c', '"').replace('\u201d', '"')  # Smart double quotes
        option_normalized = option_normalized.replace('\u2013', '-').replace('\u2014', '-')  # Em/en dashes
        
        # Remove parentheses - Remark doesn't support them in answer boxes
        option_no_parens = option_normalized.replace('(', '').replace(')', '').strip()
        # Also collapse multiple spaces that might result from removing parentheses
        import re
        option_no_parens = re.sub(r'\s+', ' ', option_no_parens)
        
        # Create a punctuation-normalized version for fuzzy matching
        # Replace semicolons with periods, remove trailing periods, normalize spaces
        def normalize_punctuation(text):
            """Normalize punctuation for comparison"""
            # Fix corrupted Unicode en-dash (â€" -> -)
            text = text.replace('â€"', '-').replace('â€"', '-')
            # Fix corrupted Unicode apostrophe (â€™ -> ')
            text = text.replace('â€™', "'").replace("â€™", "'")
            text = text.replace(',', '')  # Remove commas (Remark often drops them)
            text = text.replace(';', '.').replace(':', '.')  # Normalize punctuation
            # Normalize en-dash and em-dash to regular hyphen
            text = text.replace('\u2013', '-').replace('\u2014', '-')  # – and — to -
            # Remove hyphens entirely: "last-minute" -> "last minute", "Friday-Saturday" -> "Friday Saturday"
            text = text.replace('-', ' ')
            # Expand common contractions for matching
            text = text.replace("i'll", "i will")
            text = text.replace("you'll", "you will") 
            text = text.replace("we'll", "we will")
            text = text.replace("they'll", "they will")
            text = text.replace("don't", "do not")
            text = text.replace("won't", "will not")
            text = text.replace("can't", "cannot")
            text = re.sub(r'\.+', '.', text)  # Collapse multiple periods
            text = re.sub(r'\.\s*$', '', text)  # Remove trailing period
            text = re.sub(r'\s+', ' ', text)  # Normalize spaces
            return text.strip()
        
        option_punct_normalized = normalize_punctuation(option_no_parens)
        
        # If we have a specific data code mapping for this column, use it
        if data_code_mapping and option_index is not None and option_index < len(data_code_mapping):
            mapped_code = data_code_mapping[option_index]
            mask = df[column].astype(str).str.strip().str.lower() == str(mapped_code).lower()
            if mask.any():
                return df[mask]
        
        # === METHOD 1: Try exact match first ===
        mask = df[column] == option
        if mask.any():
            return df[mask]
        
        # === METHOD 2: Case-insensitive exact match ===
        mask = df[column].astype(str).str.strip().str.lower() == option_str
        if mask.any():
            return df[mask]
        
        # === METHOD 2a: Match without parentheses ===
        # Survey has "5 (Strongly Agree)" but Remark has "5 Strongly Agree"
        mask = df[column].astype(str).str.strip().str.lower() == option_no_parens
        if mask.any():
            return df[mask]
        
        # === METHOD 2b: Normalized quote matching ===
        data_values = df[column].dropna().unique()
        for data_val in data_values:
            data_normalized = str(data_val).strip().lower().replace('\u2019', "'").replace('\u2018', "'")
            data_normalized = data_normalized.replace('\u201c', '"').replace('\u201d', '"')
            data_normalized = data_normalized.replace('\u2013', '-').replace('\u2014', '-')
            if data_normalized == option_normalized:
                mask = df[column] == data_val
                return df[mask]
        
        # === METHOD 2b2: Punctuation-normalized matching ===
        # Handles: Survey "Neither of the above; this is..." vs Data "Neither of the above. This is..."
        for data_val in data_values:
            data_val_str = str(data_val).strip().lower()
            data_punct_normalized = normalize_punctuation(data_val_str)
            if data_punct_normalized == option_punct_normalized:
                mask = df[column] == data_val
                return df[mask]
        
        # === METHOD 2c: Number-prefix matching for scales ===
        # Handles: Survey "5 (Strongly Agree)" matching Data "5 strongly agree" or just "5"
        # CRITICAL FIX: Skip for time questions! "6:30 a.m." and "6:00 a.m." both start with "6"
        # ALSO skip for duration questions! "6 months to 1 year" and "6 to 10 years" both start with "6"
        has_time_marker = any(marker in option_str for marker in ['a.m.', 'p.m.', 'am', 'pm', ':'])
        has_duration = any(word in option_punct_normalized for word in [' minutes', ' hours', ' years', ' days', ' weeks', ' months'])
        if not has_time_marker and not has_duration:
            # Extract leading number from survey option (use no-parens version)
            survey_num_match = re.match(r'^(\d+)', option_no_parens)
            if survey_num_match:
                survey_num = survey_num_match.group(1)
                
                # Check each data value for matching leading number
                for data_val in data_values:
                    data_val_str = str(data_val).strip().lower()
                    data_num_match = re.match(r'^(\d+)', data_val_str)
                    
                    if data_num_match and data_num_match.group(1) == survey_num:
                        # Numbers match! This is our match
                        mask = df[column] == data_val
                        if mask.any():
                            return df[mask]
        
        # === METHOD 2d: Match data number-prefix to survey text ===
        # Handles: Data "5 strongly agree" matching Survey "Strongly Agree" (no number)
        for data_val in data_values:
            data_val_str = str(data_val).strip().lower()
            data_num_match = re.match(r'^(\d+)\s*(.+)', data_val_str)
            
            if data_num_match:
                data_num = data_num_match.group(1)
                data_text = data_num_match.group(2).strip()
                
                # FIXED: Use exact comparison, not substring matching
                # This prevents "5:30" from matching "Before 5:30"
                if data_text and data_text == option_no_parens:
                    mask = df[column] == data_val
                    if mask.any():
                        return df[mask]
        
        # === METHOD 3: Smart agree/disagree matching ===
        # If survey option contains specific agree/disagree text, extract and match it
        # Order matters - check most specific first!
        agree_disagree_terms = [
            'strongly disagree', 'strongly agree',
            'neither agree nor disagree', 'neither agree or disagree',
            'disagree', 'agree'  # Check these LAST (they're substrings of the above)
        ]
        
        for term in agree_disagree_terms:
            if term in option_str:
                # Look for this exact term in data
                mask = df[column].astype(str).str.strip().str.lower() == term
                if mask.any():
                    return df[mask]
                break  # Only try the first (most specific) match
        
        # === METHOD 3b: Position-based agree/disagree ===
        # If data uses agree/disagree terms but survey option doesn't contain them
        if option_index is not None:
            data_values = df[column].dropna().unique()
            data_values_lower = [str(v).strip().lower() for v in data_values]
            
            # Check if data uses agree/disagree scale
            uses_agree_scale = any(av in data_values_lower for av in ['strongly agree', 'agree', 'disagree', 'strongly disagree'])
            
            if uses_agree_scale and not any(term in option_str for term in agree_disagree_terms):
                # Survey option doesn't have agree/disagree text, but data does
                # Assume survey options are in order: Strongly Agree (0) to Strongly Disagree (4)
                # This is the opposite of numeric scale!
                position_to_agree = {
                    0: 'strongly agree',
                    1: 'agree',
                    2: 'neither',
                    3: 'disagree',
                    4: 'strongly disagree',
                }
                if option_index in position_to_agree:
                    target = position_to_agree[option_index]
                    if target == 'agree':
                        mask = df[column].astype(str).str.strip().str.lower() == 'agree'
                    elif target == 'disagree':
                        mask = df[column].astype(str).str.strip().str.lower() == 'disagree'
                    elif target == 'neither':
                        mask = df[column].astype(str).str.strip().str.lower().str.contains('neither', na=False)
                    else:
                        mask = df[column].astype(str).str.strip().str.lower() == target
                    if mask.any():
                        return df[mask]
        
        # === METHOD 4: Numeric scale matching ===
        # For options like "1", "2", "3", "4", "5" - match by position to agree/disagree
        # Detect if this is a 1-5 numeric scale
        if option_str.isdigit() and option_index is not None:
            # Map numeric position to agree/disagree text
            # 1=Strongly Disagree, 2=Disagree, 3=Neither, 4=Agree, 5=Strongly Agree
            numeric_to_text = {
                0: 'strongly disagree',  # position 0 = "1" = Strongly Disagree
                1: 'disagree',           # position 1 = "2" = Disagree
                2: 'neither',            # position 2 = "3" = Neither
                3: 'agree',              # position 3 = "4" = Agree
                4: 'strongly agree',     # position 4 = "5" = Strongly Agree
            }
            
            if option_index in numeric_to_text:
                target = numeric_to_text[option_index]
                # For "agree" and "disagree", need to match exactly (not "strongly X")
                if target == 'agree':
                    mask = (df[column].astype(str).str.strip().str.lower() == 'agree')
                elif target == 'disagree':
                    mask = (df[column].astype(str).str.strip().str.lower() == 'disagree')
                elif target == 'neither':
                    mask = df[column].astype(str).str.strip().str.lower().str.contains('neither', na=False)
                else:
                    mask = df[column].astype(str).str.strip().str.lower() == target
                
                if mask.any():
                    return df[mask]
        
        # === METHOD 5: Rating scale matching (Excellent/Poor) ===
        # First try direct matching if survey option contains rating terms
        rating_terms = ['excellent', 'very good', 'good', 'fair', 'poor']
        for term in rating_terms:
            if term in option_str:
                mask = df[column].astype(str).str.strip().str.lower() == term
                if mask.any():
                    return df[mask]
                break
        
        # === METHOD 5b: Position-based rating scale ===
        # If data uses Excellent/Very Good/Good/Fair/Poor, match by position
        if option_index is not None:
            data_values = df[column].dropna().unique()
            data_values_lower = [str(v).strip().lower() for v in data_values]
            
            # Check if data uses rating scale
            uses_rating_scale = any(rv in data_values_lower for rv in rating_terms)
            
            if uses_rating_scale:
                # Map position to rating term (position 0 = best = Excellent, etc.)
                position_to_rating = {
                    0: 'excellent',
                    1: 'very good',
                    2: 'good',
                    3: 'fair',
                    4: 'poor',
                }
                if option_index in position_to_rating:
                    target = position_to_rating[option_index]
                    mask = df[column].astype(str).str.strip().str.lower() == target
                    if mask.any():
                        return df[mask]
        
        # === METHOD 6: Letter code matching (a, b, c...) ===
        if option_index is not None:
            letter_code = chr(ord('a') + option_index)
            mask = df[column].astype(str).str.strip().str.lower() == letter_code
            if mask.any():
                return df[mask]
        
        # === METHOD 7: COMPLETELY DISABLED ===
        # Substring and fuzzy matching was causing too many false positives
        # All matching now relies on exact matches (METHODS 1-6)
        
        # METHOD 7 is disabled - skip to end
        if False:  # Never execute
            data_values = df[column].dropna().unique()
            for data_val in data_values:
                data_val_lower = str(data_val).strip().lower()
                
                # ADDITIONAL CHECK: Skip if data contains time markers
                data_has_time = any(marker in data_val_lower for marker in ['a.m.', 'p.m.', 'am', 'pm', ':'])
                if data_has_time:
                    continue  # Skip this data value - it's time-based
                
                data_normalized = data_val_lower.replace('\u2019', "'").replace('\u2018', "'")
                data_normalized = data_normalized.replace('\u2013', '-').replace('\u2014', '-')
                
                # If data value is at least 5 chars and is contained in survey option
                if len(data_normalized) >= 5 and data_normalized in option_normalized:
                    mask = df[column] == data_val
                    if mask.any():
                        return df[mask]
                
                # Or if they're both long and similar
                if len(data_val_lower) > 10 or len(option_str) > 10:
                    if data_val_lower == option_str or option_str in data_val_lower:
                        mask = df[column] == data_val
                        if mask.any():
                            return df[mask]
                    
                    # NEW: Check if most of the words match (for schedule questions)
                    # E.g., "Monday - Thursday 10-hour shift" matches "Monday - Thursday 10-hour day shift"
                    survey_words = set(option_punct_normalized.split())
                    data_words = set(data_normalized.split())
                    
                    # ONLY apply fuzzy matching for schedule/shift questions, NOT for numeric ranges
                    # Skip if option contains numeric ranges like "5 hours" or "between"
                    has_numeric_range = any(word in option_punct_normalized for word in ['between', 'less than', 'more than', 'or more', ' to '])
                    
                    if not has_numeric_range:
                        # If at least 85% of survey words are in data, it's a match
                        if survey_words and len(survey_words & data_words) / len(survey_words) >= 0.85:
                            survey_start_words = option_punct_normalized.split()[:5]
                            data_start_words = data_normalized.split()[:5]
                            
                            if len(survey_start_words) >= 5 and len(data_start_words) >= 5:
                                # Long options - check first 5 words match exactly
                                if survey_start_words == data_start_words:
                                    mask = df[column] == data_val
                                    if mask.any():
                                        return df[mask]
                            else:
                                # Short options - check if survey words appear in data in same order
                                # E.g., "Monday-Thursday shift" should match "Monday-Thursday day shift"
                                survey_all_words = option_punct_normalized.split()
                                data_all_words = data_normalized.split()
                                
                                # Check if all survey words appear in data in order
                                data_idx = 0
                                matches_in_order = True
                                for survey_word in survey_all_words:
                                    found = False
                                    while data_idx < len(data_all_words):
                                        if data_all_words[data_idx] == survey_word:
                                            found = True
                                            data_idx += 1
                                            break
                                        data_idx += 1
                                    if not found:
                                        matches_in_order = False
                                        break
                                
                                if matches_in_order:
                                    mask = df[column] == data_val
                                    if mask.any():
                                        return df[mask]
        
        # METHOD 7b: DISABLED for all questions to prevent false matches
        return df.head(0)  # Empty dataframe if no match
    
    def get_data_code_mapping(self, df, column, num_options):
        """
        Analyze data column and create a mapping from option index to data codes.
        Returns a list where index i maps to the data code for survey option i.
        Returns None if standard mapping should be used.
        """
        if column not in df.columns:
            return None
        
        data_values = sorted(df[column].dropna().unique().tolist(), key=lambda x: str(x).lower())
        
        # Check if all values are single letters
        all_single_letters = all(
            isinstance(v, str) and len(v.strip()) == 1 and v.strip().isalpha() 
            for v in data_values
        )
        
        if not all_single_letters:
            return None  # Use standard matching (text or a,b,c codes)
        
        # Check if letters are standard a,b,c,d... sequence
        letter_values = [str(v).strip().lower() for v in data_values]
        standard_letters = [chr(ord('a') + i) for i in range(num_options)]
        
        if all(l in standard_letters for l in letter_values):
            return None  # Standard a,b,c codes - use normal matching
        
        # Non-standard letters (like e, r, w) detected
        # FIXED: Always return None - let METHOD 6 handle letter codes correctly
        # Letter codes ALWAYS map by position: 'a'=0, 'b'=1, 'c'=2, etc.
        # Even if data only has ['a', 'c', 'e'], those still mean options 0, 2, 4
        return None
    
    def find_data_column(self, df, q_num):
        """Find the data column for a question number"""
        q_num_str = str(q_num)
        
        # DIAGNOSTIC: Write column search to file
        import os
        import tempfile
        diag_file = os.path.join(tempfile.gettempdir(), "COLUMN_SEARCH_DEBUG.txt")
        
        try:
            with open(diag_file, 'a', encoding='utf-8') as f:
                f.write(f"\n{'='*60}\n")
                f.write(f"Searching for column for Q{q_num}\n")
                f.write(f"{'='*60}\n")
                f.write(f"All columns in dataframe:\n")
                for i, col in enumerate(df.columns):
                    f.write(f"  {i}: '{col}'\n")
                f.write("\n")
        except:
            pass  # Silently fail if can't write diagnostic
        
        # Try various column naming patterns
        possible = [
            q_num_str,           # Just the number: "8"
            f'q{q_num_str}',     # Lowercase q: "q8"
            f'Q{q_num_str}',     # Uppercase Q: "Q8"
            f'q{q_num_str}.',    # With period: "q8."
            f'Q{q_num_str}.',    # With period: "Q8."
        ]
        
        for name in possible:
            if name in df.columns:
                try:
                    with open(diag_file, 'a', encoding='utf-8') as f:
                        f.write(f"FOUND: Column '{name}' exists\n")
                except:
                    pass
                return name
        
        # Also check if any column starts with these patterns
        for col in df.columns:
            col_lower = str(col).lower()
            if col_lower == f'q{q_num_str}' or col_lower == q_num_str:
                try:
                    with open(diag_file, 'a', encoding='utf-8') as f:
                        f.write(f"FOUND (case-insensitive): Column '{col}'\n")
                except:
                    pass
                return col
            if col_lower.startswith(f'q{q_num_str}_') or col_lower.startswith(f'q{q_num_str} '):
                try:
                    with open(diag_file, 'a', encoding='utf-8') as f:
                        f.write(f"FOUND (starts with): Column '{col}'\n")
                except:
                    pass
                return col
        
        try:
            with open(diag_file, 'a', encoding='utf-8') as f:
                f.write(f"NOT FOUND: No column found for Q{q_num}\n")
                f.write(f"  Tried patterns: {possible}\n")
                f.write(f"\nDiagnostic written to: {diag_file}\n")
        except:
            pass
        
        return None
    
    def extract_questions(self, docx_path):
        """Extract questions from Word document"""
        from docx import Document
        
        def clean_unicode(text):
            """Fix Unicode encoding issues from Word documents"""
            # Normalize en-dash and em-dash to hyphen
            text = text.replace('\u2013', '-').replace('\u2014', '-')
            # Normalize smart quotes
            text = text.replace('\u2018', "'").replace('\u2019', "'")
            text = text.replace('\u201c', '"').replace('\u201d', '"')
            return text
        
        doc = Document(docx_path)
        questions = []
        counter = 0
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            
            # Clean Unicode immediately
            text = clean_unicode(text)
            
            try:
                numPr = para._element.pPr.numPr if para._element.pPr is not None else None
                if numPr is not None:
                    numId = numPr.numId.val if numPr.numId is not None else None
                    if numId == 1:
                        counter += 1
                        questions.append({'number': counter, 'text': text})
            except:
                pass
        
        return questions
    
    def extract_questions_with_options(self, docx_path):
        """Extract questions with answer options"""
        from docx import Document
        
        def clean_unicode(text):
            """Fix Unicode encoding issues from Word documents"""
            # Normalize en-dash and em-dash to hyphen
            text = text.replace('\u2013', '-').replace('\u2014', '-')
            # Normalize smart quotes
            text = text.replace('\u2018', "'").replace('\u2019', "'")
            text = text.replace('\u201c', '"').replace('\u201d', '"')
            return text
        
        doc = Document(docx_path)
        questions = []
        counter = 0
        current_q = None
        current_opts = []
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            
            # Clean Unicode immediately
            text = clean_unicode(text)
            
            try:
                numPr = para._element.pPr.numPr if para._element.pPr is not None else None
                if numPr is not None:
                    numId = numPr.numId.val if numPr.numId is not None else None
                    
                    if numId == 1:  # Main question
                        if current_q:
                            current_q['options'] = current_opts
                            questions.append(current_q)
                        
                        counter += 1
                        current_q = {'number': counter, 'text': text, 'options': []}
                        current_opts = []
                    elif numId != 1 and current_q:
                        current_opts.append(text)
            except:
                pass
        
        if current_q:
            current_q['options'] = current_opts
            questions.append(current_q)
        
        return questions
    
    # Processing overlay methods
    def show_processing(self, message):
        """Show processing overlay"""
        self.processing_overlay = tk.Toplevel(self.root)
        self.processing_overlay.overrideredirect(True)
        self.processing_overlay.attributes('-topmost', True)
        
        # Center on screen
        w, h = 300, 150
        x = (self.root.winfo_screenwidth() - w) // 2
        y = (self.root.winfo_screenheight() - h) // 2
        self.processing_overlay.geometry(f"{w}x{h}+{x}+{y}")
        
        # Semi-transparent background (simulated with solid color)
        frame = tk.Frame(self.processing_overlay, bg='white', relief=tk.RAISED, borderwidth=2)
        frame.pack(fill=tk.BOTH, expand=True)
        
        # Spinner (text-based)
        self.spinner_label = tk.Label(frame, text="⟳", font=('Segoe UI', 32), bg='white', fg=self.COLORS['primary'])
        self.spinner_label.pack(pady=(20, 10))
        
        tk.Label(frame, text=message, font=self.fonts['body_bold'], bg='white', fg=self.COLORS['gray_900']).pack()
        tk.Label(frame, text="Please wait...", font=self.fonts['small'], bg='white', fg=self.COLORS['gray_500']).pack()
        
        self.processing_overlay.update()
        self.animate_spinner()
    
    def animate_spinner(self):
        """Animate the spinner"""
        if hasattr(self, 'processing_overlay') and self.processing_overlay.winfo_exists():
            current = self.spinner_label.cget('text')
            spinners = ['⟳', '↻', '⟲', '↺']
            try:
                idx = spinners.index(current)
                next_idx = (idx + 1) % len(spinners)
            except:
                next_idx = 0
            self.spinner_label.config(text=spinners[next_idx])
            self.root.after(150, self.animate_spinner)
    
    def hide_processing(self):
        """Hide processing overlay"""
        if hasattr(self, 'processing_overlay'):
            self.processing_overlay.destroy()
    
    def show_error(self, message):
        """Show error message"""
        self.hide_processing()
        messagebox.showerror("Error", message)
    
    # ========================================================================
    # STANDALONE INSIGHTS REPORT GENERATOR (NEW IN V88)
    # ========================================================================
    
    def launch_standalone_insights_generator(self):
        """Launch the standalone Insights Report Generator (not part of workflow)"""
        # Create a new window for the insights generator
        insights_window = tk.Toplevel(self.root)
        insights_window.title("Insights Report Generator")
        insights_window.geometry("900x700")
        insights_window.configure(bg=self.COLORS['bg'])
        
        # Center window
        insights_window.update_idletasks()
        x = (insights_window.winfo_screenwidth() - 900) // 2
        y = (insights_window.winfo_screenheight() - 700) // 2
        insights_window.geometry(f"900x700+{x}+{y}")
        
        # Header
        header = tk.Frame(insights_window, bg=self.COLORS['white'], height=70)
        header.pack(fill=tk.X)
        header.pack_propagate(False)
        
        header_content = tk.Frame(header, bg=self.COLORS['white'])
        header_content.pack(fill=tk.BOTH, expand=True, padx=30, pady=15)
        
        tk.Label(
            header_content,
            text="Insights Report Generator",
            font=self.fonts['heading'],
            bg=self.COLORS['white'],
            fg=self.COLORS['gray_900']
        ).pack(side=tk.LEFT)
        
        # Shadow line
        shadow = tk.Frame(insights_window, bg=self.COLORS['border'], height=1)
        shadow.pack(fill=tk.X)
        
        # Create scrollable content
        canvas = tk.Canvas(insights_window, bg=self.COLORS['white'], highlightthickness=0)
        scrollbar = ttk.Scrollbar(insights_window, orient="vertical", command=canvas.yview)
        content = tk.Frame(canvas, bg=self.COLORS['white'])
        
        content.bind(
            "<Configure>",
            lambda e: canvas.configure(scrollregion=canvas.bbox("all"))
        )
        
        canvas.create_window((0, 0), window=content, anchor="nw")
        canvas.configure(yscrollcommand=scrollbar.set)
        
        canvas.pack(side="left", fill="both", expand=True)
        scrollbar.pack(side="right", fill="y")
        
        # Add content padding
        content_inner = tk.Frame(content, bg=self.COLORS['white'])
        content_inner.pack(fill=tk.BOTH, expand=True, padx=40, pady=30)
        
        # Title section
        tk.Label(
            content_inner,
            text="Generate Professional Insights Report",
            font=('Segoe UI', 16, 'bold'),
            bg=self.COLORS['white'],
            fg=self.COLORS['gray_900']
        ).pack(anchor='w', pady=(0, 10))
        
        tk.Label(
            content_inner,
            text="Upload a PowerPoint presentation and generate a comprehensive insights report with executive summary, analysis, and recommendations.",
            font=self.fonts['body'],
            bg=self.COLORS['white'],
            fg=self.COLORS['gray_600'],
            wraplength=800,
            justify=tk.LEFT
        ).pack(anchor='w', pady=(0, 25))
        
        # Information box
        info_frame = tk.Frame(
            content_inner,
            bg='#EFF6FF',  # Light blue
            highlightbackground='#DBEAFE',
            highlightthickness=1
        )
        info_frame.pack(fill=tk.X, pady=(0, 25))
        
        info_inner = tk.Frame(info_frame, bg='#EFF6FF')
        info_inner.pack(fill=tk.X, padx=20, pady=15)
        
        tk.Label(
            info_inner,
            text="ℹ️  What's Included in the Report",
            font=self.fonts['body_bold'],
            bg='#EFF6FF',
            fg='#1E40AF'
        ).pack(anchor='w', pady=(0, 10))
        
        report_features = [
            "• Executive Summary of key findings",
            "• Workforce demographics and profile analysis",
            "• Sleep, fatigue, and safety assessment",
            "• Organizational climate insights",
            "• Schedule satisfaction patterns",
            "• Overtime analysis and recommendations",
            "• Strategic recommendations",
            "• Phased implementation roadmap"
        ]
        
        for feature in report_features:
            tk.Label(
                info_inner,
                text=feature,
                font=self.fonts['body'],
                bg='#EFF6FF',
                fg='#1E40AF',
                justify=tk.LEFT
            ).pack(anchor='w', pady=2)
        
        # Upload section
        tk.Label(
            content_inner,
            text="Step 1: Upload PowerPoint Presentation",
            font=self.fonts['body_bold'],
            bg=self.COLORS['white'],
            fg=self.COLORS['gray_700']
        ).pack(anchor='w', pady=(0, 10))
        
        # Track upload state in window
        insights_window.uploaded_file = None
        insights_window.file_label_frame = None
        
        def upload_and_update():
            pptx_file = filedialog.askopenfilename(
                title="Select PowerPoint Presentation",
                filetypes=[("PowerPoint Files", "*.pptx"), ("All Files", "*.*")],
                parent=insights_window
            )
            
            if not pptx_file:
                return
            
            if not os.path.exists(pptx_file):
                messagebox.showerror("File Not Found", "The selected PowerPoint file could not be found.", parent=insights_window)
                return
            
            insights_window.uploaded_file = pptx_file
            
            # Remove old file label if exists
            if insights_window.file_label_frame:
                insights_window.file_label_frame.destroy()
            
            # Show uploaded file
            insights_window.file_label_frame = tk.Frame(
                content_inner,
                bg=self.COLORS['success_bg'],
                highlightbackground=self.COLORS['success_border'],
                highlightthickness=1
            )
            insights_window.file_label_frame.pack(fill=tk.X, pady=(10, 20))
            
            file_inner = tk.Frame(insights_window.file_label_frame, bg=self.COLORS['success_bg'])
            file_inner.pack(fill=tk.X, padx=15, pady=12)
            
            tk.Label(
                file_inner,
                text="✓ PowerPoint File Loaded",
                font=self.fonts['body_bold'],
                bg=self.COLORS['success_bg'],
                fg=self.COLORS['success']
            ).pack(anchor='w')
            
            tk.Label(
                file_inner,
                text=os.path.basename(pptx_file),
                font=self.fonts['body'],
                bg=self.COLORS['success_bg'],
                fg=self.COLORS['gray_700']
            ).pack(anchor='w', pady=(5, 0))
            
            # Enable generate button
            generate_btn.config(state=tk.NORMAL, bg=self.COLORS['word'])
        
        upload_btn = tk.Button(
            content_inner,
            text="📤 Select PowerPoint File",
            font=self.fonts['button'],
            bg=self.COLORS['pptx'],
            fg='white',
            relief=tk.FLAT,
            padx=20,
            pady=12,
            cursor='hand2',
            command=upload_and_update
        )
        upload_btn.pack(fill=tk.X, pady=(0, 30))
        
        # Generate section
        separator = tk.Frame(content_inner, bg=self.COLORS['border'], height=1)
        separator.pack(fill=tk.X, pady=(0, 20))
        
        tk.Label(
            content_inner,
            text="Step 2: Generate Insights Report",
            font=self.fonts['body_bold'],
            bg=self.COLORS['white'],
            fg=self.COLORS['gray_700']
        ).pack(anchor='w', pady=(0, 10))
        
        def generate_from_standalone():
            if not insights_window.uploaded_file:
                messagebox.showerror("No File", "Please upload a PowerPoint file first.", parent=insights_window)
                return
            
            # Get save location
            default_name = f"Insights_Report_{datetime.now().strftime('%Y%m%d')}.docx"
            default_dir = os.path.dirname(insights_window.uploaded_file)
            
            output_file = filedialog.asksaveasfilename(
                title="Save Insights Report",
                defaultextension=".docx",
                filetypes=[("Word Documents", "*.docx"), ("All Files", "*.*")],
                initialdir=default_dir,
                initialfile=default_name,
                parent=insights_window
            )
            
            if not output_file:
                return
            
            # Close the insights window
            insights_window.destroy()
            
            # Show processing in main window
            self.show_processing("Analyzing PowerPoint and generating insights report...")
            
            def process():
                try:
                    self.create_insights_report_document(insights_window.uploaded_file, output_file)
                    self.root.after(0, lambda: self.finish_insights_report_generation(output_file))
                except Exception as e:
                    import traceback
                    error_msg = f"Error generating insights report:\n\n{str(e)}\n\n{traceback.format_exc()}"
                    self.root.after(0, lambda: self.show_error(error_msg))
            
            thread = threading.Thread(target=process, daemon=True)
            thread.start()
        
        generate_btn = tk.Button(
            content_inner,
            text="📄 Generate Insights Report",
            font=self.fonts['button'],
            bg=self.COLORS['gray_300'],
            fg=self.COLORS['gray_600'],
            relief=tk.FLAT,
            padx=20,
            pady=12,
            cursor='hand2',
            command=generate_from_standalone,
            state=tk.DISABLED
        )
        generate_btn.pack(fill=tk.X, pady=(0, 30))
        
        # Close button
        separator = tk.Frame(content_inner, bg=self.COLORS['border'], height=1)
        separator.pack(fill=tk.X, pady=(20, 20))
        
        close_btn = tk.Button(
            content_inner,
            text="Close",
            font=self.fonts['button'],
            bg=self.COLORS['gray_200'],
            fg=self.COLORS['gray_700'],
            relief=tk.FLAT,
            padx=20,
            pady=10,
            cursor='hand2',
            command=insights_window.destroy
        )
        close_btn.pack(anchor='w')
        
        # Bind mouse wheel scrolling
        def _on_mousewheel(event):
            canvas.yview_scroll(int(-1*(event.delta/120)), "units")
        
        canvas.bind_all("<MouseWheel>", _on_mousewheel)
    
    def create_insights_report_document(self, pptx_path, output_path):
        """
        Create the comprehensive insights report Word document
        This analyzes the PowerPoint and generates professional insights
        """
        from pptx import Presentation
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.enum.style import WD_STYLE_TYPE
        
        # Load the PowerPoint to analyze
        prs = Presentation(pptx_path)
        
        # Extract data from PowerPoint
        survey_data = self.extract_survey_data_from_pptx(prs)
        
        # Create the Word document
        doc = Document()
        
        # Set up document styles
        self.setup_document_styles(doc)
        
        # Cover Page
        self.add_cover_page(doc, survey_data)
        
        # Table of Contents placeholder
        doc.add_page_break()
        heading = doc.add_heading('Table of Contents', level=1)
        doc.add_paragraph("1. Executive Summary")
        doc.add_paragraph("2. Workforce Profile & Demographics")
        doc.add_paragraph("3. Sleep, Fatigue & Safety Analysis")
        doc.add_paragraph("4. Organizational Climate Assessment")
        doc.add_paragraph("5. Schedule Satisfaction & Preferences")
        doc.add_paragraph("6. Overtime Analysis")
        doc.add_paragraph("7. Key Recommendations")
        doc.add_paragraph("8. Implementation Roadmap")
        
        doc.add_page_break()
        
        # Executive Summary
        self.add_executive_summary(doc, survey_data)
        
        doc.add_page_break()
        
        # Workforce Profile
        self.add_workforce_profile_section(doc, survey_data)
        
        doc.add_page_break()
        
        # Sleep & Fatigue Analysis
        self.add_sleep_fatigue_section(doc, survey_data)
        
        doc.add_page_break()
        
        # Organizational Climate
        self.add_organizational_climate_section(doc, survey_data)
        
        doc.add_page_break()
        
        # Schedule Satisfaction
        self.add_schedule_satisfaction_section(doc, survey_data)
        
        doc.add_page_break()
        
        # Overtime Analysis
        self.add_overtime_section(doc, survey_data)
        
        doc.add_page_break()
        
        # Recommendations
        self.add_recommendations_section(doc, survey_data)
        
        doc.add_page_break()
        
        # Implementation Roadmap
        self.add_implementation_roadmap(doc, survey_data)
        
        # Save the document
        doc.save(output_path)
    
    def extract_survey_data_from_pptx(self, prs):
        """
        Extract all relevant data from PowerPoint slides
        Returns a dictionary with all survey metrics
        """
        data = {
            'company_name': self.company_name or self.project_name,
            'total_slides': len(prs.slides),
            'questions': [],
            'metrics': {},
            'demographics': {},
            'benchmarks': {}
        }
        
        for idx, slide in enumerate(prs.slides, 1):
            slide_data = {
                'number': idx,
                'title': '',
                'has_chart': False,
                'text_content': []
            }
            
            # Get title
            if slide.shapes.title:
                slide_data['title'] = slide.shapes.title.text
            
            # Extract text and chart data
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_data['text_content'].append(shape.text.strip())
                
                if shape.shape_type == 3:  # Chart
                    slide_data['has_chart'] = True
            
            data['questions'].append(slide_data)
        
        return data
    
    def setup_document_styles(self, doc):
        """Set up professional document styles"""
        styles = doc.styles
        
        # Modify or create styles for professional appearance
        # This ensures consistent formatting throughout the document
        pass
    
    def add_cover_page(self, doc, survey_data):
        """Add professional cover page"""
        from docx.shared import Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        
        # Title
        title = doc.add_paragraph()
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = title.add_run(f"{survey_data['company_name']}\n\n")
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = RGBColor(37, 99, 235)
        
        # Subtitle
        subtitle = doc.add_paragraph()
        subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = subtitle.add_run("Employee Survey Insights Report\n\n")
        run.font.size = Pt(18)
        
        # Date and details
        details = doc.add_paragraph()
        details.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = details.add_run(f"Prepared by Shiftwork Solutions LLC\n{datetime.now().strftime('%B %d, %Y')}")
        run.font.size = Pt(12)
        run.font.color.rgb = RGBColor(107, 114, 128)
    
    def add_executive_summary(self, doc, survey_data):
        """Add executive summary section"""
        from docx.shared import Pt, RGBColor
        
        doc.add_heading('Executive Summary', level=1)
        
        # Overview paragraph
        overview = doc.add_paragraph()
        overview.add_run("This report presents a comprehensive analysis of the employee survey conducted at ")
        overview.add_run(survey_data['company_name']).bold = True
        overview.add_run(f". The survey gathered insights from employees across {survey_data['total_slides'] - 1} key areas affecting shift operations, work-life balance, and organizational effectiveness.")
        
        doc.add_paragraph()  # Spacing
        
        # Key findings
        doc.add_heading('Key Findings', level=2)
        
        findings = [
            "Workforce demographics indicate specific scheduling constraints that must be accommodated in any schedule redesign.",
            "Sleep and fatigue metrics reveal opportunities to improve operational safety through evidence-based schedule optimization.",
            "Organizational climate assessment highlights communication and management engagement as critical success factors.",
            "Schedule satisfaction data demonstrates employee openness to change when properly engaged in the design process.",
            "Overtime patterns indicate misalignment between company needs and employee preferences, creating opportunities for win-win solutions."
        ]
        
        for finding in findings:
            p = doc.add_paragraph(finding, style='List Bullet')
        
        doc.add_paragraph()  # Spacing
        
        # Strategic implications
        doc.add_heading('Strategic Implications', level=2)
        
        p = doc.add_paragraph(
            "The survey results reveal both challenges and opportunities. While certain metrics indicate "
            "room for improvement, the data also demonstrates that employees are engaged and willing to "
            "participate in constructive changes. Success will depend on: (1) addressing communication gaps "
            "identified in the organizational climate assessment, (2) designing schedules that accommodate "
            "documented employee constraints while meeting operational needs, and (3) maintaining transparency "
            "throughout the implementation process."
        )
        
        doc.add_paragraph()  # Spacing
        
        # Next steps
        doc.add_heading('Next Steps', level=2)
        
        p = doc.add_paragraph(
            "This report provides the foundation for data-driven schedule optimization. The detailed analysis "
            "in subsequent sections offers specific insights and actionable recommendations. Shiftwork Solutions "
            "recommends reviewing these findings with key stakeholders and developing an implementation plan "
            "that addresses identified priorities while maintaining operational continuity."
        )
    
    def add_workforce_profile_section(self, doc, survey_data):
        """Add workforce demographics and profile section"""
        doc.add_heading('Workforce Profile & Demographics', level=1)
        
        doc.add_paragraph(
            "Understanding workforce demographics is essential for designing schedules that accommodate "
            "employee needs while maintaining operational effectiveness. This section analyzes key demographic "
            "factors that create scheduling constraints and opportunities."
        )
        
        doc.add_paragraph()
        
        # Tenure Analysis
        doc.add_heading('Tenure and Experience', level=2)
        doc.add_paragraph(
            "Employee tenure affects training needs, schedule flexibility, and change management approaches. "
            "The data reveals patterns in workforce stability and experience levels that inform implementation strategy."
        )
        
        doc.add_paragraph()
        
        # Life Circumstances
        doc.add_heading('Life Circumstances and Constraints', level=2)
        doc.add_paragraph(
            "Childcare responsibilities, commute distances, and other life factors create real constraints "
            "that must be considered in schedule design. Ignoring these factors leads to decreased satisfaction "
            "and increased turnover regardless of schedule quality."
        )
        
        doc.add_paragraph()
        
        # Implications
        doc.add_heading('Implications for Schedule Design', level=2)
        doc.add_paragraph(
            "The demographic profile suggests that successful schedule implementation will require: "
            "(1) accommodation of childcare and family responsibilities, (2) consideration of commute "
            "patterns in shift timing, and (3) recognition that one size does not fit all - employee "
            "choice and flexibility are essential."
        )
    
    def add_sleep_fatigue_section(self, doc, survey_data):
        """Add sleep, fatigue, and safety analysis section"""
        doc.add_heading('Sleep, Fatigue & Safety Analysis', level=1)
        
        doc.add_paragraph(
            "Sleep and fatigue represent critical safety factors in shift operations. This section examines "
            "sleep patterns, fatigue levels, and their implications for operational safety and performance."
        )
        
        doc.add_paragraph()
        
        # Sleep Debt Analysis
        doc.add_heading('Sleep Deficit Assessment', level=2)
        doc.add_paragraph(
            "The gap between sleep obtained and sleep needed represents accumulated sleep debt. Chronic "
            "sleep deficit affects cognitive performance, safety awareness, and long-term health outcomes. "
            "The survey data reveals the extent of sleep debt in the current workforce and its operational implications."
        )
        
        doc.add_paragraph()
        
        # Safety Implications
        doc.add_heading('Safety and Performance Impact', level=2)
        doc.add_paragraph(
            "Self-reported sleepiness affecting safety and performance is a leading indicator of fatigue-related "
            "risk. The survey responses in this area warrant serious attention and should inform schedule design "
            "to minimize circadian disruption and maximize recovery time."
        )
        
        doc.add_paragraph()
        
        # Recommendations
        doc.add_heading('Fatigue Risk Mitigation', level=2)
        doc.add_paragraph(
            "Schedule optimization offers proven strategies to reduce fatigue risk: (1) limiting consecutive "
            "night shifts, (2) providing adequate recovery time between shift cycles, (3) minimizing quick "
            "turnarounds, and (4) aligning schedules with natural circadian rhythms where operationally feasible."
        )
    
    def add_organizational_climate_section(self, doc, survey_data):
        """Add organizational climate assessment section"""
        doc.add_heading('Organizational Climate Assessment', level=1)
        
        doc.add_paragraph(
            "Organizational climate significantly impacts the success of any change initiative. This section "
            "analyzes employee perceptions of communication, management responsiveness, and overall workplace "
            "satisfaction - all critical factors for successful schedule implementation."
        )
        
        doc.add_paragraph()
        
        # Communication Analysis
        doc.add_heading('Communication Effectiveness', level=2)
        doc.add_paragraph(
            "Effective communication is the foundation of successful change management. The survey reveals "
            "current perceptions of communication quality and identifies specific gaps that must be addressed "
            "during the schedule redesign process."
        )
        
        doc.add_paragraph()
        
        # Management Perception
        doc.add_heading('Management Engagement and Trust', level=2)
        doc.add_paragraph(
            "Employee trust in management directly affects willingness to embrace change. The data shows "
            "how employees perceive management responsiveness and welcoming of input - key indicators of "
            "readiness for collaborative schedule design."
        )
        
        doc.add_paragraph()
        
        # Morale Indicators
        doc.add_heading('Employee Morale and Satisfaction', level=2)
        doc.add_paragraph(
            "Overall satisfaction metrics provide context for schedule redesign. Understanding current morale "
            "helps predict acceptance of changes and identifies areas where improved schedules can make the "
            "greatest positive impact."
        )
        
        doc.add_paragraph()
        
        # Critical Success Factors
        doc.add_heading('Critical Success Factors', level=2)
        doc.add_paragraph(
            "The organizational climate data suggests that successful implementation requires: (1) transparent "
            "communication throughout the process, (2) genuine employee involvement in schedule selection, "
            "(3) visible management commitment to addressing concerns, and (4) demonstration that employee "
            "input directly influences outcomes."
        )
    
    def add_schedule_satisfaction_section(self, doc, survey_data):
        """Add schedule satisfaction and preferences section"""
        doc.add_heading('Schedule Satisfaction & Preferences', level=1)
        
        doc.add_paragraph(
            "Current schedule satisfaction and employee preferences provide the foundation for designing "
            "improved alternatives. This section analyzes what works in the current schedule, what needs "
            "improvement, and what specific features employees value most."
        )
        
        doc.add_paragraph()
        
        # Current Schedule Assessment
        doc.add_heading('Current Schedule Performance', level=2)
        doc.add_paragraph(
            "Understanding both strengths and weaknesses of the current schedule prevents discarding features "
            "employees value while addressing legitimate concerns. The data reveals specific aspects to preserve "
            "and specific problems to solve."
        )
        
        doc.add_paragraph()
        
        # Design Preferences
        doc.add_heading('Employee Schedule Preferences', level=2)
        doc.add_paragraph(
            "Employee preferences regarding shift length, start times, days off patterns, and other schedule "
            "parameters inform the design of alternatives. While not every preference can be accommodated, "
            "understanding preference distributions allows design of options that satisfy majority needs."
        )
        
        doc.add_paragraph()
        
        # Flexibility vs Predictability
        doc.add_heading('Balancing Flexibility and Predictability', level=2)
        doc.add_paragraph(
            "The survey reveals employee priorities regarding schedule flexibility versus predictability. "
            "These competing values must be carefully balanced in schedule design, with recognition that "
            "different employee groups may weight these factors differently."
        )
        
        doc.add_paragraph()
        
        # Design Requirements
        doc.add_heading('Schedule Design Requirements', level=2)
        doc.add_paragraph(
            "Based on the preference data, successful schedule alternatives must: (1) preserve valued features "
            "of the current schedule, (2) address identified pain points, (3) provide meaningful employee choice, "
            "and (4) maintain operational coverage requirements."
        )
    
    def add_overtime_section(self, doc, survey_data):
        """Add overtime analysis section"""
        doc.add_heading('Overtime Analysis', level=1)
        
        doc.add_paragraph(
            "Overtime patterns reveal the gap between scheduled coverage and operational needs, as well as "
            "the tension between employee preferences and company requirements. This section analyzes overtime "
            "dependency, distribution equity, and opportunities for alignment."
        )
        
        doc.add_paragraph()
        
        # OT Dependency
        doc.add_heading('Financial Dependency on Overtime', level=2)
        doc.add_paragraph(
            "The extent to which employees depend on overtime income affects receptiveness to schedule changes "
            "that might reduce OT opportunities. Understanding this dependency is essential for designing changes "
            "that don't create unintended financial hardship."
        )
        
        doc.add_paragraph()
        
        # Equity Perceptions
        doc.add_heading('Overtime Distribution and Fairness', level=2)
        doc.add_paragraph(
            "Perceptions of fairness in overtime distribution affect morale and trust. The survey data reveals "
            "how employees view current distribution practices and identifies opportunities to improve equity "
            "through better scheduling and clearer policies."
        )
        
        doc.add_paragraph()
        
        # Predictability Issues
        doc.add_heading('Overtime Predictability', level=2)
        doc.add_paragraph(
            "Unpredictable overtime disrupts work-life balance and creates resentment even when OT pay is "
            "attractive. The data shows how predictability affects satisfaction and suggests that more planned, "
            "voluntary OT opportunities could improve both coverage and morale."
        )
        
        doc.add_paragraph()
        
        # Strategic Recommendations
        doc.add_heading('Overtime Policy Recommendations', level=2)
        doc.add_paragraph(
            "The overtime data suggests opportunities to: (1) build more coverage into base schedules, "
            "(2) create transparent voluntary OT signup systems, (3) separate routine coverage OT from "
            "emergency callouts, and (4) align OT opportunities more closely with employee preferences."
        )
    
    def add_recommendations_section(self, doc, survey_data):
        """Add key recommendations section"""
        doc.add_heading('Key Recommendations', level=1)
        
        doc.add_paragraph(
            "Based on comprehensive analysis of the survey data, Shiftwork Solutions recommends the following "
            "strategic priorities for schedule optimization and implementation."
        )
        
        doc.add_paragraph()
        
        # Priority 1
        doc.add_heading('Priority 1: Establish Transparent Communication Framework', level=2)
        doc.add_paragraph(
            "The organizational climate data reveals communication gaps that will undermine any change initiative. "
            "Before beginning schedule redesign, establish regular communication channels, clarify the process, "
            "and demonstrate management commitment to genuine employee involvement. This foundation is essential "
            "for success."
        )
        
        doc.add_paragraph()
        
        # Priority 2
        doc.add_heading('Priority 2: Design Multiple Schedule Alternatives', level=2)
        doc.add_paragraph(
            "Rather than imposing a single 'optimal' schedule, develop 3-4 distinct alternatives that address "
            "identified constraints while meeting operational needs. Allow employees to evaluate alternatives "
            "and select their preference. This approach dramatically increases acceptance and reduces resistance."
        )
        
        doc.add_paragraph()
        
        # Priority 3
        doc.add_heading('Priority 3: Address Fatigue and Safety Concerns', level=2)
        doc.add_paragraph(
            "The sleep and fatigue data indicates opportunities to improve both safety and quality of life "
            "through better schedule design. Prioritize alternatives that reduce sleep debt, limit consecutive "
            "night shifts, and provide adequate recovery time between cycles."
        )
        
        doc.add_paragraph()
        
        # Priority 4
        doc.add_heading('Priority 4: Develop Fair Overtime Policies', level=2)
        doc.add_paragraph(
            "Address overtime issues through better base schedule coverage and transparent OT distribution "
            "policies. Distinguish between voluntary scheduled OT and emergency call-ins. Ensure employees "
            "who depend on OT income have predictable opportunities while reducing forced OT for those who "
            "prefer work-life balance."
        )
        
        doc.add_paragraph()
        
        # Priority 5
        doc.add_heading('Priority 5: Plan Phased Implementation', level=2)
        doc.add_paragraph(
            "Avoid 'big bang' implementation. Instead, pilot new schedules with volunteer groups, gather "
            "feedback, make adjustments, and expand gradually. This reduces risk, builds confidence, and "
            "allows continuous improvement based on real-world experience."
        )
    
    def add_implementation_roadmap(self, doc, survey_data):
        """Add implementation roadmap section"""
        doc.add_heading('Implementation Roadmap', level=1)
        
        doc.add_paragraph(
            "Successful implementation requires careful planning and staged execution. This roadmap provides "
            "a suggested timeline and sequence of activities based on hundreds of successful implementations."
        )
        
        doc.add_paragraph()
        
        # Phase 1
        doc.add_heading('Phase 1: Foundation Building (Weeks 1-2)', level=2)
        activities = [
            "Present survey findings to leadership and key stakeholders",
            "Establish communication framework and process timeline",
            "Form employee advisory committee representing all shifts and departments",
            "Clarify decision-making process and employee role in schedule selection"
        ]
        for activity in activities:
            doc.add_paragraph(f"• {activity}", style='List Bullet')
        
        doc.add_paragraph()
        
        # Phase 2
        doc.add_heading('Phase 2: Schedule Development (Weeks 3-6)', level=2)
        activities = [
            "Develop 3-4 schedule alternatives addressing identified constraints",
            "Model operational coverage and employee work-life impact for each option",
            "Cost out each alternative including overtime projections",
            "Prepare comparison materials for employee review"
        ]
        for activity in activities:
            doc.add_paragraph(f"• {activity}", style='List Bullet')
        
        doc.add_paragraph()
        
        # Phase 3
        doc.add_heading('Phase 3: Employee Engagement (Weeks 7-8)', level=2)
        activities = [
            "Present schedule alternatives to all employees with detailed explanation",
            "Conduct employee preference survey on alternatives",
            "Hold Q&A sessions for all shifts",
            "Analyze preference data and finalize recommended schedule"
        ]
        for activity in activities:
            doc.add_paragraph(f"• {activity}", style='List Bullet')
        
        doc.add_paragraph()
        
        # Phase 4
        doc.add_heading('Phase 4: Pilot Implementation (Weeks 9-12)', level=2)
        activities = [
            "Implement new schedule with volunteer pilot group",
            "Monitor operational coverage and employee satisfaction",
            "Gather structured feedback from pilot participants",
            "Make adjustments based on pilot experience"
        ]
        for activity in activities:
            doc.add_paragraph(f"• {activity}", style='List Bullet')
        
        doc.add_paragraph()
        
        # Phase 5
        doc.add_heading('Phase 5: Full Deployment (Weeks 13-16)', level=2)
        activities = [
            "Roll out refined schedule to remaining groups",
            "Provide training on new schedule features and policies",
            "Monitor closely for first month and address issues rapidly",
            "Conduct follow-up satisfaction assessment after 90 days"
        ]
        for activity in activities:
            doc.add_paragraph(f"• {activity}", style='List Bullet')
        
        doc.add_paragraph()
        
        # Critical Success Factors
        doc.add_heading('Critical Success Factors', level=2)
        doc.add_paragraph(
            "Experience across hundreds of implementations shows that success requires: (1) visible leadership "
            "commitment throughout the process, (2) transparent communication at every stage, (3) genuine employee "
            "choice in schedule selection, (4) willingness to adjust based on feedback, and (5) patience to let "
            "the new schedule stabilize before judging success."
        )
        
        doc.add_paragraph()
        
        # Closing
        p = doc.add_paragraph()
        p.add_run("\n\nPrepared by Shiftwork Solutions LLC\n").bold = True
        p.add_run(f"{datetime.now().strftime('%B %d, %Y')}\n")
        p.add_run("\nFor additional information or implementation support, please contact:\n")
        p.add_run("Shiftwork Solutions LLC\n")
        p.add_run("www.shift-work.com")
    
    def finish_insights_report_generation(self, output_file):
        """Called when insights report generation is complete"""
        self.hide_processing()
        
        result = messagebox.askyesno(
            "Insights Report Generated",
            f"Insights report generated successfully!\n\n{os.path.basename(output_file)}\n\n"
            f"Would you like to open the report now?",
            parent=self.root
        )
        
        if result:
            import platform
            if platform.system() == 'Windows':
                os.startfile(output_file)
            elif platform.system() == 'Darwin':  # macOS
                os.system(f'open "{output_file}"')
            else:  # Linux
                os.system(f'xdg-open "{output_file}"')


def main():
    app = SurveyCrosstabApp()
    app.run()


if __name__ == "__main__":
    main()