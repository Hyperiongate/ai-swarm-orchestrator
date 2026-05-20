"""
KNOWLEDGE INGESTION ROUTES
Created: February 2, 2026
Last Updated: May 20, 2026 - ADDED knowledge_search_bp with /api/knowledge/search
              and /api/knowledge/context endpoints for live KB access (Thomas integration).

CHANGELOG:

- May 20, 2026: KNOWLEDGE SEARCH BLUEPRINT ADDED — additive change.

  WHAT WAS ADDED:
    A new Blueprint object `knowledge_search_bp` exposed at module level
    alongside the existing `ingest_bp`. The new blueprint uses url_prefix
    '/api/knowledge' (separate from ingest_bp's '/api/ingest') and exposes
    two read-only endpoints:

      GET /api/knowledge/search
        Query params: q (required), max_results (default 5, max 20),
                      category (optional)
        Returns: JSON {success, results: [{filename, title, category, score,
                 excerpt, word_count, relevance_type}, ...], count, query,
                 kb_ready, max_results}
        Calls EnhancedProjectKnowledgeBase.semantic_search() — already exists
        in knowledge_integration.py. Does NOT modify any data.

      GET /api/knowledge/context
        Query params: q (required), max_context (default 8000, max 16000),
                      max_results (default 3, max 10)
        Returns: JSON {success, context: "...formatted AI-ready string...",
                       query, kb_ready, length, max_context, max_results}
        Calls EnhancedProjectKnowledgeBase.get_context_for_task() — already
        exists. Returns the same prefixed/cited context block the Swarm uses
        internally for task work, ready to inject into a system prompt.

  WHY:
    Thomas's runtime needs live access to the project knowledge base so his
    answers reflect the current Swarm KB rather than knowledge frozen into
    his system prompt at deploy time. Mirrors Thomas's existing pattern of
    calling /api/survey/norm/search on every turn.

  GRACEFUL DEGRADATION:
    - If app.config['KNOWLEDGE_BASE'] is None (KB init failed or no files),
      both endpoints return 503 with a clear message — never crash.
    - If the KB exists but is_ready=False (still warming up in background),
      endpoints return 200 with empty results and kb_ready=false so callers
      can degrade silently rather than block. This matches the existing
      semantic_search() contract which already has a 2-second wait.

  HOW IT IS WIRED:
    app.py (in the May 20, 2026 update) imports BOTH blueprints from this
    module and registers them:
        from routes.ingest import ingest_bp, knowledge_search_bp
        app.register_blueprint(ingest_bp)
        app.register_blueprint(knowledge_search_bp)
    The first line was already there. The second line is new in app.py.
    If for any reason app.py does NOT register knowledge_search_bp, the
    new endpoints simply do not exist — every existing /api/ingest/* route
    continues to work exactly as before.

  SCOPE OF CHANGE TO THIS FILE:
    - Added `current_app` to the flask import line.
    - Added the new `knowledge_search_bp` Blueprint at the bottom of the file.
    - Added the `_get_kb_instance()` helper.
    - Added two new endpoint functions: `knowledge_search()` and
      `knowledge_context()`.
    - No existing function, endpoint, or helper was modified.
    - No existing behavior changed.
    - Rule 1 (do no harm) preserved.

- February 28, 2026 — GAP 2 FIX: 'lifestyle' keyword added to eaf detection.
- February 27, 2026 (Session 3 - Part 2): ADDED lessons_learned_md detection.
- February 27, 2026 (Session 3 - Part 1): FIXED 'engagement' triggering contract mis-classification.
- February 27, 2026 (Session 2): ADDED /api/ingest/batch endpoint.
- February 27, 2026 (Session 1): FIXED DOCX content extraction.
- February 26, 2026 (Session 2): UPDATED PPTX and Excel routes to pass file_bytes.
- February 26, 2026 (Session 1): ADDED new document types and proposal alias.
- February 22, 2026: ADDED /api/ingest/export (GET).
- February 4, 2026: FIXED PowerPoint temp file handling.

Flask API endpoints for document ingestion and knowledge search.
Part of Shoulders of Giants cumulative learning system.
Author: Jim @ Shiftwork Solutions LLC
"""

from flask import Blueprint, request, jsonify, render_template, send_file, current_app
from werkzeug.utils import secure_filename
import os
import sys
import json
from datetime import datetime
import tempfile
import io

try:
    from document_ingestion_engine import get_document_ingestor
    print("Knowledge Ingestion: Direct import succeeded")
except ImportError as e1:
    print(f"Direct import failed: {e1}")
    try:
        sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
        from document_ingestion_engine import get_document_ingestor
        print("Knowledge Ingestion: Path-adjusted import succeeded")
    except ImportError as e2:
        print(f"Path-adjusted import failed: {e2}")
        try:
            root_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
            sys.path.insert(0, root_dir)
            from document_ingestion_engine import get_document_ingestor
            print("Knowledge Ingestion: Root directory import succeeded")
        except ImportError as e3:
            print(f"All import attempts failed!")
            print(f"   Error 1: {e1}")
            print(f"   Error 2: {e2}")
            print(f"   Error 3: {e3}")
            print(f"   Current directory: {os.getcwd()}")
            print(f"   Script directory: {os.path.dirname(__file__)}")
            print(f"   Files in root: {os.listdir(root_dir) if 'root_dir' in locals() else 'N/A'}")
            raise

ingest_bp = Blueprint('ingest', __name__, url_prefix='/api/ingest')

ALLOWED_EXTENSIONS = {
    'txt', 'md', 'pdf', 'docx', 'doc',
    'xlsx', 'xls', 'csv', 'json', 'pptx', 'ppt'
}


def allowed_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


def _detect_document_type(filename):
    """
    Auto-detect document type from filename and extension.
    Mirrors detectDocumentType() in knowledge_management.html exactly.
    """
    name = filename.lower()
    ext  = name.rsplit('.', 1)[-1] if '.' in name else ''

    if ext in ('pptx', 'ppt'):
        if 'oaf' in name or 'operations' in name:
            return 'oaf'
        if 'eaf' in name or 'employee' in name or 'survey' in name or 'lifestyle' in name:
            return 'eaf'
        return 'implementation_ppt'

    if ext in ('xlsx', 'xls'):
        if 'schedule' in name or 'pattern' in name:
            return 'schedule_pattern'
        return 'data_file'

    if ext in ('docx', 'doc'):
        if 'lesson' in name:
            return 'lessons_learned'
        if 'contract' in name or 'agreement' in name:
            return 'contract'
        if 'proposal' in name:
            return 'contract'
        if 'scope' in name:
            return 'scope_of_work'
        if 'implementation' in name and 'manual' in name:
            return 'implementation_manual'
        return 'general_word'

    if ext == 'md':
        if 'lesson' in name:
            return 'lessons_learned_md'
        return 'generic'

    if ext == 'pdf':
        if 'lesson' in name:
            return 'lessons_learned'
        if 'contract' in name or 'agreement' in name:
            return 'contract'
        return 'generic'

    return 'generic'


def _extract_docx_structured(file_bytes: bytes) -> dict:
    """
    Extract structured paragraph data from a .docx file using python-docx.
    """
    result = {'paragraphs': [], 'plain_text': '', 'error': None}
    try:
        from docx import Document
        import io as _io

        doc = Document(_io.BytesIO(file_bytes))
        paragraphs = []
        plain_lines = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            style_name = ''
            if para.style and para.style.name:
                style_name = para.style.name.replace(' ', '')

            runs_with_text = [r for r in para.runs if r.text.strip()]
            is_bold = False
            if runs_with_text:
                is_bold = all(r.bold for r in runs_with_text)
            if 'heading' in style_name.lower():
                is_bold = True

            paragraphs.append({
                'style': style_name,
                'bold': is_bold,
                'text': text
            })
            plain_lines.append(text)

        result['paragraphs'] = paragraphs
        result['plain_text'] = '\n'.join(plain_lines)

    except ImportError:
        result['error'] = 'python-docx not available'
    except Exception as e:
        result['error'] = str(e)

    return result


def _process_file_for_ingest(file, document_type, metadata):
    """
    Shared file processing logic for both single and batch endpoints.
    """
    filename_lower = file.filename.lower()

    if filename_lower.endswith(('.pptx', '.ppt')):
        file_bytes = file.read()
        slide_text_content = None
        try:
            from pptx import Presentation
            tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
            tmp_path = tmp.name
            tmp.write(file_bytes)
            tmp.close()
            try:
                prs = Presentation(tmp_path)
                slide_texts = []
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_content = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text = shape.text.strip()
                            if text:
                                slide_content.append(text)
                    if slide_content:
                        slide_texts.append(f"[Slide {slide_num}]\n" + '\n'.join(slide_content))
                slide_text_content = '\n\n'.join(slide_texts)
            finally:
                try:
                    os.unlink(tmp_path)
                except Exception:
                    pass
        except ImportError:
            pass
        except Exception:
            pass
        if slide_text_content:
            metadata['slide_text_preview'] = slide_text_content[:2000]
        return ingest_document_content(
            ingestor=get_document_ingestor(),
            content=slide_text_content or '',
            document_type=document_type,
            metadata=metadata,
            file_bytes=file_bytes
        )

    elif filename_lower.endswith(('.xlsx', '.xls')):
        file_bytes = file.read()
        excel_type = document_type
        if document_type in ('generic', 'general_word', ''):
            excel_type = 'excel'
        return ingest_document_content(
            ingestor=get_document_ingestor(),
            content='',
            document_type=excel_type,
            metadata=metadata,
            file_bytes=file_bytes
        )

    elif filename_lower.endswith(('.docx', '.doc')):
        file_bytes = file.read()
        docx_data = _extract_docx_structured(file_bytes)
        if docx_data['error'] and not docx_data['paragraphs']:
            content = file_bytes.decode('utf-8', errors='ignore')
            metadata['docx_extraction_error'] = docx_data['error']
        else:
            content = json.dumps(docx_data['paragraphs'], ensure_ascii=False)
            metadata['plain_text'] = docx_data['plain_text'][:5000]
            if docx_data['error']:
                metadata['docx_partial_error'] = docx_data['error']
        return ingest_document_content(
            ingestor=get_document_ingestor(),
            content=content,
            document_type=document_type,
            metadata=metadata,
            file_bytes=file_bytes
        )

    else:
        content = file.read().decode('utf-8', errors='ignore')
        return ingest_document_content(
            ingestor=get_document_ingestor(),
            content=content,
            document_type=document_type,
            metadata=metadata
        )


@ingest_bp.route('/document', methods=['POST'])
def ingest_document():
    """
    Upload and ingest a single document into the knowledge base.
    """
    try:
        if 'file' not in request.files:
            return jsonify({'success': False, 'error': 'No file uploaded'}), 400

        file = request.files['file']
        if file.filename == '':
            return jsonify({'success': False, 'error': 'No file selected'}), 400
        if not allowed_file(file.filename):
            return jsonify({
                'success': False,
                'error': f'File type not allowed. Allowed: {", ".join(sorted(ALLOWED_EXTENSIONS))}'
            }), 400

        document_type = request.form.get('document_type', '').strip()
        if not document_type:
            document_type = _detect_document_type(file.filename)

        metadata = {
            'document_name': secure_filename(file.filename),
            'client': request.form.get('client', ''),
            'industry': request.form.get('industry', ''),
            'project_type': request.form.get('project_type', ''),
            'uploaded_by': request.form.get('uploaded_by', 'user'),
            'upload_date': datetime.now().isoformat()
        }

        result = _process_file_for_ingest(file, document_type, metadata)
        return jsonify(result), 200 if result['success'] else 400

    except Exception as e:
        import traceback
        return jsonify({
            'success': False,
            'error': f'Ingestion failed: {str(e)}',
            'traceback': traceback.format_exc()
        }), 500


# ============================================================================
# BATCH ENDPOINT
# ============================================================================
@ingest_bp.route('/batch', methods=['POST'])
def ingest_batch():
    """
    Upload and ingest multiple documents in a single request.
    """
    try:
        files = request.files.getlist('files')
        if not files or all(f.filename == '' for f in files):
            return jsonify({'success': False, 'error': 'No files uploaded'}), 400

        per_file_types = []
        raw_types = request.form.get('document_types', '')
        if raw_types:
            try:
                per_file_types = json.loads(raw_types)
            except Exception:
                per_file_types = []

        global_type = request.form.get('document_type', '').strip()
        client   = request.form.get('client', '')
        industry = request.form.get('industry', '')

        results = []
        success_count = 0
        error_count = 0

        for idx, file in enumerate(files):
            if file.filename == '':
                continue

            if not allowed_file(file.filename):
                results.append({
                    'filename': file.filename,
                    'success': False,
                    'error': f'File type not allowed ({file.filename.rsplit(".", 1)[-1] if "." in file.filename else "unknown"})'
                })
                error_count += 1
                continue

            doc_type = ''
            if idx < len(per_file_types):
                doc_type = (per_file_types[idx] or '').strip()
            if not doc_type:
                doc_type = global_type
            if not doc_type:
                doc_type = _detect_document_type(file.filename)

            metadata = {
                'document_name': secure_filename(file.filename),
                'client': client,
                'industry': industry,
                'project_type': '',
                'uploaded_by': 'user',
                'upload_date': datetime.now().isoformat()
            }

            try:
                result = _process_file_for_ingest(file, doc_type, metadata)
            except Exception as file_err:
                import traceback
                result = {
                    'success': False,
                    'error': str(file_err),
                    'traceback': traceback.format_exc()
                }

            result['filename']      = file.filename
            result['detected_type'] = doc_type
            results.append(result)

            if result.get('success'):
                success_count += 1
            else:
                error_count += 1

        total_patterns = sum(r.get('patterns_extracted', 0) for r in results)
        total_insights = sum(r.get('insights_extracted', 0) for r in results)

        total_in_kb = 0
        for r in reversed(results):
            if r.get('success') and r.get('total_patterns') is not None:
                total_in_kb = r['total_patterns']
                break

        return jsonify({
            'success': error_count == 0,
            'batch': True,
            'total_files': len(results),
            'success_count': success_count,
            'error_count': error_count,
            'total_patterns_extracted': total_patterns,
            'total_insights_extracted': total_insights,
            'total_patterns': total_in_kb,
            'results': results
        }), 200

    except Exception as e:
        import traceback
        return jsonify({
            'success': False,
            'error': f'Batch ingestion failed: {str(e)}',
            'traceback': traceback.format_exc()
        }), 500
# ============================================================================


def ingest_document_content(ingestor, content, document_type, metadata, file_bytes=None):
    """
    Helper: pass content and metadata to the ingestion engine.
    """
    try:
        if document_type == 'proposal':
            document_type = 'contract'
        kwargs = {
            'content': content,
            'document_type': document_type,
            'metadata': metadata
        }
        if file_bytes is not None:
            kwargs['file_bytes'] = file_bytes
        result = ingestor.ingest_document(**kwargs)
        return result
    except Exception as e:
        import traceback
        return {
            'success': False,
            'error': str(e),
            'traceback': traceback.format_exc()
        }


@ingest_bp.route('/status', methods=['GET'])
def get_status():
    try:
        ingestor = get_document_ingestor()
        stats = ingestor.get_knowledge_base_stats()
        return jsonify({
            'success': True,
            'stats': stats,
            'message': (
                f'Knowledge base contains {stats["total_extracts"]} documents '
                f'and {stats["total_patterns"]} patterns'
            )
        }), 200
    except Exception as e:
        return jsonify({'success': False, 'error': f'Failed to get status: {str(e)}'}), 500


@ingest_bp.route('/history', methods=['GET'])
def get_history():
    try:
        import sqlite3
        limit  = request.args.get('limit',  50, type=int)
        offset = request.args.get('offset',  0, type=int)
        ingestor = get_document_ingestor()
        db = sqlite3.connect(ingestor.db_path)
        db.row_factory = sqlite3.Row
        cursor = db.cursor()
        cursor.execute('SELECT COUNT(*) as count FROM ingestion_log')
        total = cursor.fetchone()['count']
        cursor.execute('''
            SELECT id, document_name, document_type, status,
                   patterns_extracted, insights_extracted,
                   error_message, ingested_at
            FROM ingestion_log
            ORDER BY ingested_at DESC
            LIMIT ? OFFSET ?
        ''', (limit, offset))
        history = [dict(row) for row in cursor.fetchall()]
        db.close()
        return jsonify({
            'success': True, 'total': total, 'limit': limit,
            'offset': offset, 'history': history
        }), 200
    except Exception as e:
        return jsonify({'success': False, 'error': f'Failed to get history: {str(e)}'}), 500


@ingest_bp.route('/patterns', methods=['GET'])
def get_patterns():
    try:
        import sqlite3
        pattern_type   = request.args.get('pattern_type')
        min_confidence = request.args.get('min_confidence', 0.0, type=float)
        limit          = request.args.get('limit', 100, type=int)
        ingestor = get_document_ingestor()
        db = sqlite3.connect(ingestor.db_path)
        db.row_factory = sqlite3.Row
        cursor = db.cursor()
        cursor.execute(
            "SELECT name FROM sqlite_master WHERE type='table' AND name='learned_patterns'"
        )
        if not cursor.fetchone():
            db.close()
            return jsonify({'success': True, 'count': 0, 'patterns': []}), 200
        query = '''
            SELECT id, pattern_type, pattern_name, pattern_data,
                   confidence, supporting_documents,
                   first_seen, last_updated, metadata
            FROM learned_patterns
            WHERE confidence >= ?
        '''
        params = [min_confidence]
        if pattern_type:
            query += ' AND pattern_type = ?'
            params.append(pattern_type)
        query += ' ORDER BY confidence DESC, supporting_documents DESC LIMIT ?'
        params.append(limit)
        try:
            cursor.execute(query, params)
            patterns = []
            for row in cursor.fetchall():
                pattern = dict(row)
                for field in ('pattern_data', 'metadata'):
                    try:
                        pattern[field] = (
                            json.loads(pattern[field]) if pattern.get(field) else {}
                        )
                    except Exception:
                        pattern[field] = {}
                patterns.append(pattern)
            db.close()
            return jsonify({'success': True, 'count': len(patterns), 'patterns': patterns}), 200
        except sqlite3.OperationalError:
            db.close()
            return jsonify({'success': True, 'count': 0, 'patterns': []}), 200
    except Exception as e:
        import traceback
        return jsonify({
            'success': False,
            'error': f'Failed to get patterns: {str(e)}',
            'traceback': traceback.format_exc()
        }), 500


@ingest_bp.route('/extracts', methods=['GET'])
def get_extracts():
    try:
        import sqlite3
        document_type = request.args.get('document_type')
        client        = request.args.get('client')
        industry      = request.args.get('industry')
        limit         = request.args.get('limit', 50, type=int)
        ingestor = get_document_ingestor()
        db = sqlite3.connect(ingestor.db_path)
        db.row_factory = sqlite3.Row
        cursor = db.cursor()
        query = '''
            SELECT id, document_type, document_name,
                   client, industry, project_type,
                   extracted_at, file_size
            FROM knowledge_extracts
            WHERE 1=1
        '''
        params = []
        if document_type:
            query += ' AND document_type = ?'
            params.append(document_type)
        if client:
            query += ' AND client = ?'
            params.append(client)
        if industry:
            query += ' AND industry = ?'
            params.append(industry)
        query += ' ORDER BY extracted_at DESC LIMIT ?'
        params.append(limit)
        cursor.execute(query, params)
        extracts = [dict(row) for row in cursor.fetchall()]
        db.close()
        return jsonify({'success': True, 'count': len(extracts), 'extracts': extracts}), 200
    except Exception as e:
        return jsonify({'success': False, 'error': f'Failed to get extracts: {str(e)}'}), 500


@ingest_bp.route('/extract/<int:extract_id>', methods=['GET'])
def get_extract_detail(extract_id):
    try:
        import sqlite3
        ingestor = get_document_ingestor()
        db = sqlite3.connect(ingestor.db_path)
        db.row_factory = sqlite3.Row
        cursor = db.cursor()
        cursor.execute('SELECT * FROM knowledge_extracts WHERE id = ?', (extract_id,))
        row = cursor.fetchone()
        db.close()
        if not row:
            return jsonify({'success': False, 'error': 'Extract not found'}), 404
        extract = dict(row)
        for field in ('extracted_data', 'metadata'):
            try:
                extract[field] = json.loads(extract[field]) if extract.get(field) else {}
            except Exception:
                pass
        return jsonify({'success': True, 'extract': extract}), 200
    except Exception as e:
        return jsonify({'success': False, 'error': f'Failed to get extract: {str(e)}'}), 500


# ============================================================================
# EXPORT ENDPOINT
# ============================================================================
@ingest_bp.route('/export', methods=['GET'])
def export_knowledge():
    try:
        import sqlite3
        ingestor = get_document_ingestor()
        db = sqlite3.connect(ingestor.db_path)
        db.row_factory = sqlite3.Row
        cursor = db.cursor()
        cursor.execute('''
            SELECT id, document_type, document_name, client, industry,
                   project_type, extracted_at, file_size, extracted_data, metadata
            FROM knowledge_extracts ORDER BY extracted_at DESC
        ''')
        extracts = []
        for row in cursor.fetchall():
            extract = dict(row)
            for field in ('extracted_data', 'metadata'):
                try:
                    extract[field] = json.loads(extract[field]) if extract.get(field) else {}
                except Exception:
                    extract[field] = {}
            extracts.append(extract)
        patterns = []
        try:
            cursor.execute('''
                SELECT id, pattern_type, pattern_name, pattern_data,
                       confidence, supporting_documents, first_seen, last_updated, metadata
                FROM learned_patterns ORDER BY confidence DESC, supporting_documents DESC
            ''')
            for row in cursor.fetchall():
                pattern = dict(row)
                for field in ('pattern_data', 'metadata'):
                    try:
                        pattern[field] = json.loads(pattern[field]) if pattern.get(field) else {}
                    except Exception:
                        pattern[field] = {}
                patterns.append(pattern)
        except sqlite3.OperationalError:
            pass
        ingestion_log = []
        try:
            cursor.execute('''
                SELECT id, document_name, document_type, status,
                       patterns_extracted, insights_extracted, ingested_at
                FROM ingestion_log ORDER BY ingested_at DESC LIMIT 500
            ''')
            ingestion_log = [dict(row) for row in cursor.fetchall()]
        except sqlite3.OperationalError:
            pass
        db.close()
        by_doc_type = {}
        by_industry = {}
        for e in extracts:
            dt  = e.get('document_type') or 'unknown'
            ind = e.get('industry') or 'unspecified'
            by_doc_type[dt]  = by_doc_type.get(dt, 0) + 1
            by_industry[ind] = by_industry.get(ind, 0) + 1
        export_payload = {
            'export_metadata': {
                'exported_at': datetime.now().isoformat(),
                'system': 'Shiftwork Solutions AI Swarm - Knowledge Base',
                'version': '1.0',
                'total_extracts': len(extracts),
                'total_patterns': len(patterns),
                'total_ingestion_log_entries': len(ingestion_log)
            },
            'statistics': {
                'total_documents': len(extracts),
                'total_patterns': len(patterns),
                'by_document_type': by_doc_type,
                'by_industry': by_industry
            },
            'knowledge_extracts': extracts,
            'learned_patterns': patterns,
            'ingestion_log': ingestion_log
        }
        json_bytes = json.dumps(export_payload, indent=2, default=str).encode('utf-8')
        buffer = io.BytesIO(json_bytes)
        buffer.seek(0)
        filename = f"shiftwork_knowledge_export_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json"
        return send_file(
            buffer, mimetype='application/json',
            as_attachment=True, download_name=filename
        )
    except Exception as e:
        import traceback
        return jsonify({
            'success': False,
            'error': f'Export failed: {str(e)}',
            'traceback': traceback.format_exc()
        }), 500
# ============================================================================


# ============================================================================
# CLEAR ENDPOINT
# ============================================================================
@ingest_bp.route('/clear', methods=['POST'])
def clear_knowledge_base():
    try:
        import sqlite3
        body = request.get_json(silent=True) or {}
        if body.get('confirm') != 'CLEAR':
            return jsonify({
                'success': False,
                'error': 'Must send { "confirm": "CLEAR" } in request body'
            }), 400
        ingestor = get_document_ingestor()
        db = sqlite3.connect(ingestor.db_path)
        cursor = db.cursor()
        counts = {}
        for table in ('knowledge_extracts', 'learned_patterns', 'ingestion_log'):
            try:
                cursor.execute(f'SELECT COUNT(*) FROM {table}')
                counts[table] = cursor.fetchone()[0]
                cursor.execute(f'DELETE FROM {table}')
            except sqlite3.OperationalError:
                counts[table] = 0
        db.commit()
        db.close()
        return jsonify({
            'success': True,
            'message': 'Knowledge base cleared successfully',
            'deleted': counts
        }), 200
    except Exception as e:
        import traceback
        return jsonify({
            'success': False,
            'error': f'Clear failed: {str(e)}',
            'traceback': traceback.format_exc()
        }), 500
# ============================================================================


# ============================================================================
# KNOWLEDGE SEARCH BLUEPRINT — Added May 20, 2026
#
# This is a SEPARATE blueprint from ingest_bp because we want the URL prefix
# '/api/knowledge' (not '/api/ingest/knowledge'). It is exported at module
# level so app.py can register it alongside ingest_bp:
#
#     from routes.ingest import ingest_bp, knowledge_search_bp
#     app.register_blueprint(ingest_bp)
#     app.register_blueprint(knowledge_search_bp)
#
# The blueprint exposes two read-only endpoints that let Thomas (and any
# future AI client) query the live project knowledge base on every turn.
# Both endpoints fail gracefully if the KB is unavailable.
# ============================================================================

knowledge_search_bp = Blueprint(
    'knowledge_search',
    __name__,
    url_prefix='/api/knowledge'
)


def _get_kb_instance():
    """
    Return the EnhancedProjectKnowledgeBase singleton from app.config, or
    None if it has not been initialized. The Swarm's app.py stores the
    instance in app.config['KNOWLEDGE_BASE'] immediately after creating it
    (added May 20, 2026). This helper centralizes the lookup so both
    endpoints handle the missing-KB case the same way.
    """
    try:
        return current_app.config.get('KNOWLEDGE_BASE')
    except Exception:
        return None


@knowledge_search_bp.route('/search', methods=['GET'])
def knowledge_search():
    """
    Search the project knowledge base and return structured results.

    Query parameters
    ----------------
    q : str (required)
        The search query string. Forms the basis of semantic + keyword
        matching against indexed documents.
    max_results : int (optional, default 5, max 20)
        Maximum number of result objects to return.
    category : str (optional)
        Filter to a single document category as classified by
        knowledge_integration.py._categorize_document(). Valid values
        include (but are not limited to):
          "Contract", "Implementation Guide", "Survey & Assessment",
          "Schedule Library", "Company Profile", "Lessons Learned",
          "Best Practices Guide", "Executive Summary", "Scope of Work",
          "Reference Material"
        Omit to search across all categories.

    Returns
    -------
    200 OK with JSON body:
        {
            "success":     true,
            "query":       "<query>",
            "category":    "<category or null>",
            "kb_ready":    bool,
            "count":       int,
            "max_results": int,
            "results": [
                {
                    "filename":       "...",
                    "title":          "...",
                    "category":       "...",
                    "score":          float,
                    "excerpt":        "...",
                    "word_count":     int,
                    "relevance_type": "Highly Relevant" | "Very Relevant" | ...
                },
                ...
            ]
        }

    400 BAD REQUEST if 'q' is missing or empty.

    503 SERVICE UNAVAILABLE if the KB singleton does not exist in
    app.config (initialization failed at startup). The response body
    is JSON with success:false and a clear error message.
    """
    query = (request.args.get('q') or '').strip()
    if not query:
        return jsonify({
            'success': False,
            'error':   "Missing required query parameter 'q'."
        }), 400

    try:
        max_results = int(request.args.get('max_results', 5))
    except (TypeError, ValueError):
        max_results = 5
    max_results = max(1, min(max_results, 20))

    category = request.args.get('category') or None
    if category is not None:
        category = category.strip() or None

    kb = _get_kb_instance()
    if kb is None:
        return jsonify({
            'success':  False,
            'error':    'Knowledge base not initialized. '
                        'Check Swarm startup logs and /api/admin/kb-diagnose.',
            'kb_ready': False
        }), 503

    try:
        results = kb.semantic_search(
            query=query,
            max_results=max_results,
            category_filter=category
        )
        return jsonify({
            'success':     True,
            'query':       query,
            'category':    category,
            'kb_ready':    bool(getattr(kb, 'is_ready', False)),
            'count':       len(results),
            'max_results': max_results,
            'results':     results
        }), 200
    except Exception as e:
        import traceback
        return jsonify({
            'success':   False,
            'error':     f'Knowledge search failed: {str(e)}',
            'traceback': traceback.format_exc()
        }), 500


@knowledge_search_bp.route('/context', methods=['GET'])
def knowledge_context():
    """
    Return an AI-ready formatted context block for a query.

    This is the endpoint Thomas's runtime calls on every conversation turn.
    The returned `context` string is already wrapped with a "SHIFTWORK
    SOLUTIONS PROJECT KNOWLEDGE" header and per-source citations — Thomas
    appends it directly to his system prompt without further formatting.

    Query parameters
    ----------------
    q : str (required)
        The search query (typically the visitor's latest message or a
        short summary derived from it).
    max_context : int (optional, default 8000, max 16000)
        Maximum length of the returned context string in characters.
        The KB's get_context_for_task() truncates internally if needed.
    max_results : int (optional, default 3, max 10)
        Maximum number of source documents to include in the context.

    Returns
    -------
    200 OK with JSON body:
        {
            "success":     true,
            "query":       "<query>",
            "kb_ready":    bool,
            "length":      int,
            "max_context": int,
            "max_results": int,
            "context":     "...the formatted context string..."
        }

    The "context" field will be an empty string if the KB is still warming
    up (kb_ready:false) or if no documents matched the query. Callers
    should treat empty context as a normal degradation case.

    400 BAD REQUEST if 'q' is missing or empty.

    503 SERVICE UNAVAILABLE if the KB singleton does not exist in
    app.config.
    """
    query = (request.args.get('q') or '').strip()
    if not query:
        return jsonify({
            'success': False,
            'error':   "Missing required query parameter 'q'."
        }), 400

    try:
        max_context = int(request.args.get('max_context', 8000))
    except (TypeError, ValueError):
        max_context = 8000
    max_context = max(500, min(max_context, 16000))

    try:
        max_results = int(request.args.get('max_results', 3))
    except (TypeError, ValueError):
        max_results = 3
    max_results = max(1, min(max_results, 10))

    kb = _get_kb_instance()
    if kb is None:
        return jsonify({
            'success':  False,
            'error':    'Knowledge base not initialized. '
                        'Check Swarm startup logs and /api/admin/kb-diagnose.',
            'kb_ready': False,
            'context':  ''
        }), 503

    try:
        context = kb.get_context_for_task(
            task_description=query,
            max_context=max_context,
            max_results=max_results
        )
        return jsonify({
            'success':     True,
            'query':       query,
            'kb_ready':    bool(getattr(kb, 'is_ready', False)),
            'length':      len(context),
            'max_context': max_context,
            'max_results': max_results,
            'context':     context
        }), 200
    except Exception as e:
        import traceback
        return jsonify({
            'success':   False,
            'error':     f'Knowledge context failed: {str(e)}',
            'context':   '',
            'traceback': traceback.format_exc()
        }), 500


# I did no harm and this file is not truncated
